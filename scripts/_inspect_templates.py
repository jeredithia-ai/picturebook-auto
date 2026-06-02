"""临时：解析 Worksheet A4 模板 (pptx) + Reading Report 样本 (docx) 的精确版式。

提取：页面尺寸、边距、占位符/文本框位置、字体字号、表格结构。
用于校准底层逻辑里的 A4 规格 + 出血/安全边距。
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from docx import Document
from docx.shared import Emu as DEmu

DL = Path(r"C:\Users\Jered\下载\VIPKID")


def emu_to_mm(v) -> float:
    return round(int(v) / 36000, 2)  # 1 mm = 36000 EMU


def emu_to_cm(v) -> float:
    return round(int(v) / 360000, 2)


def inspect_pptx(path: Path) -> None:
    print("=" * 70)
    print(f"WORKSHEET PPTX: {path.name}")
    print("=" * 70)
    prs = Presentation(str(path))
    w, h = prs.slide_width, prs.slide_height
    print(f"幻灯片尺寸: {emu_to_mm(w)} x {emu_to_mm(h)} mm  ({emu_to_cm(w)} x {emu_to_cm(h)} cm)")
    print(f"  (A4 = 210 x 297 mm 竖版 / 297 x 210 横版)")
    print(f"共 {len(prs.slides)} 页\n")

    for si, slide in enumerate(prs.slides):
        print(f"--- Slide {si + 1} ---")
        for shape in slide.shapes:
            kind = shape.shape_type
            name = shape.name
            try:
                l, t = emu_to_mm(shape.left), emu_to_mm(shape.top)
                wd, ht = emu_to_mm(shape.width), emu_to_mm(shape.height)
                pos = f"L={l} T={t} W={wd} H={ht} mm"
            except Exception:
                pos = "(no geometry)"
            txt = ""
            if shape.has_text_frame:
                runs_info = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        sz = run.font.size.pt if run.font.size else "?"
                        nm = run.font.name or "?"
                        bold = run.font.bold
                        runs_info.append(f"[{nm} {sz}pt b={bold}] {run.text[:30]!r}")
                txt = " | ".join(runs_info[:4])
            print(f"  · {name} ({kind}) {pos}")
            if txt:
                print(f"      text: {txt}")
        print()


def inspect_docx(path: Path) -> None:
    print("=" * 70)
    print(f"READING REPORT DOCX: {path.name}")
    print("=" * 70)
    doc = Document(str(path))
    sec = doc.sections[0]
    print(f"页面尺寸: {emu_to_mm(sec.page_width)} x {emu_to_mm(sec.page_height)} mm")
    print(f"  (A4 = 210 x 297 mm)")
    print(f"边距: 上={emu_to_mm(sec.top_margin)} 下={emu_to_mm(sec.bottom_margin)} "
          f"左={emu_to_mm(sec.left_margin)} 右={emu_to_mm(sec.right_margin)} mm")
    print(f"页眉距={emu_to_mm(sec.header_distance)} 页脚距={emu_to_mm(sec.footer_distance)} mm\n")

    print(f"段落数={len(doc.paragraphs)}  表格数={len(doc.tables)}\n")

    print("--- 前 25 段（含字体）---")
    for i, p in enumerate(doc.paragraphs[:25]):
        if not p.text.strip():
            continue
        runs_info = []
        for run in p.runs[:3]:
            sz = run.font.size.pt if run.font.size else "?"
            nm = run.font.name or "?"
            runs_info.append(f"[{nm} {sz}pt b={run.font.bold}]")
        align = p.alignment
        print(f"  P{i} align={align} {' '.join(runs_info)} {p.text[:50]!r}")

    print("\n--- 表格结构（全部行）---")
    for ti, tbl in enumerate(doc.tables):
        print(f"  Table{ti}: {len(tbl.rows)} 行 x {len(tbl.columns)} 列")
        for ri, row in enumerate(tbl.rows):
            # 去重合并单元格的重复文本
            seen = []
            for c in row.cells:
                t = c.text.strip()
                if t and (not seen or seen[-1] != t):
                    seen.append(t)
            print(f"    R{ri}: {seen}")


def main() -> None:
    ws = next(DL.glob("**/Worksheet  A4模版*.pptx"), None)
    rr = next(DL.glob("**/阅读报告-Level 1-Book 1Final.docx"), None)
    if ws:
        inspect_pptx(ws)
    else:
        print("!! 没找到 worksheet 模板")
    print("\n\n")
    if rr:
        inspect_docx(rr)
    else:
        print("!! 没找到 reading report 样本")


if __name__ == "__main__":
    main()
