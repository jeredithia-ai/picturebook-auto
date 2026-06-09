"""Streamlit 单页 Web App：从老师输入到 4 件套 ZIP 下载。

启动：
    cd picturebook-auto
    streamlit run scripts/web_app.py

流程：
    1. 表单：Title / Level / Book# / Theme / 7 句故事原文 / IP 年龄
    2. AI 抽取（Claude 文本）→ 可编辑预览：词表/分页/语法/拼读/RR 题/Worksheet 题
    3. 一键 Generate All → Picture Book PPT + 7 张图 + Worksheet PPTX + RR DOCX + TG DOCX
    4. 全部打包 ZIP，供老师下载
"""
from __future__ import annotations

import base64
import contextlib
import hashlib
import hmac
import io
import json
import os
import re
import secrets
import shutil
import sys
import time
import zipfile
from concurrent.futures import (
    ThreadPoolExecutor, FIRST_COMPLETED, wait, as_completed,
    TimeoutError as FuturesTimeoutError,
)
from collections.abc import Mapping
from datetime import datetime
from pathlib import Path

# 让 streamlit run 也能找到同级模块
sys.path.insert(0, str(Path(__file__).resolve().parent))

import streamlit as st
import streamlit.components.v1 as components

from ai_extractor import (
    QUESTION_INSTRUCTIONS, QUESTION_POOL, QUESTION_TITLES,
    apply_extracted_to_outline, extract_all, generate_one_worksheet_question,
    generate_story_draft,
)
from auto_fill import auto_summary
from config import (
    DOUBAO_API_KEY, JIMENG_API_KEY, MOCK_AI_EXTRACT, MOCK_IMAGES,
    OUTPUTS_DIR, brand_color_hex, resolve_ip_age, rr_question_distribution,
    COMPOSITION_POLICY, BRAND_DIR,
    render_global_standards_md, render_deliverable_spec_md, DELIVERABLE_SPECS,
)
from cn_prompt_builder import build_cn_page_prompt, page_display_name
from doc_preview import render_to_images, extract_text, has_visual_preview
from parser import BookOutline, PageSpec, enrich_from_syllabus
from ppt_builder import build_picturebook_pptx, safe_filename
from prompt_builder import build_page_prompt  # legacy: fallback only
from reading_report_builder import attach_rr_questions, build_reading_report
from seedream_client import generate_image, generate_image_for_level
from teacher_guide_builder import build_teacher_guide
from worksheet_builder import attach_worksheet_questions, build_worksheet
from curriculum_display import (
    curriculum_section_tables,
    level_metrics_rows,
    mini_map_png_bytes,
    render_level_cards_html,
    render_mini_map_html,
    section_to_rows,
)


def run_with_live_timer(label: str, fn, *args, tick: float = 0.4, done_note: str = "", **kwargs):
    """在子线程执行 fn，主线程每 tick 秒刷新「⏱ 已用 Xs」，完成后显示总用时。

    用于让所有耗时操作（AI 抽取 / 单图重生 / 文档组装等）都能实时看到用时。
    注意：fn 必须是纯计算（不要在里面调用 st.*），st 写入请留在主线程。
    """
    holder = st.empty()
    start = time.time()
    result = None
    with ThreadPoolExecutor(max_workers=1) as ex:
        fut = ex.submit(fn, *args, **kwargs)
        while True:
            try:
                result = fut.result(timeout=tick)
                break
            except FuturesTimeoutError:
                el = time.time() - start
                holder.markdown(
                    f"<div class='timer-pill'><span class='timer-spin'></span>"
                    f"{label} · 已用 <b>{el:.1f}s</b></div>",
                    unsafe_allow_html=True,
                )
    el = time.time() - start
    note = f" · {done_note}" if done_note else ""
    holder.markdown(
        f"<div class='timer-pill done'>✅ {label} · 用时 <b>{el:.1f}s</b>{note}</div>",
        unsafe_allow_html=True,
    )
    return result


LEVEL_OPTIONS = ["Smart", "1", "2", "3", "4", "5", "6"]
SHOT_OPTIONS = ["close", "medium", "full", "wide"]
# v5 机位角度（俯视/仰视/平视/越肩…），eye=平视为默认
ANGLE_OPTIONS = ["eye", "high", "low", "birdseye", "over_shoulder"]
ANGLE_LABELS = {
    "eye": "平视 eye",
    "high": "俯视 high",
    "low": "仰视 low",
    "birdseye": "鸟瞰 birdseye",
    "over_shoulder": "越肩/主角视角 over_shoulder",
}


# =============================== v2.0 严格解锁框架 ===============================
#
# 绘本组装 7 步工作流（每步必须点 ✓ 确认才能进入下一步）：
#   1 📚 输入 + AI 抽取
#   2 🎭 人物 IP 锁定
#   3 🎨 风格背景设定
#   4 📖 分页事件编辑
#   5 📝 提示词预览
#   6 🖼️ 单页生图（Phase 2 升级为 3 候选）
#   7 📦 组装 PPT
# ============================================================================

BOOK_STEPS = [
    (1, "📚 输入 + AI 抽取"),
    (3, "🎭 IP + 🎨 画风锁定"),
    (4, "🖼️ 生图工作台"),
    (5, "📦 组装 4 件套"),
]


# 内部步骤号（历史遗留 1/3/4/5，2 已并入）→ 给用户看的连续序号 1/2/3/4
_STEP_DISPLAY = {num: i + 1 for i, (num, _t) in enumerate(BOOK_STEPS)}


def _step_display(step_num: int) -> int:
    return _STEP_DISPLAY.get(step_num, step_num)


def _step_status(step_num: int) -> str:
    """返回步骤状态：'done' / 'active' / 'locked'。"""
    unlocked = st.session_state.get("book_unlocked_step", 1)
    if step_num < unlocked:
        return "done"
    if step_num == unlocked:
        return "active"
    return "locked"


def _step_icon(status: str) -> str:
    return {"done": "✅", "active": "⏳", "locked": "🔒"}[status]


def _render_progress_bar() -> None:
    """渲染 7 步进度条（顶部）。"""
    cols = st.columns(len(BOOK_STEPS))
    for i, (num, title) in enumerate(BOOK_STEPS):
        status = _step_status(num)
        icon = _step_icon(status)
        color = {"done": "#10b981", "active": "#f59e0b", "locked": "#9ca3af"}[status]
        bg = {"done": "#d1fae5", "active": "#fef3c7", "locked": "#f3f4f6"}[status]
        cols[i].markdown(
            f"<div style='text-align:center;padding:0.5rem 0.2rem;"
            f"background:{bg};border-radius:8px;border:2px solid {color};'>"
            f"<div style='font-size:1.2rem'>{icon}</div>"
            f"<div style='font-size:0.7rem;color:{color};font-weight:600;'>"
            f"Step {_step_display(num)}</div>"
            f"<div style='font-size:0.65rem;color:#6b7280;'>{title}</div>"
            f"</div>",
            unsafe_allow_html=True,
        )


def _locked_step_expander(step_num: int, title: str, default_expanded: bool = True):
    """带状态徽章 + 锁定逻辑的 expander 上下文管理器替身。

    locked → 渲染灰色不可点提示，返回 None（调用方应跳过内容）。
    done/active → 返回 expander 上下文（with ... as exp:）。
    """
    status = _step_status(step_num)
    icon = _step_icon(status)

    if status == "locked":
        st.markdown(
            f"<div style='padding:0.7rem 1rem;background:#f9fafb;border-radius:8px;"
            f"color:#9ca3af;margin:0.4rem 0;border:1px dashed #d1d5db;'>"
            f"{icon} <b>Step {_step_display(step_num)}：{title}</b> &nbsp;&nbsp;"
            f"<span style='font-size:0.85rem'>（请先完成上一步并点 ✅ 确认才能解锁）</span>"
            f"</div>",
            unsafe_allow_html=True,
        )
        return None

    expanded = (status == "active") and default_expanded
    label_color = {"done": "#10b981", "active": "#f59e0b"}[status]
    return st.expander(f"{icon} **Step {_step_display(step_num)}：{title}**", expanded=expanded)


def _confirm_next_step(step_num: int, label: str = "", help_text: str = "") -> None:
    """在某步底部渲染 ✓ 确认按钮，点击后解锁下一步。"""
    if _step_status(step_num) != "active":
        return
    st.markdown("&nbsp;", unsafe_allow_html=True)
    cols = st.columns([3, 1])
    cols[0].caption(help_text or "👉 确认无误后，点击右侧按钮解锁下一步")
    if cols[1].button(
        label or f"✅ 确认 Step {_step_display(step_num)} 并进入下一步",
        type="primary",
        key=f"confirm_step_{step_num}",
        width="stretch",
    ):
        cur = st.session_state.get("book_unlocked_step", 1)
        st.session_state["book_unlocked_step"] = max(cur, step_num + 1)
        # 向导式：确认后自动把视图前进到刚解锁的下一步
        st.session_state["book_view_step"] = step_num + 1
        st.rerun()


def _reset_workflow_button() -> None:
    """顶部重置按钮 — 回到 Step 1。"""
    if st.button("🔄 重置工作流（回到 Step 1）", key="reset_workflow"):
        st.session_state["book_unlocked_step"] = 1
        st.session_state["book_view_step"] = 3
        st.rerun()


def _stage_shell(step_num: int, title: str, wizard: bool):
    """阶段内容容器：向导模式用卡片容器（始终展开）；旧模式用带锁 expander。"""
    if wizard:
        return st.container(border=True)
    return _locked_step_expander(step_num, title)


# 绘本内层「向导」三阶段（内部步骤号 3/4/5 → 展示 1/2/3）
BOOK_WIZARD_STAGES: list[tuple[int, str, str, str]] = [
    (3, "🎭", "锁人物 + 定画风", "选故事角色形象、定全本画风背景"),
    (4, "🖼️", "生成插图", "一键出图、逐页重生 / 锁定满意的图"),
    (5, "📦", "打包交付", "一键产出绘本 + 练习册 + RR + 教师指南"),
]


def _render_book_wizard() -> None:
    """向导式绘本工作流：顶部可点击步骤条 + 一次只显示一个阶段（最易上手）。"""
    unlocked = max(3, st.session_state.get("book_unlocked_step", 3))
    nums = [s[0] for s in BOOK_WIZARD_STAGES]
    view = st.session_state.get("book_view_step", unlocked)
    if view not in nums:
        view = unlocked if unlocked in nums else 3
    if view > unlocked:
        view = unlocked
    st.session_state["book_view_step"] = view

    # —— 顶部进度条（按已解锁进度，Element Plus 风格细条）——
    _total = len(BOOK_WIZARD_STAGES)
    _done = sum(1 for (num, *_r) in BOOK_WIZARD_STAGES if num < unlocked)
    _pct = int(min(_done, _total) / _total * 100)
    st.markdown(
        f"<div class='wiz-track'><div class='wiz-fill' style='width:{_pct}%'></div></div>",
        unsafe_allow_html=True,
    )

    # —— 顶部步骤条（可点击导航；未解锁的禁用）——
    cols = st.columns(len(BOOK_WIZARD_STAGES))
    for i, (num, icon, title, sub) in enumerate(BOOK_WIZARD_STAGES):
        status = "done" if num < unlocked else ("active" if num == unlocked else "locked")
        badge = "✅" if status == "done" else ("🔒" if status == "locked" else icon)
        is_cur = (num == view)
        prefix = "● " if is_cur else ""
        if cols[i].button(
            f"{prefix}{badge}  {i + 1}. {title}",
            key=f"bookwiz_nav_{num}",
            type=("primary" if is_cur else "secondary"),
            width="stretch",
            disabled=(status == "locked"),
            help=sub,
        ):
            st.session_state["book_view_step"] = num
            st.rerun()

    cur_meta = next(s for s in BOOK_WIZARD_STAGES if s[0] == view)
    pos = nums.index(view) + 1
    st.caption(f"第 {pos} / {len(BOOK_WIZARD_STAGES)} 步 — **{cur_meta[2]}**：{cur_meta[3]}")

    # —— 只渲染当前阶段（卡片容器）——
    if view == 3:
        _render_step_ip_style(step_num=3, wizard=True)
    elif view == 4:
        _render_step_workbench(step_num=4, wizard=True)
    elif view == 5:
        _render_step7_assemble(step_num=5, wizard=True)

    # —— 底部上一步 / 重置 ——
    st.markdown("<div style='height:0.4rem'></div>", unsafe_allow_html=True)
    bcol1, bcol2, _ = st.columns([1, 1, 4])
    with bcol1:
        if pos > 1 and st.button("← 上一步", key="bookwiz_prev", width="stretch"):
            st.session_state["book_view_step"] = nums[pos - 2]
            st.rerun()
    with bcol2:
        _reset_workflow_button()


# ============================================================================
# Step 3：🎨 风格背景设定（6 字段）
# ============================================================================

VISUAL_STYLE_OPTIONS = [
    "温暖水彩童书风（推荐 L0-L4）",
    "细腻插画书风（推荐 L5-L6）",
    "扁平卡通风",
    "写实插画风",
    "黑白线稿风",
]

MAIN_SCENE_OPTIONS = [
    "🤖 由故事自动推断",
    "学校教室",
    "家中客厅 / 卧室",
    "户外公园 / 操场",
    "旅行风景",
    "森林 / 动物世界",
    "海边 / 海洋",
    "城市街道",
    "魔法 / 童话场景",
]

COLOR_TONE_OPTIONS = [
    "温暖明亮（暖黄+橘粉，适合家/朋友主题）",
    "清新柔和（淡蓝+薄荷绿，适合学校/成长）",
    "低饱和柔和（米白+灰粉，文艺感）",
    "鲜艳活泼（高饱和，适合冒险/动物）",
    "冷调宁静（蓝紫+雪白，适合夜晚/思考）",
]

COMPOSITION_OPTIONS = [
    "人物为主 50-60%（人物特写）",
    "人物中等 30-40%（推荐 — 留白配文字）",
    "环境为主 20-30%（远景叙事）",
]

LIGHT_SOURCE_OPTIONS = [
    "柔和自然光（推荐）",
    "温暖阳光（侧光，暖色调）",
    "逆光剪影（情绪感）",
    "室内灯光（暖黄）",
    "黄昏 / 黎明（梦幻）",
]

# v3.2 B 层：视角 + 焦点（全局默认，可调）
VIEW_ANGLE_OPTIONS = [
    "平视（推荐 — 与儿童视线齐平，最自然）",
    "微俯视（成人视角看儿童，温柔感）",
    "微仰视（儿童视角看世界，开阔感）",
    "侧面视角（叙事感，多用于互动场景）",
    "🤖 由 AI 按每页内容自动选",
]

FOCUS_HANDLING_OPTIONS = [
    "主体清晰 + 背景轻度虚化（推荐 — 突出主角）",
    "全画面清晰（信息密集型，多元素并存）",
    "主体清晰 + 背景大幅虚化（特写情绪，少元素）",
]

# v3.2 A 层：每页可调的人物状态字段
POSE_OPTIONS = [
    "🤖 由 AI 按故事推断", "站立", "坐着", "蹲下", "趴着",
    "走动", "奔跑", "拥抱", "对话中", "思考中", "举手",
    "手指指向", "回头", "弯腰", "趴在桌上",
]
EMOTION_OPTIONS = [
    "🤖 由 AI 按故事推断", "开心微笑", "紧张不安", "认真专注",
    "惊讶张嘴", "平静放松", "略带难过", "兴奋大笑",
    "好奇凝视", "害羞低头", "自豪挺胸",
]
GAZE_OPTIONS = [
    "🤖 自动", "看镜头（与读者眼神交流）", "看向对话角色",
    "看向手中物品", "看向远方", "看向另一角色背影",
    "低头看地面",
]
POSITION_OPTIONS = [
    "🤖 自动", "画面正中（焦点）", "左前景", "右前景",
    "左中景", "右中景", "中远景", "左侧远景", "右侧远景",
]
INTERACTION_OPTIONS = [
    "🤖 自动", "面对面对话", "并肩同行 / 同伴",
    "一方主动一方被动", "陌生人擦肩", "正在合作做事",
    "对立 / 争执", "围观 / 围聚",
]
TEXT_SAFE_POSITION = [
    "🤖 自动（跟全局留白）", "顶部留白（文字放上方）",
    "底部留白（文字放下方）", "左侧留白", "右侧留白",
    "无文字 / 不留",
]

# 注意：这些会拼进发给 gpt-image-2 的负向文本，Azure 安全审核只看词本身，
# 不能写敏感词（暴力/血腥/裸露等），否则误判 safety_violations 拦截整页。
_DEFAULT_GLOBAL_AVOID = (
    "丑陋 / 畸形\n"
    "多手指 / 错位关节\n"
    "字幕 / 水印 / logo\n"
    "阴暗压抑画风\n"
    "浓妆 / 成熟感\n"
    "低分辨率 / 模糊"
)


def _render_style_panel_step(step_num: int, embed: bool = False) -> None:
    """Step 3：风格背景设定（6 字段）。embed=True 时内联渲染（无 expander/confirm）。"""
    if embed:
        exp = contextlib.nullcontext()
    else:
        exp = _locked_step_expander(step_num, "🎨 风格背景设定（全局应用到所有 7 页）")
        if exp is None:
            return

    with exp:
        st.caption(
            "🤖 **默认全自动**：视觉风格 / 主场景 / 色调 / 构图 / 视角 等都由 AI 按本页故事自动决定，"
            "并自动注入每页提示词。一般**不用动**；只在需要时补下面的「必出现元素」，或展开「高级覆盖」微调。"
        )

        defaults = st.session_state.get("style_config", {})

        # 默认全自动取值（不展开高级覆盖时直接用这些默认）
        visual_style = defaults.get("visual_style", VISUAL_STYLE_OPTIONS[0])
        main_scene = defaults.get("main_scene", MAIN_SCENE_OPTIONS[0])
        color_tone = defaults.get("color_tone", COLOR_TONE_OPTIONS[0])
        composition = defaults.get("composition", COMPOSITION_OPTIONS[1])
        light_source = defaults.get("light_source", LIGHT_SOURCE_OPTIONS[0])
        text_safe_zone = defaults.get("text_safe_zone", True)
        view_angle = defaults.get("view_angle", VIEW_ANGLE_OPTIONS[0])
        focus_handling = defaults.get("focus_handling", FOCUS_HANDLING_OPTIONS[0])

        # 高级覆盖用复选框开关（不能用嵌套 expander）
        if st.checkbox("⚙️ 显示高级覆盖（默认全自动，一般不用动）", value=False, key="cfg_show_adv"):
            col1, col2, col3 = st.columns(3)
            with col1:
                visual_style = st.selectbox(
                    "① 视觉风格",
                    VISUAL_STYLE_OPTIONS,
                    index=VISUAL_STYLE_OPTIONS.index(defaults.get("visual_style", VISUAL_STYLE_OPTIONS[0]))
                    if defaults.get("visual_style") in VISUAL_STYLE_OPTIONS else 0,
                    key="cfg_visual",
                    help="决定全本水彩童书画风的关键关键词",
                )
            with col2:
                main_scene = st.selectbox(
                    "② 主场景",
                    MAIN_SCENE_OPTIONS,
                    index=MAIN_SCENE_OPTIONS.index(defaults.get("main_scene", MAIN_SCENE_OPTIONS[0]))
                    if defaults.get("main_scene") in MAIN_SCENE_OPTIONS else 0,
                    key="cfg_scene",
                    help="选 🤖 自动 → AI 按每页文本逐页推断；选具体场景 → 强制锁定",
                )
            with col3:
                color_tone = st.selectbox(
                    "③ 色调",
                    COLOR_TONE_OPTIONS,
                    index=COLOR_TONE_OPTIONS.index(defaults.get("color_tone", COLOR_TONE_OPTIONS[0]))
                    if defaults.get("color_tone") in COLOR_TONE_OPTIONS else 0,
                    key="cfg_tone",
                )

            col4, col5, col6 = st.columns(3)
            with col4:
                composition = st.selectbox(
                    "④ 构图（人物占比）",
                    COMPOSITION_OPTIONS,
                    index=1,  # 默认推荐
                    key="cfg_comp",
                )
            with col5:
                light_source = st.selectbox(
                    "⑤ 光源",
                    LIGHT_SOURCE_OPTIONS,
                    index=LIGHT_SOURCE_OPTIONS.index(defaults.get("light_source", LIGHT_SOURCE_OPTIONS[0]))
                    if defaults.get("light_source") in LIGHT_SOURCE_OPTIONS else 0,
                    key="cfg_light",
                )
            with col6:
                text_safe_zone = st.checkbox(
                    "⚠️ 顶/底 25% 留白给文字",
                    value=defaults.get("text_safe_zone", True),
                    key="cfg_textsafe",
                    help="开启后，提示词会强制顶部和底部留白，避免文字遮住主体",
                )

            # v3.2 B 层：视角 + 焦点
            col7, col8 = st.columns(2)
            with col7:
                view_angle = st.selectbox(
                    "⑦ 视角（B 层 · 镜头角度）",
                    VIEW_ANGLE_OPTIONS,
                    index=VIEW_ANGLE_OPTIONS.index(defaults.get("view_angle", VIEW_ANGLE_OPTIONS[0]))
                    if defaults.get("view_angle") in VIEW_ANGLE_OPTIONS else 0,
                    key="cfg_view_angle",
                    help="推荐用平视，与儿童视线齐平最自然。AI 自动模式 = Claude 按每页内容选",
                )
            with col8:
                focus_handling = st.selectbox(
                    "⑧ 焦点处理（B 层 · 景深）",
                    FOCUS_HANDLING_OPTIONS,
                    index=FOCUS_HANDLING_OPTIONS.index(defaults.get("focus_handling", FOCUS_HANDLING_OPTIONS[0]))
                    if defaults.get("focus_handling") in FOCUS_HANDLING_OPTIONS else 0,
                    key="cfg_focus",
                    help="主体清晰 + 背景虚化 → 突出主角；全清晰 → 适合信息密集页面",
                )

        st.markdown("---")
        col7, col8 = st.columns(2)
        with col7:
            global_must = st.text_area(
                "📌 全局必出现元素（每行一个，会追加到每页正向提示词）",
                value=defaults.get("global_must", ""),
                height=120,
                key="cfg_must",
                placeholder="例如：\nAnna 头戴白色发箍\nMia 始终扎单束高马尾\n所有场景必有教室元素",
            )
        with col8:
            global_avoid = st.text_area(
                "🚫 全局必避免元素（每行一个，会追加到每页反向提示词）",
                value=defaults.get("global_avoid", _DEFAULT_GLOBAL_AVOID),
                height=120,
                key="cfg_avoid",
            )

        new_cfg = {
            "visual_style": visual_style,
            "main_scene": main_scene,
            "color_tone": color_tone,
            "composition": composition,
            "light_source": light_source,
            "text_safe_zone": text_safe_zone,
            "global_must": global_must,
            "global_avoid": global_avoid,
            # v3.2 B 层新增
            "view_angle": view_angle,
            "focus_handling": focus_handling,
        }
        st.session_state["style_config"] = new_cfg

        # 实时预览：把当前 6 字段拼成一段中文 style block，给老师看效果
        style_preview = _format_style_config_preview(new_cfg)
        with st.expander("🔍 实时风格设定预览（会注入到每页提示词）", expanded=False):
            st.markdown(f"**全局正向（追加到每页正向末尾）：**")
            st.code(style_preview["positive_block"], language="text")
            st.markdown(f"**全局反向（追加到每页反向末尾）：**")
            st.code(style_preview["negative_block"], language="text")

        if not embed:
            _confirm_next_step(
                step_num,
                label="✅ 确认风格设定，进入分页编辑",
                help_text="👉 风格设定一旦确认，会自动注入到 7 页的提示词。后续仍可回到本步重设。",
            )


# ============================================================================
# Step 2：🎭 人物 IP 锁定（包住已有的 auto_summary_panel）
# ============================================================================

def _render_step2_ip_lock(step_num: int, embed: bool = False) -> None:
    if embed:
        exp = contextlib.nullcontext()
    else:
        exp = _locked_step_expander(step_num, "🎭 人物 IP 锁定（确认故事里的角色形象）")
        if exp is None:
            return
    with exp:
        st.caption(
            "💡 系统已从故事里识别到的角色 + 让你**勾选 IP 库**里的形象做参考图。"
            "如果故事里出现 \"a girl\" / \"a boy\" 这类未命名角色，可以在下方"
            "把它们映射到具体 IP（如 Mia / Tommy / Anna）。"
        )
        if st.session_state.get("auto") is not None:
            _render_auto_summary_panel()
        else:
            st.info("ℹ️ Step 1 抽取完成后会显示主角识别面板。")

        if not embed:
            _confirm_next_step(
                step_num,
                label="✅ 确认人物 IP，进入风格设定",
                help_text="👉 IP 一旦确认，生图时会自动加载对应的参考图（多角色场景最多 3 张）",
            )


# ============================================================================
# Step 4（合并 4+5）：📖 分页事件 + 提示词（事实和 prompt 同页编辑，实时联动）
# ============================================================================

def _render_step4_combined(step_num: int, embed: bool = False) -> None:
    """合并版：每页一张大卡片，左列编辑事实 / 右列实时显示 prompt + DeepSeek 按钮。"""
    if embed:
        exp = contextlib.nullcontext()
    else:
        exp = _locked_step_expander(
            step_num, "📖 分页事件 + 提示词（左编事实 / 右看 prompt / Claude 智能润色）"
        )
        if exp is None:
            return
    with exp:
        st.caption(
            "💡 一张卡片 = 一页绘本。"
            "**左列**编辑故事/画面事实，**右列**实时显示 prompt（点 🔄 重建按钮同步）。"
            "**🩺 Claude 润色** 把当前 prompt 喂给 Claude 优化（需配置 IMAROUTER_API_KEY）。"
        )

        # 检查 DeepSeek 可用性
        try:
            from scene_cn_writer import is_available as ds_available
            deepseek_ok = ds_available()
        except Exception:
            deepseek_ok = False

        if not deepseek_ok:
            st.warning(
                "⚠️ 文本模型未配置。🤖 写 scene_cn / 🩺 润色按钮将不可用。"
                "请在 .env 设置 `IMAROUTER_API_KEY`（文本默认走 `claude-opus-4-7`）。"
            )

        # 全局重建按钮
        col_rebuild, col_ds_all = st.columns([1, 1])
        with col_rebuild:
            if st.button(
                "🔄 基于最新事实/风格重建全部 prompt",
                type="secondary",
                key="s4_rebuild_all",
                width="stretch",
                help="把 Step 3 风格设定 + 当前分页事实 重新组装成提示词（覆盖已编辑的 prompt）",
            ):
                _rebuild_all_cn_prompts()
                st.success("✅ 已重建全部 prompt")
                st.rerun()
        with col_ds_all:
            if st.button(
                "🤖 用 Claude 重写全部 scene_cn（强烈推荐）",
                type="secondary",
                key="s4_ds_all_scenes",
                width="stretch",
                disabled=not deepseek_ok,
                help="让 Claude 给 7 页全部写新的 120-220 字中文画面描述",
            ):
                _deepseek_rewrite_all_scenes()
                st.rerun()

        # 科普非虚构：科学事实校验（核查文字+画面正确性/比例/图文一致）
        outline = st.session_state.get("outline")
        is_nf = False
        if outline is not None:
            try:
                from cn_prompt_builder import _is_nonfiction
                is_nf = _is_nonfiction(outline)
            except Exception:
                is_nf = False
        if is_nf:
            if st.button(
                "🔬 科学事实校验（科普必做）",
                type="secondary",
                key="s4_factcheck",
                width="stretch",
                disabled=not deepseek_ok,
                help="核查每页文字+画面的科学事实正确性、真实比例、图文一致；可一键应用修正",
            ):
                _run_fact_check()
                st.rerun()
            _render_fact_check_results()

        st.divider()

        ec = st.session_state.extracted
        page_prompts = st.session_state.get("page_prompts") or {}
        ec_pages_by_idx = {p.get("index"): p for p in (ec.pages or [])}

        for idx in sorted(page_prompts.keys()):
            entry = page_prompts[idx]
            ec_page = ec_pages_by_idx.get(idx) or {}
            display_name = entry.get("display_name") or page_display_name(idx)
            label_kind = "封面" if idx == 0 else "故事页"
            refs = entry.get("references", [])
            ref_badge = f"🖼️ {len(refs)} 张参考图" if refs else "⚠️ 无参考图"

            st.markdown(f"#### {display_name} · {label_kind} · {ref_badge}")

            # 单页 DeepSeek 操作按钮
            bcol1, bcol2, bcol3 = st.columns([1, 1, 1])
            with bcol1:
                if st.button(
                    "🤖 Claude 写 scene_cn",
                    key=f"s4_ds_scene_{idx}",
                    disabled=not deepseek_ok,
                    width="stretch",
                    help="让 Claude 重写本页 120-220 字中文画面描述",
                ):
                    _deepseek_rewrite_one_scene(idx)
                    st.rerun()
            with bcol2:
                if st.button(
                    "🩺 Claude 润色 prompt",
                    key=f"s4_ds_polish_{idx}",
                    disabled=not deepseek_ok,
                    width="stretch",
                    help="把本页正向 prompt 喂给 Claude 按画面结构最佳实践重写",
                ):
                    _deepseek_polish_one_prompt(idx)
                    st.rerun()
            with bcol3:
                if st.button(
                    "🔄 仅重建本页 prompt（用模板）",
                    key=f"s4_rebuild_one_{idx}",
                    width="stretch",
                    help="只用 cn_prompt_builder 模板重建本页（不调 Claude）",
                ):
                    _rebuild_one_cn_prompt(idx)
                    st.rerun()

            # 双栏布局
            left, right = st.columns([1, 1])

            # === 左列：事实编辑 ===
            with left:
                st.caption("① 故事原文（英文）")
                new_text = st.text_area(
                    f"text_{idx}",
                    value=ec_page.get("text", ""),
                    height=100,
                    key=f"s4_text_{idx}",
                    label_visibility="collapsed",
                )
                if new_text != ec_page.get("text", ""):
                    ec_page["text"] = new_text

                st.caption("② 中文画面描述 scene_cn（120-220 字）")
                new_scene_cn = st.text_area(
                    f"scene_cn_{idx}",
                    value=ec_page.get("scene_cn", ""),
                    height=150,
                    key=f"s4_scene_cn_{idx}",
                    label_visibility="collapsed",
                    placeholder=(
                        "主体（人物+具体外观）+ 动作（具体动词姿势）+ 环境（可见物品）+ 氛围（光照）\n"
                        "👆 点上方 🤖 Claude 写 scene_cn 让它帮你写"
                    ),
                )
                if new_scene_cn != ec_page.get("scene_cn", ""):
                    ec_page["scene_cn"] = new_scene_cn

                fc1, fc2 = st.columns(2)
                with fc1:
                    st.caption("③ 必出现（每行一条）")
                    entry["must_include"] = st.text_area(
                        f"must_inc_{idx}",
                        value=entry.get("must_include", ""),
                        height=110,
                        key=f"s4_must_{idx}",
                        label_visibility="collapsed",
                        placeholder="Anna 头戴白色发箍\n桌上 5 本课本",
                    )
                with fc2:
                    st.caption("④ 必避免（每行一条）")
                    entry["page_avoid"] = st.text_area(
                        f"page_avoid_{idx}",
                        value=entry.get("page_avoid", ""),
                        height=110,
                        key=f"s4_avoid_{idx}",
                        label_visibility="collapsed",
                        placeholder="本页不应有宠物\nMia 不能散发",
                    )

                cur_shot = ec_page.get("shot", "medium") or "medium"
                new_shot = st.selectbox(
                    "镜头",
                    SHOT_OPTIONS,
                    index=max(0, SHOT_OPTIONS.index(cur_shot) if cur_shot in SHOT_OPTIONS else 1),
                    key=f"s4_shot_{idx}",
                )
                if new_shot != cur_shot:
                    ec_page["shot"] = new_shot

                # v3.2 A 层：🎭 人物状态字段（每页可调）
                with st.expander("🎭 A 层 · 人物状态（推荐每页都设，强约束）", expanded=False):
                    p_cur = entry.get("page_constraints") or {}

                    cc1, cc2 = st.columns(2)
                    with cc1:
                        pose = st.selectbox(
                            "主角姿势",
                            POSE_OPTIONS,
                            index=POSE_OPTIONS.index(p_cur.get("pose", POSE_OPTIONS[0]))
                            if p_cur.get("pose") in POSE_OPTIONS else 0,
                            key=f"s4a_pose_{idx}",
                        )
                        emotion = st.selectbox(
                            "情绪表达",
                            EMOTION_OPTIONS,
                            index=EMOTION_OPTIONS.index(p_cur.get("emotion", EMOTION_OPTIONS[0]))
                            if p_cur.get("emotion") in EMOTION_OPTIONS else 0,
                            key=f"s4a_emo_{idx}",
                        )
                        gaze = st.selectbox(
                            "视线方向",
                            GAZE_OPTIONS,
                            index=GAZE_OPTIONS.index(p_cur.get("gaze", GAZE_OPTIONS[0]))
                            if p_cur.get("gaze") in GAZE_OPTIONS else 0,
                            key=f"s4a_gaze_{idx}",
                        )
                    with cc2:
                        position = st.selectbox(
                            "画面位置",
                            POSITION_OPTIONS,
                            index=POSITION_OPTIONS.index(p_cur.get("position", POSITION_OPTIONS[0]))
                            if p_cur.get("position") in POSITION_OPTIONS else 0,
                            key=f"s4a_position_{idx}",
                        )
                        interaction = st.selectbox(
                            "互动关系（多角色时）",
                            INTERACTION_OPTIONS,
                            index=INTERACTION_OPTIONS.index(p_cur.get("interaction", INTERACTION_OPTIONS[0]))
                            if p_cur.get("interaction") in INTERACTION_OPTIONS else 0,
                            key=f"s4a_inter_{idx}",
                        )
                        text_pos = st.selectbox(
                            "文字预留位置",
                            TEXT_SAFE_POSITION,
                            index=TEXT_SAFE_POSITION.index(p_cur.get("text_pos", TEXT_SAFE_POSITION[0]))
                            if p_cur.get("text_pos") in TEXT_SAFE_POSITION else 0,
                            key=f"s4a_textpos_{idx}",
                        )
                    focus_char = st.text_input(
                        "焦点角色（输入名字，留空 = AI 自动判断）",
                        value=p_cur.get("focus_char", ""),
                        key=f"s4a_focus_{idx}",
                        placeholder="例：Anna",
                    )
                    entry["page_constraints"] = {
                        "pose": pose, "emotion": emotion, "gaze": gaze,
                        "position": position, "interaction": interaction,
                        "text_pos": text_pos, "focus_char": focus_char,
                    }

            # === 右列：prompt 编辑 ===
            with right:
                st.caption("✅ 正向 Prompt（火山风单段流畅自然语言）")
                cur_pos = entry.get("positive", entry.get("prompt", ""))
                new_pos = st.text_area(
                    label="正向 Prompt",
                    value=cur_pos,
                    height=200,
                    key=f"s4_positive_{idx}",
                    label_visibility="collapsed",
                )
                entry["positive"] = new_pos

                st.caption("❌ 反向 Prompt（每行一条）")
                cur_neg = entry.get("negative", "")
                new_neg = st.text_area(
                    label="反向 Prompt",
                    value=cur_neg,
                    height=180,
                    key=f"s4_negative_{idx}",
                    label_visibility="collapsed",
                )
                entry["negative"] = new_neg

            # 折叠显示最终组装效果
            with st.expander(
                f"🔍 {display_name} · 最终组装效果（含全局风格 + 必出现/避免）",
                expanded=False,
            ):
                final_pos, final_neg = _preview_final_prompt(idx)
                st.markdown("**✅ 正向（最终）：**")
                st.code(final_pos, language="text")
                st.markdown("**❌ 反向（最终）：**")
                st.code(final_neg, language="text")

            if refs:
                st.caption(
                    "🎨 参考图: " + " | ".join(Path(r).name for r in refs[:4])
                )

            st.divider()

        if not embed:
            _confirm_next_step(
                step_num,
                label="✅ 事实和提示词都 OK，进入生图",
                help_text="👉 确认后进入生图步骤。建议至少跑过 Claude 润色（如已配置）",
            )


def _preview_final_prompt(idx: int) -> tuple[str, str]:
    """返回某页（已加全局风格 + A 层人物状态 + 必出现/避免）的最终正向/反向 prompt 文本，供 UI 预览。"""
    page_prompts = st.session_state.get("page_prompts") or {}
    entry = page_prompts.get(idx) or {}
    final_pos = entry.get("positive", "")
    final_neg = entry.get("negative", "")
    style_cfg = st.session_state.get("style_config") or {}
    blocks = _format_style_config_preview(style_cfg) if style_cfg else {
        "positive_block": "", "negative_block": ""
    }
    if blocks["positive_block"]:
        final_pos = f"【全局风格】{blocks['positive_block']}\n\n{final_pos}"
    if blocks["negative_block"]:
        final_neg = (final_neg.rstrip() + "\n" if final_neg else "") + blocks["negative_block"]
    # v3.2 A 层注入
    pc_block = _format_page_constraints(entry.get("page_constraints"))
    if pc_block:
        final_pos = final_pos.rstrip() + "\n\n【A 层 · 本页人物状态约束】\n" + pc_block
    must_inc = (entry.get("must_include") or "").strip()
    if must_inc:
        final_pos = final_pos.rstrip() + "\n\n【教师锁定 · 必须出现】\n" + must_inc
    page_avoid = (entry.get("page_avoid") or "").strip()
    if page_avoid:
        final_neg = (final_neg.rstrip() + "\n" if final_neg else "") + page_avoid
    # v3.2 反馈学习的负向（来自问题反馈面板）
    learned_neg = (entry.get("learned_negatives") or "").strip()
    if learned_neg:
        final_neg = (final_neg.rstrip() + "\n" if final_neg else "") + learned_neg
    return final_pos, final_neg


def _rebuild_one_cn_prompt(page_idx: int) -> None:
    """只重建单页的 prompt（用 cn_prompt_builder 模板）。"""
    ec = st.session_state.extracted
    outline = st.session_state.outline
    if not ec or not outline:
        return
    apply_extracted_to_outline(outline, ec)
    ip_age = outline.ip_age or resolve_ip_age(outline.level)
    page_prompts = st.session_state.get("page_prompts") or {}
    cast_pool = st.session_state.get("story_cast_pool") or None
    generic_overrides = st.session_state.get("generic_overrides") or None
    for page in outline.pages:
        if page.index != page_idx:
            continue
        built = build_cn_page_prompt(
            page, outline, ip_age,
            cast_pool=cast_pool, generic_overrides=generic_overrides,
        )
        prev_must = (page_prompts.get(page_idx) or {}).get("must_include", "")
        prev_avoid = (page_prompts.get(page_idx) or {}).get("page_avoid", "")
        page_prompts[page_idx] = {
            "positive": built.positive,
            "negative": built.negative,
            "prompt": built.prompt,
            "references": [str(r) for r in built.references],
            "must_include": prev_must,
            "page_avoid": prev_avoid,
            "label": page.label,
            "display_name": page_display_name(page.index),
        }
        break
    st.session_state.page_prompts = page_prompts
    st.success(f"✅ 已重建 {page_display_name(page_idx)} 的 prompt")


def _deepseek_rewrite_one_scene(page_idx: int) -> None:
    """让 DeepSeek 重写单页的 scene_cn。"""
    from scene_cn_writer import write_scene_cn, DeepSeekError

    ec = st.session_state.extracted
    outline = st.session_state.outline
    if not ec or not outline:
        return
    page = next((p for p in outline.pages if p.index == page_idx), None)
    ec_page = next((p for p in (ec.pages or []) if p.get("index") == page_idx), None)
    if not page or not ec_page:
        return

    page_prompts = st.session_state.get("page_prompts") or {}
    entry = page_prompts.get(page_idx) or {}
    style_cfg = st.session_state.get("style_config") or {}
    blocks = _format_style_config_preview(style_cfg) if style_cfg else {"positive_block": ""}
    cast_pool = st.session_state.get("story_cast_pool") or []

    # 取人物描述
    cast_descs = _cast_descriptions_from_pool(cast_pool, int(outline.ip_age or 12))

    # 取前几页 scene_cn 做摘要（保持连续性）
    prev_summary = ""
    for prev_idx in sorted(p.get("index") for p in (ec.pages or [])):
        if prev_idx >= page_idx or prev_idx == 0:
            continue
        prev_text = next((p.get("scene_cn", "")[:80] for p in ec.pages if p.get("index") == prev_idx), "")
        if prev_text:
            prev_summary += f"P{prev_idx}: {prev_text}...\n"

    try:
        with st.spinner(f"Claude 正在为 {page_display_name(page_idx)} 写画面描述..."):
            new_scene = write_scene_cn(
                story_sentence=ec_page.get("text", ""),
                page_idx=page_idx,
                book_title=outline.title,
                level=outline.level,
                ip_age=int(outline.ip_age or 12),
                cast_descriptions=cast_descs,
                style_summary=blocks.get("positive_block", ""),
                must_include=entry.get("must_include", ""),
                must_avoid=entry.get("page_avoid", ""),
                previous_pages_summary=prev_summary,
            )
        ec_page["scene_cn"] = new_scene
        st.success(f"✅ Claude 已重写 {page_display_name(page_idx)} 的 scene_cn（{len(new_scene)} 字）")
    except DeepSeekError as e:
        st.error(f"❌ Claude 调用失败：{e}")
    except Exception as e:
        st.error(f"❌ 未预期错误：{e}")


def _deepseek_rewrite_all_scenes() -> None:
    """让 DeepSeek 重写全部 7 页 scene_cn。"""
    ec = st.session_state.extracted
    if not ec or not ec.pages:
        return
    progress = st.progress(0.0)
    total = len(ec.pages)
    for i, ec_page in enumerate(ec.pages, 1):
        idx = ec_page.get("index")
        try:
            _deepseek_rewrite_one_scene(idx)
        except Exception:
            pass
        progress.progress(i / total, text=f"Claude 写第 {i}/{total} 页 scene_cn...")
    progress.empty()
    st.success(f"🎉 Claude 已重写全部 {total} 页 scene_cn。建议下一步点 🔄 重建全部 prompt 让新 scene_cn 注入到 prompt。")


def _run_fact_check() -> None:
    """科普非虚构：对 ec.pages 的 text+scene_cn 跑科学事实校验，结果存入 session。"""
    ec = st.session_state.extracted
    outline = st.session_state.get("outline")
    if not ec or not ec.pages:
        return
    try:
        from fact_check import fact_check_pages
        with st.spinner("Claude 正在核查每页的科学事实、比例与图文一致..."):
            issues = fact_check_pages(
                ec.pages,
                title=getattr(outline, "title", "") if outline else "",
                level=getattr(outline, "level", "") if outline else "",
                fiction_type=(getattr(outline, "fiction_type", "") or getattr(outline, "reader_type", "")) if outline else "",
            )
        st.session_state["fact_issues"] = issues
        if not issues:
            st.success("✅ 科学事实校验：未发现问题。")
        else:
            st.warning(f"⚠️ 发现 {len(issues)} 处需要确认的问题，见下方。")
    except Exception as e:
        st.error(f"❌ 科学事实校验失败：{e}")


def _render_fact_check_results() -> None:
    """展示上次科学事实校验结果，并提供一键应用修正。"""
    issues = st.session_state.get("fact_issues")
    if not issues:
        return
    ec = st.session_state.extracted
    sev_icon = {"high": "🔴", "medium": "🟠", "low": "🟡"}
    with st.expander(f"🔬 科学事实校验结果（{len(issues)} 处）", expanded=True):
        for it in issues:
            loc = "Cover" if it["index"] == 0 else page_display_name(it["index"])
            st.markdown(
                f"{sev_icon.get(it['severity'], '🟠')} **{loc}**（{it['field']}）：{it['problem']}"
            )
            if it.get("suggestion"):
                st.caption(f"建议：{it['suggestion']}")
            if it.get("fixed_text"):
                st.caption(f"修正后文字：{it['fixed_text']}")
            if it.get("fixed_scene_cn"):
                st.caption(f"修正后画面：{it['fixed_scene_cn']}")
        if st.button("✅ 一键应用全部修正（写回文字/画面）", key="s4_apply_factfix", type="primary"):
            from fact_check import apply_fixes_to_ec_pages
            n = apply_fixes_to_ec_pages(ec.pages, issues)
            st.session_state["fact_issues"] = None
            st.success(f"已应用 {n} 处修正。建议重写/重建 prompt 让修正生效。")
            st.rerun()


def _deepseek_polish_one_prompt(page_idx: int) -> None:
    """让 DeepSeek 润色单页的正向 prompt。"""
    from scene_cn_writer import polish_image_prompt, DeepSeekError

    page_prompts = st.session_state.get("page_prompts") or {}
    entry = page_prompts.get(page_idx) or {}
    cur_pos = entry.get("positive", "")
    if not cur_pos:
        st.warning("当前页没有正向 prompt，请先点 🔄 重建本页")
        return

    ec = st.session_state.extracted
    ec_page = next((p for p in (ec.pages or []) if p.get("index") == page_idx), None)
    story_sent = ec_page.get("text", "") if ec_page else ""

    style_cfg = st.session_state.get("style_config") or {}
    blocks = _format_style_config_preview(style_cfg) if style_cfg else {"positive_block": ""}

    try:
        with st.spinner(f"Claude 正在润色 {page_display_name(page_idx)} 的 prompt..."):
            polished = polish_image_prompt(
                current_prompt=cur_pos,
                story_sentence=story_sent,
                style_summary=blocks.get("positive_block", ""),
                must_include=entry.get("must_include", ""),
                must_avoid=entry.get("page_avoid", ""),
            )
        entry["positive"] = polished
        page_prompts[page_idx] = entry
        st.session_state.page_prompts = page_prompts
        st.success(f"✅ Claude 已润色 {page_display_name(page_idx)} 的正向 prompt（{len(polished)} 字）")
    except DeepSeekError as e:
        st.error(f"❌ Claude 调用失败：{e}")
    except Exception as e:
        st.error(f"❌ 未预期错误：{e}")


def _cast_descriptions_from_pool(cast_pool: list[str], ip_age: int) -> list[str]:
    """从 cast_pool (IP key 列表) 转成完整的人物形象描述，喂给 DeepSeek。

    优先用 cn_prompt_builder._key_lock_phrase（已包含完整中文外观），
    fallback 到 ip_library 的 description 字段。
    """
    out: list[str] = []
    # 1) 优先：cn_prompt_builder._key_lock_phrase
    try:
        from cn_prompt_builder import _key_lock_phrase
        for key in cast_pool:
            try:
                phrase = _key_lock_phrase(key, ip_age)
                if phrase and phrase.strip():
                    out.append(phrase.strip())
            except Exception:
                continue
        if out:
            return out
    except ImportError:
        pass

    # 2) Fallback: ip_library
    try:
        from ip_library import get_ip
        for key in cast_pool:
            ip = get_ip(key)
            if ip and ip.get("description"):
                out.append(f"{key}：{ip['description']}")
    except ImportError:
        pass
    return out


def _render_step4_page_facts(step_num: int) -> None:
    exp = _locked_step_expander(step_num, "📖 分页事件编辑（7 页的故事文本 + 画面事实）")
    if exp is None:
        return
    with exp:
        st.caption(
            "💡 这一步**只编辑事实**：每页讲什么故事、画面里有什么、必须出现/必避免什么道具。"
            "提示词在下一步 Step 5 由这些事实自动重建。"
        )
        ec = st.session_state.extracted
        page_prompts = st.session_state.get("page_prompts") or {}
        ec_pages_by_idx = {p.get("index"): p for p in (ec.pages or [])}

        for idx in sorted(page_prompts.keys()):
            entry = page_prompts[idx]
            ec_page = ec_pages_by_idx.get(idx) or {}
            display_name = entry.get("display_name") or page_display_name(idx)
            label_kind = "封面" if idx == 0 else "故事页"

            st.markdown(f"#### {display_name} · {label_kind}")

            c1, c2 = st.columns([1, 1])
            with c1:
                st.caption("① 故事原文（英文）")
                new_text = st.text_area(
                    f"text_{idx}",
                    value=ec_page.get("text", ""),
                    height=120,
                    key=f"s4_text_{idx}",
                    label_visibility="collapsed",
                )
                if new_text != ec_page.get("text", ""):
                    ec_page["text"] = new_text
            with c2:
                st.caption("② 中文画面描述（120-220 字，AI 已生成可改）")
                new_scene_cn = st.text_area(
                    f"scene_cn_{idx}",
                    value=ec_page.get("scene_cn", ""),
                    height=120,
                    key=f"s4_scene_cn_{idx}",
                    label_visibility="collapsed",
                    placeholder=(
                        "AI 应该生成 120-220 字连贯描述：\n"
                        "主体（人物+具体外观）+ 动作（具体动词姿势）+ 环境（可见物品）+ 氛围（光照）"
                    ),
                )
                if new_scene_cn != ec_page.get("scene_cn", ""):
                    ec_page["scene_cn"] = new_scene_cn

            c3, c4 = st.columns([1, 1])
            with c3:
                st.caption("③ 必须出现（每行一条，会追加到正向提示词）")
                entry["must_include"] = st.text_area(
                    f"must_inc_{idx}",
                    value=entry.get("must_include", ""),
                    height=100,
                    key=f"s4_must_{idx}",
                    label_visibility="collapsed",
                    placeholder="Anna 头戴白色发箍、穿绿色毛衣\nAnna 黑色齐下巴bob短发\n桌上 5 本课本\n教室背景：暖米白空墙、单侧窗柔光",
                )
            with c4:
                st.caption("④ 必避免（每行一条，会追加到反向提示词）")
                entry["page_avoid"] = st.text_area(
                    f"page_avoid_{idx}",
                    value=entry.get("page_avoid", ""),
                    height=100,
                    key=f"s4_avoid_{idx}",
                    label_visibility="collapsed",
                    placeholder="本页不应有宠物\n不要 Tommy 戴眼镜\nMia 不能散发",
                )

            bc1, bc2 = st.columns([1, 3])
            with bc1:
                cur_shot = ec_page.get("shot", "medium") or "medium"
                new_shot = st.selectbox(
                    f"镜头_{idx}",
                    SHOT_OPTIONS,
                    index=max(0, SHOT_OPTIONS.index(cur_shot) if cur_shot in SHOT_OPTIONS else 1),
                    key=f"s4_shot_{idx}",
                )
                if new_shot != cur_shot:
                    ec_page["shot"] = new_shot
            with bc2:
                refs = entry.get("references", [])
                if refs:
                    st.caption(
                        "🎨 参考图: " + " | ".join(Path(r).name for r in refs[:4])
                    )

            st.divider()

        _confirm_next_step(
            step_num,
            label="✅ 确认分页事件，重建提示词",
            help_text="👉 确认后会基于新事实**重建所有 7 页提示词**，进入 Step 5 审核",
        )


# ============================================================================
# Step 5：📝 提示词预览（基于事实重建 + 单页正向/反向编辑）
# ============================================================================

def _render_step5_prompts(step_num: int) -> None:
    exp = _locked_step_expander(step_num, "📝 提示词预览与微调（正向 + 反向 双段）")
    if exp is None:
        return
    with exp:
        st.caption(
            "💡 提示词 = Step 3 风格设定 + Step 4 分页事实 自动组装。"
            "可以手动微调；如果改回事实，请回 Step 4。"
        )

        page_prompts = st.session_state.get("page_prompts") or {}

        col_rebuild, _ = st.columns([1, 3])
        with col_rebuild:
            if st.button(
                "🔄 基于最新事实/风格重建全部 prompt",
                type="secondary",
                key="s5_rebuild_all",
                width="stretch",
                help="把 Step 3 风格设定 + Step 4 分页事实 重新组装成提示词，会**覆盖**当前已编辑的 prompt",
            ):
                _rebuild_all_cn_prompts()
                st.success("✅ 已重建全部 prompt")
                st.rerun()

        for idx in sorted(page_prompts.keys()):
            entry = page_prompts[idx]
            display_name = entry.get("display_name") or page_display_name(idx)
            refs = entry.get("references", [])
            ref_badge = f"🖼️ {len(refs)} 张参考图" if refs else "⚠️ 无参考图"
            st.markdown(f"#### {display_name} · {ref_badge}")

            c3a, c3b = st.columns([3, 2])
            with c3a:
                st.caption("✅ 正向 Prompt（火山风单段流畅自然语言）")
                cur_pos = entry.get("positive", entry.get("prompt", ""))
                new_pos = st.text_area(
                    f"s5_positive_{idx}",
                    value=cur_pos,
                    height=200,
                    key=f"s5_pos_{idx}",
                    label_visibility="collapsed",
                )
                entry["positive"] = new_pos
            with c3b:
                st.caption("❌ 反向 Prompt（每行一条，不要什么）")
                cur_neg = entry.get("negative", "")
                new_neg = st.text_area(
                    f"s5_negative_{idx}",
                    value=cur_neg,
                    height=200,
                    key=f"s5_neg_{idx}",
                    label_visibility="collapsed",
                )
                entry["negative"] = new_neg

            with st.expander(f"🔍 {display_name} 最终组装效果（含 Step 3 全局设定）", expanded=False):
                final_pos = entry.get("positive", "")
                final_neg = entry.get("negative", "")
                style_cfg = st.session_state.get("style_config") or {}
                blocks = _format_style_config_preview(style_cfg) if style_cfg else {
                    "positive_block": "", "negative_block": ""
                }
                if blocks["positive_block"]:
                    final_pos = f"【全局风格】{blocks['positive_block']}\n\n{final_pos}"
                if blocks["negative_block"]:
                    final_neg = (final_neg.rstrip() + "\n" if final_neg else "") + blocks["negative_block"]
                must_inc = (entry.get("must_include") or "").strip()
                if must_inc:
                    final_pos = final_pos.rstrip() + "\n\n【教师锁定 · 必须出现】\n" + must_inc
                page_avoid = (entry.get("page_avoid") or "").strip()
                if page_avoid:
                    final_neg = (final_neg.rstrip() + "\n" if final_neg else "") + page_avoid

                st.markdown("**✅ 正向（最终）：**")
                st.code(final_pos, language="text")
                st.markdown("**❌ 反向（最终）：**")
                st.code(final_neg, language="text")

            st.divider()

        _confirm_next_step(
            step_num,
            label="✅ 提示词 OK，进入生图",
            help_text="👉 确认后进入生图步骤。生图时会自动用最终组装的提示词。",
        )


# ============================================================================
# Step 6：🖼️ 单页生图（Phase 1 保持当前 1 图，Phase 2 升级 3 候选）
# ============================================================================

def _lexile_display(outline) -> tuple[str, str]:
    """块11：返回 (展示值, 来源说明)。禁止编造——没有官方/人工值时显示 N/A 待核。"""
    if outline is None:
        return "—", "未知"
    lex = (getattr(outline, "lexile", "") or "").strip()
    src = (getattr(outline, "lexile_source", "") or "").strip()
    if lex and src in ("syllabus", "manual", "analyzer"):
        label = {"syllabus": "大纲官方", "manual": "人工核值", "analyzer": "官方分析器"}.get(src, src)
        return lex, label
    if lex:  # 有值但来源不明（如旧数据）→ 也展示但提示待核
        return lex, "待核实"
    return "N/A（待官方核值）", "缺失·禁编造"


def _render_lexile_panel(outline) -> None:
    """块11：大纲没有官方 Lexile 时，给老师【一键打开官方分析器 + 回填】的半自动取值（绝不编造）。"""
    if outline is None:
        return
    lex = (getattr(outline, "lexile", "") or "").strip()
    src = (getattr(outline, "lexile_source", "") or "").strip()
    has_real = bool(lex and src in ("syllabus", "manual", "analyzer"))
    with st.expander("📊 蓝思 Lexile 取值（禁止编造 · 每个值都要有依据）", expanded=not has_real):
        if has_real:
            st.success(f"✅ 当前 Lexile = **{lex}**（来源：{src}）。如需更正可在下方重新回填。")
        else:
            st.warning(
                "⚠️ 大纲未收录该书官方 Lexile。**不会自动编造**。请用官方分析器对正文取真实值后回填；"
                "未回填时 RR/TG 的 Lexile 显示为 `N/A（待官方核值）`。"
            )
        st.markdown(
            "🔗 打开官方分析器：[Lexile Text Analyzer](https://hub.lexile.com/text-analyzer/) "
            "（登录后粘贴本书英文正文 → 取得 Measure）"
        )
        new_val = st.text_input(
            "把官方分析器测得的 Lexile 回填到这里（如 450L / BR100L）",
            value=lex if src in ("manual", "analyzer") else "",
            key="lexile_manual_input",
            placeholder="例如 450L",
        )
        if st.button("💾 保存 Lexile（标记为人工核值）", key="lexile_save_btn"):
            v = (new_val or "").strip()
            if not v:
                st.error("请输入官方分析器测得的 Lexile 值，或留空保持 N/A。")
            else:
                outline.lexile = v
                outline.lexile_source = "manual"
                try:
                    st.session_state.auto["lexile"] = v
                except Exception:
                    pass
                st.success(f"✅ 已保存 Lexile = {v}（人工核值）。RR/TG 将使用该值。")
                st.rerun()


def _run_storyboard_drafts(mock_images: bool) -> None:
    """生图前【分镜草图】：为全部 8 页各出一张快速低分草图（1024²、不自审、不放大），
    供老师对照故事描述审稿。重生草图会清除"已确认"标记（需重新确认）。"""
    ec = st.session_state.extracted
    outline: BookOutline = st.session_state.outline
    apply_extracted_to_outline(outline, ec)
    if st.session_state.get("_syllabus_hit"):
        enrich_from_syllabus(outline)
        _sync_ec_from_syllabus(ec, outline)

    _run_dir, img_dir, _ = _ensure_run_dir()
    ip_age = outline.ip_age or resolve_ip_age(outline.level)
    pages = list(outline.pages)
    sb: dict = st.session_state.get("storyboard") or {}

    def _build(page):
        prompt, refs = _build_final_prompt_for_page(page, outline, ip_age)
        ver = int((sb.get(page.index) or {}).get("version", 0)) + 1
        dest = img_dir / f"draft_{page.index:02d}_v{ver}.png"
        return (page, prompt, list(refs)[:1], ver, dest)

    concurrency = max(1, int(os.getenv("IMAGE_CONCURRENCY", "2")))
    n = len(pages)
    progress = st.progress(0, "准备分镜 prompt…")
    status = st.empty()
    errors: dict[int, str] = {}
    tasks = []
    for i, p in enumerate(pages, 1):
        status.text(f"组装 prompt {i}/{n}…")
        progress.progress(i / (n + 1), f"组装 prompt {i}/{n}")
        tasks.append(_build(p))
    if mock_images:
        status.text("占位图模式：快速生成 8 张草图…")
    else:
        status.text(f"调用 gpt-image-2（{min(concurrency, n)} 张并发，首张约 30–60 秒）…")
    progress.progress(0, "生成分镜草图…")

    def _work(t):
        page, prompt, refs, ver, dest = t
        generate_image(prompt=prompt, dest=dest, references=refs,
                       mock=mock_images, label=f"draft P{page.index}",
                       size="1024x1024", deliver_print=False)
        return t

    done = 0
    with ThreadPoolExecutor(max_workers=min(concurrency, max(1, n))) as ex:
        futs = {ex.submit(_work, t): t for t in tasks}
        for fut in as_completed(futs):
            page, prompt, refs, ver, dest = futs[fut]
            done += 1
            try:
                fut.result()
                sb[page.index] = {"path": str(dest), "version": ver}
            except Exception as e:  # noqa: BLE001
                errors[page.index] = str(e)
            status.text(f"分镜草图 {done}/{n}…")
            progress.progress(done / n, f"分镜草图 {done}/{n}")

    st.session_state["storyboard"] = sb
    st.session_state["storyboard_confirmed"] = False  # 草图变了 → 必须重新确认
    progress.progress(1.0, "✅ 分镜草图完成")
    status.empty()
    if errors:
        st.warning("部分草图失败：" + "；".join(f"P{k}: {v[:40]}" for k, v in errors.items()))


def storyboard_ready() -> bool:
    """硬门：是否已确认全部分镜。未确认禁止正式出图。"""
    return bool(st.session_state.get("storyboard_confirmed"))


def _render_storyboard_panel(mock_images: bool) -> None:
    """生图前【分镜确认台】（用户拍板 2026-06-08 · 方案A+硬门）：
      ① 一键出 8 张快速草图分镜；
      ② 每页：草图 + 可编辑的中文故事描述（安全线）；
      ③ 硬门「✅ 确认全部分镜」——不确认禁止进入正式出图。
    """
    outline = st.session_state.get("outline")
    if outline is None or not getattr(outline, "pages", None):
        return
    try:
        from cn_prompt_builder import safety_line_default
    except Exception:
        return

    sb: dict = st.session_state.get("storyboard") or {}
    confirmed = storyboard_ready()

    title = "🎬 生图前·分镜确认台（先出 8 张草图 + 改故事描述 → 确认后才正式出图）"
    with st.expander(title, expanded=not confirmed):
        st.caption(
            "💡 流程：① 先出 8 张【快速草图】看分镜与故事描述对不对 → ② 逐页改中文描述（谁+在哪+做什么）"
            " → ③ 点「✅ 确认全部分镜」。**只有确认后才允许正式逐张出高清图。**"
            "草图为低分快稿（不放大、不自审），仅供审稿。"
        )
        if mock_images:
            st.caption("🟡 当前为占位图模式：分镜草图几秒内出完，不调 API。")
        else:
            st.caption("⏳ 未勾选占位图：分镜草图会调 gpt-image-2，8 张约 2–8 分钟，进度条走满前请耐心等待。")
        cda, cdb = st.columns([1, 3])
        with cda:
            if st.button("🎬 生成 / 重出 8 张分镜草图", type="primary", key="sb_gen_btn"):
                _run_storyboard_drafts(mock_images)
                st.rerun()
        with cdb:
            if not sb:
                st.info("还没有草图。先点左边「🎬 生成 8 张分镜草图」。")

        story_pages = [p for p in outline.pages]
        with st.form("storyboard_form"):
            edited: dict[int, str] = {}
            ncols = 2
            for row_start in range(0, len(story_pages), ncols):
                cols = st.columns(ncols)
                for j, p in enumerate(story_pages[row_start:row_start + ncols]):
                    with cols[j]:
                        st.markdown(f"**{page_display_name(p.index)}**")
                        rec = sb.get(p.index) or {}
                        dpath = rec.get("path")
                        if dpath and Path(dpath).exists():
                            st.image(dpath, use_container_width=True)
                        else:
                            st.caption("（暂无草图）")
                        is_cover = (p.page_type == "cover" or p.index == 0)
                        default = (p.text if is_cover else safety_line_default(p)) or ""
                        edited[p.index] = st.text_area(
                            "故事描述（可改）", value=default,
                            key=f"sb_desc_{p.index}", height=80,
                        )
            c1, c2 = st.columns([1, 1])
            with c1:
                ok = st.form_submit_button("✅ 确认全部分镜（写回 + 重建 prompt + 放行出图）",
                                           type="primary", width="stretch")
            with c2:
                save_only = st.form_submit_button("💾 仅保存描述（不确认）", width="stretch")

        if ok or save_only:
            for p in story_pages:
                if p.page_type == "cover" or p.index == 0:
                    continue
                p.safety_line = (edited.get(p.index, "") or "").strip()
            _rebuild_all_cn_prompts()
            if ok:
                st.session_state["storyboard_confirmed"] = True
                st.success("✅ 分镜已确认，提示词已重建。现在可以进入正式出图。")
            else:
                st.session_state["storyboard_confirmed"] = False
                st.info("已保存描述（未确认）。确认后才能正式出图。")
            st.rerun()

    if not confirmed:
        st.warning("⛔ 尚未确认分镜——「正式出图」按钮已锁定。请先出草图、改描述并点「✅ 确认全部分镜」。")
    else:
        st.success("✅ 分镜已确认，可正式出图。")


def _render_safety_line_panel() -> None:
    """块4（用户拍板 2026-06-08）：生图前逐页【简体中文场景安全线】确认。

    每页给老师一句"谁+在哪+做什么"的中文安全线（默认从 scene_cn 提炼），
    老师可逐页编辑；确认后写回 page.safety_line（注入 prompt 最前作为权威画面）并重建全部 prompt。
    """
    outline = st.session_state.get("outline")
    if outline is None or not getattr(outline, "pages", None):
        return
    try:
        from cn_prompt_builder import safety_line_default
    except Exception:
        return

    confirmed = st.session_state.get("safety_lines_confirmed", False)
    title = "🛡️ 生图前·场景安全线确认（逐页一句中文：谁+在哪+做什么）"
    with st.expander(title, expanded=not confirmed):
        st.caption(
            "💡 老师先读一遍每页画面要画什么，可直接改这句话。"
            "**确认后这句会作为该页画面的权威核心**注入提示词最前（图严格照它画），并重建全部 prompt。"
            "不确认也能出图（默认用 AI 原句）。"
        )
        story_pages = [p for p in outline.pages if (p.page_type != "cover" and p.index != 0)]
        with st.form("safety_line_form"):
            edited: dict[int, str] = {}
            for p in story_pages:
                default = safety_line_default(p)
                edited[p.index] = st.text_input(
                    f"{page_display_name(p.index)}",
                    value=default,
                    key=f"safety_line_{p.index}",
                )
            c1, c2 = st.columns([1, 1])
            with c1:
                ok = st.form_submit_button("✅ 确认并写回（重建 prompt）", type="primary", width="stretch")
            with c2:
                reset = st.form_submit_button("↩️ 重置为 AI 原句", width="stretch")
        if ok:
            for p in story_pages:
                p.safety_line = (edited.get(p.index, "") or "").strip()
            _rebuild_all_cn_prompts()
            st.session_state["safety_lines_confirmed"] = True
            st.success("✅ 场景安全线已确认并写回，提示词已重建。可以开始生图了。")
            st.rerun()
        if reset:
            for p in story_pages:
                p.safety_line = ""
            st.session_state["safety_lines_confirmed"] = False
            _rebuild_all_cn_prompts()
            st.info("已清空安全线，恢复使用 AI 原句。")
            st.rerun()

    if not confirmed:
        st.caption("ℹ️ 尚未确认场景安全线——可先确认再出图（也可直接出图，使用 AI 原句）。")


def _render_step6_image_gen(step_num: int, embed: bool = False) -> None:
    if embed:
        exp = contextlib.nullcontext()
    else:
        exp = _locked_step_expander(step_num, "🖼️ 单页生图（生成 → 审核 → 单图重生）")
        if exp is None:
            return
    with exp:
        st.caption(
            "💡 每页生 1 张图（gpt-image-2，主角恒为首位参考图保持一致），不满意可单页重生。"
        )

        _force_mock = not bool(JIMENG_API_KEY)
        mock_imgs = st.checkbox(
            "🟡 仅调试版式时勾选（不调用 gpt-image-2 API，用占位图代替）",
            value=_force_mock,
            disabled=_force_mock,
            help="**正式出图请保持不勾**。",
            key="s6_mock_imgs",
        )
        if mock_imgs:
            st.warning("⚠️ 当前为 **占位图模式**。要出正式 PPT 请先取消勾选。")

        # 生图前【分镜确认台】（方案A+硬门）：先出 8 张草图 + 改故事描述 → 确认后才放行正式出图。
        #   （内含逐页中文故事描述编辑，等价并取代原独立的场景安全线面板）
        _render_storyboard_panel(mock_imgs)
        _sb_ok = storyboard_ready()

        col_a, col_b = st.columns([4, 1])
        with col_a:
            n_pages = len(st.session_state.outline.pages) if st.session_state.outline else 8
            _conc = max(1, int(os.getenv("IMAGE_CONCURRENCY", "2")))
            # 真实耗时：单张 gpt-image-2 约 30-60 秒（模型固有），并发 _conc 张同时跑。
            _batches = -(-n_pages // _conc)  # 向上取整的批次数
            est_lo = _batches * 30
            est_hi = _batches * 60
            if not mock_imgs:
                st.info(
                    f"🎨 将调用 gpt-image-2 出 {n_pages} 张水彩绘本插画"
                    f"（{_conc} 张并发），约需 {est_lo//60}–{est_hi//60} 分钟。"
                    f"单张模型生成本身就要 30–60 秒，属正常。"
                )
        with col_b:
            if st.button("🎨 生成 / 重新生成所有图", type="primary",
                         width="stretch", key="s6_gen_btn", disabled=not _sb_ok,
                         help=None if _sb_ok else "需先在上方「分镜确认台」确认全部分镜才能正式出图"):
                _run_image_generation_only(mock_imgs)
        if not _sb_ok:
            st.info("ℹ️「正式出图」已锁定：请先在上方「🎬 分镜确认台」出草图、改描述并点「✅ 确认全部分镜」。")

        if st.session_state.get("image_results"):
            st.divider()
            _render_image_review_panel(mock_imgs)
            n_locked = sum(1 for r in st.session_state.image_results.values() if r.get("locked"))
            n_total = len(st.session_state.image_results)
            if n_locked < n_total:
                st.warning(f"还有 {n_total - n_locked} 张图未锁定 ✅。")
            else:
                st.success(f"🎉 全部 {n_total} 张图已锁定。")

            if not embed:
                _confirm_next_step(
                    step_num,
                    label="✅ 图都满意，进入组装",
                    help_text="👉 建议先把全部图勾 ✅ 锁定再进入下一步",
                )


# ============================================================================
# Step 7：📦 组装 4 件套
# ============================================================================

def _render_step7_assemble(step_num: int, wizard: bool = False) -> None:
    exp = _stage_shell(step_num, "📦 组装 4 件套（PPT / Worksheet / RR / TG）", wizard)
    if exp is None:
        return
    with exp:
        st.caption("💡 点击下方按钮，会用 Step 6 锁定的图 + Step 1-5 的数据，一键打包成 ZIP。")

        if not st.session_state.get("image_results"):
            st.error("⚠️ 还没有生图结果，请回 Step 6 先生图")
            return

        n_locked = sum(1 for r in st.session_state.image_results.values() if r.get("locked"))
        n_total = len(st.session_state.image_results)
        st.info(f"📊 当前 {n_locked}/{n_total} 张图已锁定")

        if st.button("📦 组装 4 件套 + 打包 ZIP", type="primary",
                     width="stretch", key="s7_assemble"):
            _run_docs_assembly()


def _format_page_constraints(pc: dict | None) -> str:
    """把 A 层 page_constraints 字段拼成一段可注入正向 prompt 的文字。

    自动跳过 🤖 字段（让 AI 自己推断）。
    """
    if not pc:
        return ""
    bits: list[str] = []
    pose = pc.get("pose", "")
    if pose and not pose.startswith("🤖"):
        bits.append(f"主角姿势：{pose}")
    emo = pc.get("emotion", "")
    if emo and not emo.startswith("🤖"):
        bits.append(f"情绪：{emo}")
    gaze = pc.get("gaze", "")
    if gaze and not gaze.startswith("🤖"):
        bits.append(f"视线：{gaze.split('（')[0]}")
    pos = pc.get("position", "")
    if pos and not pos.startswith("🤖"):
        bits.append(f"画面位置：{pos.split('（')[0]}")
    inter = pc.get("interaction", "")
    if inter and not inter.startswith("🤖"):
        bits.append(f"互动关系：{inter}")
    text_pos = pc.get("text_pos", "")
    if text_pos and not text_pos.startswith("🤖") and not text_pos.startswith("无"):
        bits.append(f"留白：{text_pos.split('（')[0]}")
    focus = (pc.get("focus_char") or "").strip()
    if focus:
        bits.append(f"焦点角色：{focus}（应当更大、更清晰、位置更接近中心）")
    return "；".join(bits)


def _format_style_config_preview(cfg: dict) -> dict:
    """把 6 字段拼成可注入提示词的两段（正向 + 反向）。"""
    pos_parts: list[str] = []
    # 视觉风格
    vs = cfg.get("visual_style", "")
    if vs:
        pos_parts.append(f"风格：{vs.split('（')[0]}")
    # 主场景（自动则不强制）
    scene = cfg.get("main_scene", "")
    if scene and not scene.startswith("🤖"):
        pos_parts.append(f"全本主场景：{scene}")
    # 色调
    tone = cfg.get("color_tone", "")
    if tone:
        pos_parts.append(f"色调：{tone.split('（')[0]}")
    # 构图
    comp = cfg.get("composition", "")
    if comp:
        pos_parts.append(f"构图：{comp.split('（')[0]}")
    # 光源
    light = cfg.get("light_source", "")
    if light:
        pos_parts.append(f"光源：{light.split('（')[0]}")
    # 留白
    if cfg.get("text_safe_zone"):
        pos_parts.append("画面顶部和底部 25% 留白给文字，主体居中偏下")
    # v3.2 B 层：视角
    view = cfg.get("view_angle", "")
    if view and not view.startswith("🤖"):
        pos_parts.append(f"视角：{view.split('（')[0]}")
    # v3.2 B 层：焦点处理
    focus = cfg.get("focus_handling", "")
    if focus:
        pos_parts.append(f"景深：{focus.split('（')[0]}")
    # 全局必出现
    must = (cfg.get("global_must") or "").strip()
    if must:
        items = [l.strip() for l in must.splitlines() if l.strip()]
        if items:
            pos_parts.append("全局必出现：" + "；".join(items))

    positive_block = "；".join(pos_parts)

    avoid = (cfg.get("global_avoid") or "").strip()
    avoid_items = [l.strip() for l in avoid.splitlines() if l.strip()]
    negative_block = "；".join(avoid_items)

    return {"positive_block": positive_block, "negative_block": negative_block}


# =============================== 主入口 ===============================
# ============================================================================
# v3.3：5 步工作流（合并 IP+画风、新增只读底层逻辑、合并提示词+生图为工作台）
# ============================================================================

def _render_step_ip_style(step_num: int, wizard: bool = False) -> None:
    """Step 2（合并旧 IP 锁定 + 画风设定）：先锁人，再定调。"""
    exp = _stage_shell(step_num, "🎭 人物 IP + 🎨 画风背景（先锁人，再定调）", wizard)
    if exp is None:
        return
    with exp:
        st.markdown("#### 🎭 A · 人物 IP 锁定")
        _render_step2_ip_lock(step_num, embed=True)
        st.divider()
        st.markdown("#### 🎨 B · 画风 / 背景设定（全局应用到所有页）")
        _render_style_panel_step(step_num, embed=True)
        _confirm_next_step(
            step_num,
            label="✅ 人物 + 画风都锁定，进入生图工作台",
            help_text="👉 IP 决定参考图、画风注入每页提示词；确认后进入生图工作台",
        )


def _render_step_rules(step_num: int) -> None:
    """Step 3（新）：底层逻辑 / 系统硬规则（只读，无需输入）。"""
    exp = _locked_step_expander(step_num, "📐 底层逻辑（系统自动套用的硬规则 · 只读）")
    if exp is None:
        return
    with exp:
        outline: BookOutline = st.session_state.outline
        level_key = (outline.level_key if outline else "5")
        is_dual = level_key in ("smart", "0", "1", "2")
        rr_dist = rr_question_distribution(outline.level if outline else "5")
        star = " + ".join(
            f"{rr_dist.count(s)}×{'⭐' * s}" for s in sorted(set(rr_dist))
        )
        st.caption("💡 以下规则**全自动套用**到所有交付物，无需你输入。列在这里是为了让你心里有数。")

        cp = COMPOSITION_POLICY
        st.markdown(
            "**🎯 画面构图 / 比例（AI 按故事自动决定，自动注入每页提示词）**\n"
            f"- 主角是**画面唯一视觉中心**，占画面 **{cp['protagonist_pct']}**（清晰饱满）\n"
            f"- 同框人物按**真实身高比例**（谁都不能比同框同龄人明显大一圈）\n"
            f"- 动物按**真实比例**（仓鼠≈手掌大，不会画成猫狗大小）\n"
            f"- 背景占 **{cp['background_pct']}**，环境清晰但不喧宾夺主\n"
            f"- 视角：{cp['perspective']}　·　画风：{cp['style']}"
        )
        st.markdown(
            "**🧼 画面平滑 / 统一（自动注入每页正向 + 反向提示词）**\n"
            "- 整体**干净、平滑、统一**，强调**大色块叙事**与整体轮廓\n"
            "- **不要细碎噪点 / 高频纹理 / 脏污颗粒 / 密集小装饰**\n"
            "- 边缘清晰利落、表面干净、画面呼吸感强、一目了然"
        )
        st.markdown(
            "**🎭 系列默认 IP + 全本一致性（自动套用）**\n"
            "- 系列固定主角 **Mia（女孩）/ Tommy（男孩）**\n"
            "- 故事**没出现命名主角时**：女孩默认套 **Mia** 形象、男孩默认套 **Tommy** 形象\n"
            "- 一旦定下角色形象，**发型 / 服装 / 配饰 / 五官每页保持一致**，绝不跨页跳变"
        )
        st.divider()

        c1, c2 = st.columns(2)
        with c1:
            st.markdown(
                "**🖼️ 绘本图片**\n"
                "- 比例 **3:2 横版**（1536×1024，gpt-image-2）\n"
                "- 主角**恒为首位参考图**，保证全本不跳帧\n"
                "- 每页预留 **~20% 浅色区**给文字，模型自动选边\n"
                "- 页数 = **封面 + 7 故事 = 8 页**（PPT 补足 4 的倍数）"
            )
            st.markdown(
                "**📝 词汇显示**\n"
                + ("- 当前 **L0-L2 双行**：Mastery + Exposure\n"
                   if is_dual else
                   "- 当前 **L3-L6 单行**：Vocabulary 4 词\n")
                + "- 全部 **lemma 原型 + 小写 + 无标点**"
            )
        with c2:
            st.markdown(
                "**📄 文档规格**\n"
                "- Reading Report：**A4 竖版**，控制在 **1 页**\n"
                f"- RR 阅读表达题：**{len(rr_dist)} 题**（{star}），标 (P×) 出处\n"
                "- Worksheet：固定 **2 词汇 + 2 句子 + 2 阅读** 方向\n"
                "- 阅读字数 = **纯故事字数**（不含题目）"
            )
            st.markdown(
                "**📦 命名 / 打包**\n"
                "- 文件名：`Level X_BookXX_品类_标题.后缀`\n"
                "- 字体：标题 Poppins 40pt / 副标题 22pt\n"
                "- 一键打包 **绘本+练习册+RR+TG → ZIP**"
            )
        _confirm_next_step(
            step_num,
            label="✅ 我了解这些规则，进入 IP + 画风锁定",
            help_text="👉 这些是系统硬规则（含构图/比例），自动生效；确认后进入 IP + 画风锁定",
        )


def _render_step_workbench(step_num: int, wizard: bool = False) -> None:
    """Step 4（合并旧 分页提示词 + 单页生图）：生图工作台。"""
    exp = _stage_shell(
        step_num, "🖼️ 生图工作台（一键出 8 图 → 图旁看故事 → 逐页重生/锁定）", wizard
    )
    if exp is None:
        return
    with exp:
        st.caption(
            "💡 **图在前**：先一键出 8 张图；每张图旁边直接显示「该页故事原文 + 剧情要点」，"
            "提示词默认折叠。不满意就点 🔁 重生 / 展开提示词改了再出（像即梦逐页调）。"
        )
        # 🎨 主区：生成 + 逐页审核（图大、文字在旁、提示词折叠）
        _render_step6_image_gen(step_num, embed=True)

        st.divider()
        # ✍️ 次要：分页提示词总览（默认收起；用 checkbox 而非 expander，避免与内部 expander 嵌套）
        if st.checkbox(
            "✍️ 展开高级：分页提示词总览（批量查看 / Claude 润色，一般不用动）",
            value=False,
            key="wb_show_prompt_overview",
        ):
            _render_step4_combined(step_num, embed=True)

        _confirm_next_step(
            step_num,
            label="✅ 图都满意，进入组装",
            help_text="👉 建议先把全部图勾 ✅ 锁定再进入组装",
        )


def _render_batch_mode() -> None:
    """📚 批量生产：N 个大纲 → N×4 件套。详细实现见 batch_runner。"""
    st.subheader("📚 批量生产（输入 N 个大纲 → 每本自动产出 4 件套）")
    st.caption(
        "每本之间数据严格隔离；绘本图全自动生成（标记『待人工抽查』，可事后回单本模式逐页重生）。"
    )
    st.info(
        "粘贴多本大纲，每本用 `===` 分隔。每本第一行 = `Title | Level | Book#`，其后为故事原文。"
    )
    st.text_area(
        "批量大纲",
        height=220,
        key="batch_outlines_raw",
        placeholder=(
            "Spring Days | L1 | 02\n"
            "Mia sees a flower. ...\n"
            "===\n"
            "What Makes a Good Friend? | L5 | 01\n"
            "Anna felt nervous ...\n"
        ),
    )
    c1, c2, c3, c4 = st.columns(4)
    with c1:
        st.number_input("本级并发", min_value=1, max_value=8, value=2, key="batch_concurrency",
                        help="同时跑几本（本与本之间隔离）。")
    with c2:
        st.number_input("出图全局并发", min_value=1, max_value=8, value=4,
                        key="batch_image_concurrency",
                        help="① 所有本子在途出图请求的总上限，贴 API RPM 设置——提速的关键钮。")
    with c3:
        st.selectbox("输出方式", ["每本子文件夹", "平铺 + 规范命名"], key="batch_output_mode")
    with c4:
        st.checkbox("打包 ZIP", value=True, key="batch_zip")

    c5, c6 = st.columns([1, 2])
    with c5:
        st.checkbox("② 断点续跑（已出图的页不重做）", value=False, key="batch_resume")
    with c6:
        st.text_input("♻️ 重跑失败本：填上次 batch_log.json 路径（留空=全量新跑）",
                      key="batch_rerun_log", placeholder=r"outputs\batch_YYYYmmdd_HHMMSS\batch_log.json")

    b1, b2 = st.columns([1, 1])
    with b1:
        if st.button("🔍 预检 Dry-run（不花 API）", width="stretch"):
            try:
                from batch_runner import preflight_from_ui
                preflight_from_ui()
            except Exception as e:
                st.error(f"预检失败：{e}")
    with b2:
        if st.button("🚀 开始批量生产", type="primary", width="stretch"):
            try:
                from batch_runner import run_batch_from_ui
                run_batch_from_ui()
            except Exception as e:
                st.error(f"批量生产启动失败：{e}")


def _extract_uploaded_page_images(uploaded_files, dest_dir: Path) -> list[Path]:
    """块5：把上传的成品绘本（PDF / PPTX / 散图）抽成有序页图 PNG 列表。

    - 图片：直接落盘（按文件名排序）。
    - PDF：用 PyMuPDF 逐页渲染成图。
    - PPTX：取每页幻灯片里面积最大的图片（适配本系统自产绘本）。
    """
    dest_dir.mkdir(parents=True, exist_ok=True)
    out: list[Path] = []
    img_seq = sorted(uploaded_files, key=lambda f: f.name.lower())
    n = 0
    for f in img_seq:
        name = (f.name or "").lower()
        data = f.getvalue()
        if name.endswith(".pdf"):
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(stream=data, filetype="pdf")
                for pi in range(doc.page_count):
                    pg = doc.load_page(pi)
                    pix = pg.get_pixmap(matrix=fitz.Matrix(2, 2))
                    n += 1
                    p = dest_dir / f"page_{n:02d}.png"
                    pix.save(str(p))
                    out.append(p)
                doc.close()
            except Exception as e:
                st.warning(f"PDF 解析失败（{f.name}）：{e}")
        elif name.endswith(".pptx"):
            try:
                from pptx import Presentation
                from pptx.util import Emu
                import io as _io
                prs = Presentation(_io.BytesIO(data))
                for slide in prs.slides:
                    best = None
                    best_area = -1
                    for shape in slide.shapes:
                        if shape.shape_type == 13 or getattr(shape, "image", None) is not None:  # PICTURE
                            try:
                                area = int(shape.width) * int(shape.height)
                            except Exception:
                                area = 0
                            if area > best_area:
                                best_area = area
                                best = shape
                    if best is not None:
                        try:
                            blob = best.image.blob
                            ext = best.image.ext or "png"
                            n += 1
                            p = dest_dir / f"page_{n:02d}.{ext}"
                            p.write_bytes(blob)
                            out.append(p)
                        except Exception:
                            pass
            except Exception as e:
                st.warning(f"PPTX 解析失败（{f.name}）：{e}")
        elif name.endswith((".png", ".jpg", ".jpeg", ".webp")):
            n += 1
            p = dest_dir / f"page_{n:02d}.png"
            try:
                from PIL import Image as _Image
                import io as _io
                im = _Image.open(_io.BytesIO(data)).convert("RGB")
                im.save(p, "PNG")
                out.append(p)
            except Exception:
                p2 = dest_dir / f"page_{n:02d}{Path(name).suffix}"
                p2.write_bytes(data)
                out.append(p2)
    return out[:8]


def _render_upload_mode() -> None:
    """块5（用户拍板 2026-06-08）：上传成品绘本(PDF/PPT/散图) → 只生成 Worksheet + RR + TG（不重出绘本本体）。"""
    st.subheader("📤 上传成品绘本 → 生成教辅（Worksheet + Reading Report + Teacher's Guide）")
    st.caption(
        "💡 已有成品绘本时用这里：上传 **PDF / PPTX / 散图**（最多 8 张：封面+7 页，允许缺页）。"
        "系统**不重出绘本本体**，只产出 WS + RR + TG。单词/拼读 **逐字取自该 Level 大纲**。"
    )

    c1, c2 = st.columns([3, 1])
    with c1:
        up_title = st.text_input("📕 Book Title *", value="", key="up_title",
                                 help="用于匹配大纲、文档大标题、文件名。")
    with c2:
        up_level = st.selectbox("🎚️ Level *", LEVEL_OPTIONS, index=4, key="up_level")

    files = st.file_uploader(
        "上传成品绘本（PDF / PPTX / 多张图片）",
        type=["pdf", "pptx", "png", "jpg", "jpeg", "webp"],
        accept_multiple_files=True,
        key="up_files",
    )
    up_text = st.text_area(
        "📝 故事英文原文（强烈建议粘贴：用于出阅读题/正文；留空则尽量用大纲）",
        height=160, key="up_raw_text",
        placeholder="把这本绘本的英文正文粘进来，用于生成阅读理解题与 RR 正文。",
    )
    rr_answers = st.checkbox("RR 生成示例答案版（教师版）", value=False, key="up_rr_answers")

    if not st.button("🚀 生成教辅三件套（WS + RR + TG）", type="primary", key="up_run_btn"):
        return

    if not up_title.strip():
        st.error("请填写 Book Title。")
        return
    if not files:
        st.error("请至少上传 1 个文件（PDF / PPTX / 图片）。")
        return

    run_dir, img_dir, name_prefix = _ensure_run_dir()
    with st.spinner("解析上传文件，抽取页图…"):
        image_paths = _extract_uploaded_page_images(files, img_dir)
    if not image_paths:
        st.error("没有从上传文件里解析出任何页图。请确认文件内容（PDF/PPTX 是否含图）。")
        return
    st.success(f"✅ 已解析出 {len(image_paths)} 张页图。")

    ip_age = resolve_ip_age(up_level)
    raw = (up_text or "").strip()
    # 1) 构建 outline（有正文→AI 抽题/拆页；无正文→空页骨架，靠大纲补词/题）
    outline = _build_outline(
        title=up_title, level=up_level, book_number="", cefr="", theme="",
        ip_age=int(ip_age), raw_story=raw, custom_chars_text="",
    )
    ec = None
    if raw:
        try:
            ec = run_with_live_timer("AI 抽取（抽词 · 拆段 · 出题）", extract_all,
                                     raw_story=raw, title=up_title, level=up_level,
                                     cefr="", theme="")
            apply_extracted_to_outline(outline, ec)
        except Exception as e:
            st.warning(f"AI 抽取失败，改用大纲/兜底：{e}")
    # 2) 大纲权威逐字覆盖词表/拼读（块6）
    hit = enrich_from_syllabus(outline)
    if hit:
        st.info("📚 命中 S&S 大纲：单词/拼读已按大纲逐字覆盖。")
    elif not raw:
        st.warning("⚠️ 未命中大纲且未粘贴正文 —— 题目/词表可能为占位，建议粘贴正文或确认书名与大纲一致。")
    # 3) 挂题
    if ec is not None:
        attach_rr_questions(outline, ec.rr_questions)
        rqc = int(st.session_state.get("ws_reading_q_count", 4))
        attach_worksheet_questions(outline, ec.worksheet_questions, reading_q_count=rqc)

    # 4) 组装 WS + RR + TG（不出绘本本体）
    progress = st.progress(0, "组装教辅中…")
    ws_path = run_dir / f"{name_prefix}_Worksheet.pptx"
    rr_path = run_dir / f"{name_prefix}_Reading_Report.docx"
    tg_path = run_dir / f"{name_prefix}_Teachers_Guide.docx"
    try:
        progress.progress(1 / 3, "生成 Worksheet…")
        build_worksheet(outline, ws_path, image_paths=image_paths)
        progress.progress(2 / 3, "生成 Reading Report…")
        build_reading_report(outline, rr_path, with_answers=rr_answers)
        progress.progress(3 / 3, "生成 Teacher's Guide…")
        build_teacher_guide(outline, tg_path)
    except Exception as e:
        st.error(f"组装失败：{e}")
        return

    zip_path = run_dir / f"{name_prefix}_Teaching_Kit.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:
        for p in (ws_path, rr_path, tg_path):
            if p.exists():
                z.write(p, arcname=p.name)
        for img in image_paths:
            z.write(img, arcname=f"images/{img.name}")

    progress.progress(1.0, "完成 ✅")
    st.success("✅ 教辅三件套已生成。")
    with open(zip_path, "rb") as fh:
        st.download_button("⬇️ 下载 ZIP（WS + RR + TG + 页图）", fh.read(),
                           file_name=zip_path.name, mime="application/zip")


def _render_global_standards_panel() -> None:
    """Section 0：全局底层逻辑（只读）。放在最前，让老师一眼看清今天的硬标准。

    数据来自 config 单一数据源（render_global_standards_md），与各 builder 同源。
    """
    outline = st.session_state.get("outline")
    level = outline.level if outline else None
    with st.expander("📐 全局底层逻辑（只读 · 所有交付物 100% 强制套用）", expanded=(outline is None)):
        st.caption("💡 这些是系统**全自动套用**的硬标准，无需你输入。批量生产时同样对每一本生效。")
        st.markdown(render_global_standards_md(level))
        st.divider()
        st.markdown("##### 五、4 大交付物尺寸 / 规格（点开看每件细则）")
        cols = st.columns(4)
        for i, key in enumerate(["book", "worksheet", "rr", "tg"]):
            spec = DELIVERABLE_SPECS[key]
            with cols[i]:
                with st.popover(f"{spec['icon']} {spec['name']}"):
                    st.markdown(render_deliverable_spec_md(key, level))


_DELIVERABLES = [
    ("book", "📖 绘本"),
    ("ws", "📝 练习"),
    ("rr", "📄 RR"),
    ("tg", "👩‍🏫 TG"),
]
_KIT_LABEL = "绘本 · 练习 · RR · TG"


def _render_deliverable_nav() -> str:
    """主区顶部：4 交付物横向切换。抽取完成前不显示，默认返回 'book'。"""
    if st.session_state.get("extracted") is None:
        return "book"
    labels = [lbl for _, lbl in _DELIVERABLES]
    st.markdown("<div class='nav-title-inline'>🗂️ 交付物</div>", unsafe_allow_html=True)
    sel = st.radio(
        "交付物导航", labels,
        horizontal=True,
        label_visibility="collapsed",
        key="deliverable_nav",
    )
    for k, lbl in _DELIVERABLES:
        if lbl == sel:
            return k
    return "book"


def _render_evals_panel() -> None:
    """一键体检（evals）— 对当前 outline 跑纯规则检查，出红/黄/绿报告。"""
    st.caption("对当前故事跑规则检查：词汇 lemma/小写/专有名词、绘本 IP+年龄、Worksheet 结构、RR 星级。")
    if not st.button("运行体检", width="stretch", key="run_evals_btn"):
        return
    outline = st.session_state.get("outline")
    ec = st.session_state.get("extracted")
    if not outline:
        st.warning("请先完成 AI 抽取，再运行体检。")
        return
    try:
        from evals import run_all, OK, WARN, ERROR
        ws = getattr(ec, "worksheet_questions", None) if ec else None
        rr = (getattr(ec, "reading_report", None)
              or getattr(ec, "rr_items", None)) if ec else None
        report = run_all(
            outline=outline,
            worksheet_questions=ws,
            rr_items=rr,
            cast_pool=st.session_state.get("story_cast_pool") or None,
            generic_overrides=st.session_state.get("generic_overrides") or None,
        )
        if report.passed and report.n_warn == 0:
            st.success("✅ 全部通过，无问题")
        elif report.passed:
            st.info(f"✅ 无硬性错误，{report.n_warn} 条提醒")
        else:
            st.error(f"❌ {report.n_error} 条硬性问题需修，{report.n_warn} 条提醒")
        for it in report.issues:
            if it.level == ERROR:
                st.markdown(f"❌ **[{it.category}]** {it.msg}")
            elif it.level == WARN:
                st.markdown(f"⚠️ [{it.category}] {it.msg}")
            else:
                st.markdown(f"✅ <span style='color:#888'>[{it.category}] {it.msg}</span>",
                            unsafe_allow_html=True)
    except Exception as e:
        st.warning(f"体检执行出错：{e}")


_DINO_ICON = BRAND_DIR / "dino_head_icon.png"
_KIDDE_HERO = BRAND_DIR / "dino_reading_logo.png"
_KIDDE_FALLBACK = Path(__file__).resolve().parent.parent / "assets" / "ip_library" / "dino.png"


def _img_b64(path: Path) -> str:
    """读取品牌图并转 base64（嵌入 HTML）。"""
    try:
        if path.exists():
            return base64.b64encode(path.read_bytes()).decode("ascii")
    except Exception:
        pass
    return ""


def _icon_b64() -> str:
    """Kidde 小图标（优先官方 reading logo / ip_library dino）。"""
    return _img_b64(_KIDDE_HERO) or _img_b64(_KIDDE_FALLBACK) or _img_b64(_DINO_ICON)


def _hero_kidde_b64() -> str:
    """Hero 区大号 Kidde 形象（VIPKID 官方黄色恐龙，非通用绿恐龙）。"""
    return _img_b64(_KIDDE_HERO) or _img_b64(_KIDDE_FALLBACK)


def _render_hero() -> None:
    """Hero：AI 赋能 badge + 标题 + 三特性 pill + Kidde 官方恐龙图。"""
    kidde_b64 = _hero_kidde_b64()
    kidde_img = (
        f"<img src='data:image/png;base64,{kidde_b64}' class='hero-kidde' alt='Kidde'/>"
        if kidde_b64
        else "<span class='hero-kidde-fallback'>🦕</span>"
    )
    st.markdown(
        f"""
        <section class='hero-section hero-gradient'>
          <div class='hero-grid'>
            <div class='hero-copy'>
              <div class='hero-pill'>
                <span class='hero-pill-icon'>✨</span>
                <span>AI 赋能 线下教学</span>
              </div>
              <h1 class='hero-headline'>
                VIPKID Dino<br/><span class='hero-accent'>线下绘本教学</span>
              </h1>
              <p class='hero-lead'>
                革新化的绘本教学流程：从 AI 智能抽取内容，到教师个性化精修，
                最后实现 PPT、练习册与教案的一键生成。
              </p>
              <div class='hero-features'>
                <div class='hero-feat'>
                  <span class='feat-ic feat-ic-teal'>🧠</span>
                  <div><p class='feat-t'>AI Extraction</p><p class='feat-s'>智能识别分级内容</p></div>
                </div>
                <div class='hero-feat'>
                  <span class='feat-ic feat-ic-orange'>✏️</span>
                  <div><p class='feat-t'>Fine-tuning</p><p class='feat-s'>灵活的人工微调</p></div>
                </div>
                <div class='hero-feat'>
                  <span class='feat-ic feat-ic-blue'>🚀</span>
                  <div><p class='feat-t'>One-click PPT</p><p class='feat-s'>课件资源秒速生成</p></div>
                </div>
              </div>
            </div>
            <div class='hero-visual'>
              <div class='hero-glow'></div>
              <div class='hero-glass'>
                {kidde_img}
                <div class='hero-glass-cap'>
                  <div>
                    <p class='cap-t'>Interactive Learning</p>
                    <p class='cap-s'>Dino series interactive curriculum</p>
                  </div>
                  <span class='cap-play'>▶</span>
                </div>
              </div>
            </div>
          </div>
        </section>
        """,
        unsafe_allow_html=True,
    )


_FRAMEWORK_DIR = OUTPUTS_DIR / "_framework"
_CMAP_PDF = _FRAMEWORK_DIR / "课程地图_L0-L6.pdf"
_CMAP_HTML = _FRAMEWORK_DIR / "课程地图_L0-L6.html"
_CMAP_XLSX = _FRAMEWORK_DIR / "课程对标总表_L0-L6.xlsx"
_CMAP_LONGIMG = _FRAMEWORK_DIR / "课程级别长图_L0-L6.png"
_CMAP_PREVIEW = _FRAMEWORK_DIR / "_preview_onepager.png"
_READING_LOGO = BRAND_DIR / "dino_reading_logo.png"
_BUNDLED_CURRICULUM = Path(__file__).resolve().parent.parent / "assets" / "curriculum"


def _seed_curriculum_bundle(*, force: bool = False) -> None:
    """从仓库内置资料复制到 outputs（Streamlit Cloud 无 Playwright，靠此同步本地长图/Excel）。"""
    if not _BUNDLED_CURRICULUM.is_dir():
        return
    for name, dest in (
        ("课程对标总表_L0-L6.xlsx", _CMAP_XLSX),
        ("课程地图_L0-L6.html", _CMAP_HTML),
        ("课程地图_L0-L6.pdf", _CMAP_PDF),
        ("课程级别长图_L0-L6.png", _CMAP_LONGIMG),
    ):
        src = _BUNDLED_CURRICULUM / name
        if src.exists() and (force or not dest.exists()):
            dest.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(src, dest)


def _ensure_curriculum_assets(*, force: bool = False) -> None:
    """按需生成 Excel / HTML / PDF / 竖版长图（与 build_curriculum_xlsx.DATA 同源）。"""
    if st.session_state.get("_cmap_build_tried") and not force:
        return
    st.session_state["_cmap_build_tried"] = True
    st.session_state.pop("_cmap_build_err", None)
    _FRAMEWORK_DIR.mkdir(parents=True, exist_ok=True)
    _seed_curriculum_bundle(force=force)
    errs: list[str] = []
    try:
        if force or not _CMAP_XLSX.exists():
            from build_curriculum_xlsx import build as build_xlsx
            build_xlsx()
    except Exception as e:
        errs.append(f"Excel: {e}")
    try:
        if force or not _CMAP_HTML.exists():
            from build_curriculum_onepager import build as build_html
            build_html()
    except Exception as e:
        errs.append(f"HTML/PDF: {e}")
    if force or not _CMAP_LONGIMG.exists():
        try:
            from build_curriculum_longimg import build as build_longimg
            build_longimg()
        except Exception as e:
            if not _CMAP_LONGIMG.exists():
                errs.append(str(e))
    if errs and not (_CMAP_XLSX.exists() and _CMAP_LONGIMG.exists()):
        st.session_state["_cmap_build_err"] = "；".join(errs)


def _render_curriculum_map() -> None:
    """0–6 课程地图：迷你梯度图（HTML·永远显示）+ 展开区下载长图/PDF/Excel（与教研总表同源）。"""
    logo_b64 = ""
    try:
        logo_b64 = base64.b64encode(_READING_LOGO.read_bytes()).decode("ascii")
    except Exception:
        pass
    logo_html = (f"<img src='data:image/png;base64,{logo_b64}' style='height:30px'/>"
                 if logo_b64 else "<span style='font-size:22px'>🦕</span>")
    st.markdown(render_mini_map_html(logo_html), unsafe_allow_html=True)

    with st.expander(
        "🗺️ 详细课程级别介绍（竖版长图 · 一页纸 PDF · 教研 Excel · 完整维度表）",
        expanded=False,
    ):
        _ensure_curriculum_assets()
        if _CMAP_LONGIMG.exists():
            st.image(str(_CMAP_LONGIMG), use_container_width=True,
                     caption="竖版长图：级别为列、维度为行（与北美外教阅读课对标表同源，适合分享）")
        elif _CMAP_PREVIEW.exists():
            st.image(str(_CMAP_PREVIEW), use_container_width=True,
                     caption="对外一页纸预览（A4 横向）")
        else:
            st.markdown(render_mini_map_html(logo_html), unsafe_allow_html=True)
            st.caption("竖版长图尚未生成；上方为 HTML 课程地图。可下载下方 PNG 或点「生成完整资料包」。")
            st.image(mini_map_png_bytes(), use_container_width=True,
                     caption="课程迷你梯度图 PNG（已内置 Noto 中文字体，云端可正常显示）")

        c0, c1, c2, c3, c4 = st.columns(5)
        with c0:
            st.download_button(
                "⬇️ 迷你梯度图 PNG",
                mini_map_png_bytes(),
                file_name="VIPKID_Dino_课程地图迷你图_L0-L6.png",
                mime="image/png",
                width="stretch",
                key="dl_mini_map_png",
            )
        with c1:
            if _CMAP_LONGIMG.exists():
                st.download_button("⬇️ 竖版长图 PNG", _CMAP_LONGIMG.read_bytes(),
                                   file_name="VIPKID_Dino_课程级别长图_L0-L6.png",
                                   mime="image/png", width="stretch")
        with c2:
            if _CMAP_PDF.exists():
                st.download_button("⬇️ 一页纸 PDF", _CMAP_PDF.read_bytes(),
                                   file_name="VIPKID_Dino_课程地图_L0-L6.pdf",
                                   mime="application/pdf", width="stretch")
        with c3:
            if _CMAP_XLSX.exists():
                st.download_button("⬇️ Excel 对标总表", _CMAP_XLSX.read_bytes(),
                                   file_name="VIPKID_Dino_课程对标总表_L0-L6.xlsx",
                                   mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                   width="stretch")
        with c4:
            if _CMAP_HTML.exists():
                st.download_button("⬇️ 网页版 HTML", _CMAP_HTML.read_bytes(),
                                   file_name="VIPKID_Dino_课程地图_L0-L6.html",
                                   mime="text/html", width="stretch")
            elif st.button("🔄 生成完整资料包", key="build_cmap_btn", width="stretch"):
                st.session_state.pop("_cmap_build_tried", None)
                _ensure_curriculum_assets(force=True)
                st.rerun()

        err = st.session_state.get("_cmap_build_err")
        if err:
            st.caption(f"⚠️ 完整资料包生成部分失败（迷你图 PNG 始终可下）：{err}")
        st.caption(
            "ᴿ = 对标参考值（依 CEFR 推导，可校准）；其余为官方 S&S / TG 权威值。"
            "长图 / PDF / Excel / HTML 四份内容同源。"
        )


def _sec_head(icon: str, title: str) -> None:
    """统一区块标题：图标徽标 + 标题 + 渐隐细线（字体层级清晰、对齐有序）。"""
    st.markdown(
        f"<div class='sec-head'><span class='ic'>{icon}</span>{title}<span class='line'></span></div>",
        unsafe_allow_html=True,
    )


def _chips(items, kind: str = "") -> None:
    """关键词标签行（chip/战格）。kind ∈ {'', 'gray', 'ok'}。"""
    cls = ("chip " + kind).strip()
    spans = "".join(f"<span class='{cls}'>{x}</span>" for x in items if x)
    st.markdown(f"<div class='chip-row'>{spans}</div>", unsafe_allow_html=True)


_MAIN_NAV = {
    "overview": "概览",
    "background": "背景",
    "features": "功能",
    "onboarding": "新手引导",
    "metrics": "指标",
    "work": "开始制作",
    "settings": "设置",
}
_LEGACY_NAV = {"guide": "overview", "faq": "onboarding"}
_AUTH_COOKIE = "dino_auth"
_NAV_COOKIE = "dino_tab"
_AUTH_MSG = b"vipkid-dino-auth-v1"


def _secret_str(key: str) -> str:
    try:
        return str(st.secrets.get(key, "")).strip()
    except Exception:
        return ""


def _parse_users_text(text: str) -> dict[str, str]:
    """解析 JSON 或多行 user:pass 格式。"""
    text = (text or "").strip()
    if not text:
        return {}
    if text.startswith("{"):
        try:
            obj = json.loads(text)
            if isinstance(obj, dict):
                return {
                    str(k).strip(): str(v).strip()
                    for k, v in obj.items()
                    if str(k).strip() and str(v).strip()
                }
        except json.JSONDecodeError:
            pass
    users: dict[str, str] = {}
    for line in text.splitlines():
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        if ":" not in line:
            continue
        user, pwd = line.split(":", 1)
        user, pwd = user.strip(), pwd.strip()
        if user and pwd:
            users[user] = pwd
    return users


def _get_app_users() -> dict[str, str]:
    """返回 username → password。未配置任何账户时返回空 dict（不启用登录门）。"""
    users: dict[str, str] = {}
    try:
        raw = st.secrets.get("APP_USERS")
        if isinstance(raw, Mapping):
            users = {
                str(k).strip(): str(v).strip()
                for k, v in raw.items()
                if str(k).strip() and str(v).strip()
            }
    except Exception:
        pass
    if not users:
        users = _parse_users_text(_secret_str("APP_USERS") or os.getenv("APP_USERS", ""))
    if not users:
        pwd = _secret_str("APP_PASSWORD") or os.getenv("APP_PASSWORD", "").strip()
        if pwd:
            user = _secret_str("APP_USER") or os.getenv("APP_USER", "").strip() or "admin"
            users = {user: pwd}
    return users


def _auth_enabled() -> bool:
    return bool(_get_app_users())


def _make_auth_token(username: str, password: str) -> str:
    """HMAC 令牌：username:digest，可存浏览器，不含明文密码。"""
    digest = hmac.new(
        password.encode(),
        f"{username}:".encode() + _AUTH_MSG,
        hashlib.sha256,
    ).hexdigest()
    return f"{username}:{digest}"


def _verify_auth_token(token: str, users: dict[str, str]) -> str | None:
    """校验令牌，成功返回用户名；兼容旧版仅 digest 的单用户令牌。"""
    if not token or not users:
        return None
    if ":" not in token:
        for user, pwd in users.items():
            legacy = hmac.new(pwd.encode(), _AUTH_MSG, hashlib.sha256).hexdigest()
            try:
                if hmac.compare_digest(token, legacy):
                    return user
            except Exception:
                continue
        return None
    username, digest = token.split(":", 1)
    pwd = users.get(username)
    if not pwd or not digest:
        return None
    expected = hmac.new(
        pwd.encode(),
        f"{username}:".encode() + _AUTH_MSG,
        hashlib.sha256,
    ).hexdigest()
    try:
        return username if hmac.compare_digest(digest, expected) else None
    except Exception:
        return None


def _clear_auth_query() -> None:
    if "auth" in st.query_params:
        del st.query_params["auth"]


def _storage_bridge(storage_key: str, query_key: str) -> None:
    """把 localStorage 里的值同步到 URL query，刷新后可恢复状态。"""
    components.html(
        f"""
        <script>
        (function() {{
          try {{
            const saved = localStorage.getItem({storage_key!r});
            const url = new URL(window.location);
            if (saved && !url.searchParams.has({query_key!r})) {{
              url.searchParams.set({query_key!r}, saved);
              window.location.replace(url);
            }}
          }} catch (e) {{}}
        }})();
        </script>
        """,
        height=0,
    )


def _persist_storage(storage_key: str, value: str) -> None:
    safe = value.replace("\\", "\\\\").replace("'", "\\'")
    components.html(
        f"<script>try{{localStorage.setItem('{storage_key}','{safe}');}}catch(e){{}}</script>",
        height=0,
    )


def _clear_storage(storage_key: str) -> None:
    components.html(
        f"<script>try{{localStorage.removeItem('{storage_key}');}}catch(e){{}}</script>",
        height=0,
    )


def _init_main_nav() -> None:
    _storage_bridge(_NAV_COOKIE, "tab")
    tab = st.query_params.get("tab", "")
    tab = _LEGACY_NAV.get(tab, tab)
    if tab in _MAIN_NAV:
        st.session_state["main_nav"] = tab
    elif "main_nav" not in st.session_state:
        st.session_state["main_nav"] = "overview"
    nav_key = st.session_state.get("main_nav", "overview")
    if nav_key in _MAIN_NAV:
        st.session_state["main_nav_radio"] = _MAIN_NAV[nav_key]


def _go_to_nav(key: str) -> None:
    if key not in _MAIN_NAV:
        return
    st.session_state["main_nav"] = key
    st.session_state["main_nav_radio"] = _MAIN_NAV[key]
    st.query_params["tab"] = key
    _persist_storage(_NAV_COOKIE, key)
    st.rerun()


def _render_app_header_compact() -> None:
    """顶栏：Kidde logo + pill 主导航；非制作页时右侧显示快捷「开始制作」。"""
    b64 = _icon_b64()
    img = (
        f"<img src='data:image/png;base64,{b64}' class='nav-dino-logo' alt='Kidde'/>"
        if b64 else "<span class='nav-dino-fallback'>🦕</span>"
    )
    nav = st.session_state.get("main_nav", "overview")
    labels = list(_MAIN_NAV.values())
    keys = list(_MAIN_NAV.keys())
    idx = keys.index(nav) if nav in keys else 0

    st.markdown('<div id="app-header-anchor"></div>', unsafe_allow_html=True)
    with st.container(border=True):
        c_brand, c_nav, c_act = st.columns([1.05, 1.65, 0.65], vertical_alignment="center")
        with c_brand:
            st.markdown(
                f"""
                <div class='app-topbar-brand'>
                  {img}
                  <div>
                    <div class='app-topbar-title'>VIPKID Dino</div>
                    <div class='app-topbar-sub'>线下绘本教学系统</div>
                  </div>
                </div>
                """,
                unsafe_allow_html=True,
            )
        with c_nav:
            st.markdown('<div class="main-nav-wrap">', unsafe_allow_html=True)
            selected = st.radio(
                "主导航",
                labels,
                index=idx,
                horizontal=True,
                label_visibility="collapsed",
                key="main_nav_radio",
            )
            st.markdown("</div>", unsafe_allow_html=True)
        with c_act:
            st.markdown('<div class="header-actions">', unsafe_allow_html=True)
            if _auth_enabled() and not st.session_state.get("_authed"):
                if st.button("登录", key="hdr_login", use_container_width=True):
                    _clear_auth_query()
                    st.session_state.pop("_authed", None)
                    st.session_state.pop("_authed_user", None)
                    _clear_storage(_AUTH_COOKIE)
                    st.rerun()
            elif nav != "work":
                if st.button("开始制作 →", key="hdr_start", type="primary", use_container_width=True):
                    _go_to_nav("work")
            st.markdown("</div>", unsafe_allow_html=True)

    selected_key = keys[labels.index(selected)]
    if selected_key != nav:
        st.session_state["main_nav"] = selected_key
        st.session_state["main_nav_radio"] = _MAIN_NAV[selected_key]
        st.query_params["tab"] = selected_key
        _persist_storage(_NAV_COOKIE, selected_key)
        st.rerun()


def _render_overview_section() -> None:
    """概览：Hero + 课程地图卡片 + 关键数字 + CTA。"""
    _render_hero()
    st.markdown(render_level_cards_html(), unsafe_allow_html=True)
    m1, m2, m3, m4 = st.columns(4)
    m1.metric("4 大交付物", _KIT_LABEL)
    m2.metric("必填输入", "书名 + Level")
    m3.metric("AI 抽取", "约 2–3 分钟")
    m4.metric("级别覆盖", "L0 – L6")
    with st.expander(
        "📊 详细梯度图（HTML 迷你地图 · 下载长图/PDF/Excel）",
        expanded=False,
    ):
        _render_curriculum_map()
    st.markdown("---")
    c1, c2 = st.columns([1, 3])
    with c1:
        if st.button("👉 开始制作", type="primary", use_container_width=True, key="cta_work"):
            _go_to_nav("work")
    with c2:
        st.caption("第一次使用？建议先看 **新手引导** 走查一遍流程。")
    with st.expander("版本信息", expanded=False):
        st.markdown(
            "**v3.3**：文本走 Claude · 生图走 gpt-image-2 · 中文 prompt（主体+行为+环境）。"
        )


def _render_background_section() -> None:
    """背景：为什么做、教学定位、全局底层逻辑。"""
    st.subheader("背景")
    st.markdown(
        "VIPKID Dino 线下绘本教学系统，面向 **0–13 岁** 分级英文故事课。"
        "把老师从「重复排版、拆词、出题、做课件」里解放出来，"
        "让你把精力放在 **故事审核与课堂设计** 上。"
    )
    st.markdown("##### 设计原则")
    b1, b2, b3 = st.columns(3)
    with b1:
        st.markdown("**📐 标准焊死**  \nCEFR / 词表格式 / 页数 / 留白等由系统强制套用，批量生产不跑偏。")
    with b2:
        st.markdown("**🎭 IP 一致**  \n官方角色库 + 参考图锁定，全本生图主角不跳帧。")
    with b3:
        st.markdown("**✏️ 人机协同**  \nAI 先抽、老师后改；每一步可审核、可回退。")
    st.divider()
    _render_global_standards_panel()


def _render_features_section() -> None:
    """功能：各模块能做什么（信息层，不含操作台）。"""
    st.subheader("功能")
    st.caption("以下为系统能力说明；实际操作请切到 **开始制作**。")
    feats = [
        ("🤖 AI 智能抽取", "从故事原文自动拆 7 页、抽词表/拼读/语法、出 Worksheet 与 RR 题。"),
        ("🖼️ 绘本", "IP 锁定 → 画风设定 → 分页编辑 → 单页生图 → 组装 PPT Reader。"),
        ("📝 练习", "逐题换型/改难度/AI 重出/配图，一键生成 Worksheet PPTX。"),
        ("📄 RR", "阅读表达题预览与单题修改，空白版/示例版可选。"),
        ("👩‍🏫 TG", "自动生成 Teacher's Guide DOCX。"),
        ("📚 批量生产", "一次跑多本大纲，适合系列化备课。"),
        ("📤 上传成品绘本", "已有 PDF/PPT/散图 → 只出 WS + RR + TG 教辅三件套。"),
    ]
    for title, desc in feats:
        with st.expander(title, expanded=False):
            st.markdown(desc)
    st.divider()
    st.markdown("##### 4 大交付物规格速查")
    cols = st.columns(4)
    for i, key in enumerate(["book", "worksheet", "rr", "tg"]):
        spec = DELIVERABLE_SPECS[key]
        with cols[i]:
            with st.popover(f"{spec['icon']} {spec['name']}"):
                st.markdown(render_deliverable_spec_md(key))


def _faq_items() -> list[tuple[str, str]]:
    return [
        (
            "刷新后又要重新登录？",
            "登录成功后会在浏览器 **localStorage** 保存 HMAC 校验令牌（**不含明文密码**）。"
            "若仍失效，请检查浏览器是否禁用本地存储。",
        ),
        (
            "账户怎么配置？",
            "管理员在 Secrets / 环境变量配置 **APP_USERS**（多用户）或 **APP_USER + APP_PASSWORD**（单用户）。"
            "登录页填写 **用户名 + 密码**。",
        ),
        (
            "AI 抽取大概要多久？",
            "通常 **2–3 分钟**。**抽取过程中请勿关闭或刷新页面**，否则会中断任务。",
        ),
        (
            "哪些是必填项？",
            "**Book Title** 与 **Level**；故事由 AI 根据书名自动生成，可在文本框微调。",
        ),
        (
            "故事怎么分页？",
            "AI 生成或粘贴的故事 → 自动均分到 7 页；若标注了 `Page 1`… 则按你的分页走。",
        ),
        (
            "4 件套分别是什么？",
            f"{_KIT_LABEL}；共享同一份 AI 抽取数据。",
        ),
        (
            "生图 API 不可用？",
            "无 Key 时走 **mock 占位图**，文本抽取仍可用。详见 **设置 → 服务连接**。",
        ),
    ]


def _render_faq_accordion(open_first: bool = True) -> None:
    """FAQ 手风琴（HTML details，样式对齐 mockup）。"""
    blocks = []
    for i, (q, a) in enumerate(_faq_items()):
        open_attr = " open" if open_first and i == 0 else ""
        a_html = a.replace("\n", "<br/>")
        blocks.append(
            f"<details class='faq-item'{open_attr}>"
            f"<summary><span class='faq-q'>{q}</span><span class='faq-chevron'>▼</span></summary>"
            f"<div class='faq-a'>{a_html}</div></details>"
        )
    st.markdown(
        f"""
        <section class='faq-section' id='faq'>
          <div class='faq-head'>
            <h2>常见问题 (FAQ)</h2>
            <p>了解更多关于 AI 抽取流程与平台的使用细节</p>
          </div>
          <div class='faq-list'>{"".join(blocks)}</div>
        </section>
        """,
        unsafe_allow_html=True,
    )


def _render_footer() -> None:
    """页脚链接区（mockup footer）。"""
    st.markdown(
        """
        <footer class='site-footer'>
          <div class='footer-brand'>
            <span class='footer-logo'>VIPKID Dino</span>
            <p>© 2024 VIPKID Dino. Empowering young learners through AI storytelling.</p>
          </div>
          <div class='footer-links'>
            <span>Privacy Policy</span>
            <span>Terms of Service</span>
            <span>Contact Us</span>
            <span>About</span>
          </div>
        </footer>
        """,
        unsafe_allow_html=True,
    )


def _render_onboarding_section() -> None:
    """新手引导：5 步走查 + FAQ。"""
    st.subheader("新手引导")
    steps = [
        ("填写书名与级别", "在 **开始制作** 填 **Book Title** + **Level**（必填）。"),
        ("生成故事草稿", "点 **AI 生成故事草稿**，或抽取时自动生成。"),
        ("人工微调", "在故事文本框里改句子、换人名或细节。"),
        ("AI 抽取", "约 2–3 分钟，期间 **勿刷新页面**。"),
        ("微调并下载", "切换 **绘本 / 练习 / RR / TG**，组装后下载 ZIP。"),
    ]
    for i, (title, desc) in enumerate(steps, 1):
        st.markdown(f"**第 {i} 步 · {title}**  \n{desc}")
    st.divider()
    if st.button("去开始制作 →", type="primary", key="onboard_go_work"):
        _go_to_nav("work")
    _render_faq_accordion(open_first=False)


def _render_metrics_section() -> None:
    """指标：与教研 Excel 同源的完整对标 + 绘本产出相关字段。"""
    st.subheader("指标")
    st.markdown(
        "对标 **VIPKID Dino 北美外教阅读能力达成** 与线下绘本 0–6 体系。"
        "下表与《课程对标总表_L0-L6.xlsx》同源；带 **ᴿ** 的维度为参考对标值，其余为官方权威值。"
    )
    st.markdown("##### 级别速览（选 Level 时对照）")
    st.dataframe(level_metrics_rows(), use_container_width=True, hide_index=True)

    st.markdown("##### 完整维度表（与 Excel 六段一致）")
    for section_name, section_data in curriculum_section_tables():
        with st.expander(section_name, expanded=(section_name.startswith("①"))):
            st.dataframe(section_to_rows(section_data), use_container_width=True, hide_index=True)

    st.markdown("##### 绘本产出质量指标（抽取后 / 设置 → 体检）")
    st.markdown(
        f"- **{_KIT_LABEL}**：共享同一份 AI 抽取数据  \n"
        "- **词汇**：lemma 小写、L0–2 双行词表 / L3–6 单行  \n"
        "- **绘本**：IP 年龄档、参考图、7 页分页  \n"
        "- **练习**：6 页固定结构  \n"
        "- **RR**：L0–2 共 4 题 / L3–6 共 5 题（星级分布）  \n"
        "- **TG**：8 模块 + 与 Worksheet 答案一致"
    )
    _ensure_curriculum_assets()
    if any(p.exists() for p in (_CMAP_PDF, _CMAP_XLSX, _CMAP_LONGIMG)):
        st.markdown("##### 下载完整对标资料")
        dl = st.columns(4)
        with dl[0]:
            if _CMAP_LONGIMG.exists():
                st.download_button("⬇️ 竖版长图", _CMAP_LONGIMG.read_bytes(),
                                   file_name=_CMAP_LONGIMG.name, mime="image/png",
                                   key="metrics_dl_long")
        with dl[1]:
            if _CMAP_PDF.exists():
                with open(_CMAP_PDF, "rb") as f:
                    st.download_button("⬇️ PDF 一页纸", f.read(), _CMAP_PDF.name, key="metrics_dl_pdf")
        with dl[2]:
            if _CMAP_XLSX.exists():
                with open(_CMAP_XLSX, "rb") as f:
                    st.download_button("⬇️ Excel 总表", f.read(), _CMAP_XLSX.name, key="metrics_dl_xlsx")
        with dl[3]:
            st.download_button("⬇️ 迷你图 PNG", mini_map_png_bytes(),
                               file_name="VIPKID_Dino_课程地图迷你图_L0-L6.png",
                               mime="image/png", key="metrics_dl_mini")


def _render_settings_section() -> None:
    st.subheader("设置")
    st.markdown("##### 🔑 服务连接")
    _key_status_banner()
    st.caption("API Key 通过环境变量或 Streamlit Secrets 配置，页面不会显示或缓存 Key 本身。")

    st.divider()
    st.markdown("##### 🎨 默认画风（新绘本自动套用，工作台内仍可再改）")
    _render_style_panel_step(3, embed=True)

    st.divider()
    st.markdown("##### 🔬 质量体检")
    _render_evals_panel()

    st.divider()
    if _auth_enabled():
        user = st.session_state.get("_authed_user", "")
        if user:
            st.caption(f"已登录账户：**{user}**")
        if st.button("退出登录", type="secondary", key="logout_btn"):
            st.session_state.pop("_authed", None)
            st.session_state.pop("_authed_user", None)
            _clear_auth_query()
            _clear_storage(_AUTH_COOKIE)
            st.rerun()


def _validate_required_inputs(title: str, raw_story: str = "") -> list[str]:
    errors: list[str] = []
    if not (title or "").strip():
        errors.append("Book Title 不能为空")
    story = (raw_story or "").strip()
    if story and len(story.split()) < 15:
        errors.append("故事过短（建议至少 15 个英文词），请补充或重新生成")
    return errors


def _render_login_gate() -> None:
    """未登录时展示的登录页（含品牌区，不含主导航）。"""
    b64 = _icon_b64()
    img = (
        f"<img src='data:image/png;base64,{b64}' class='nav-dino-logo' alt='Kidde'/>"
        if b64 else "<span class='nav-dino-fallback'>🦕</span>"
    )
    st.markdown(
        f"""
        <div class='app-topbar-brand' style='margin-bottom:1rem;'>
          {img}
          <div>
            <div class='app-topbar-title'>VIPKID Dino</div>
            <div class='app-topbar-sub'>线下绘本教学系统</div>
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.markdown("## 🔒 登录")
    st.caption("请输入用户名与密码（向管理员获取）。登录状态会记住本机，**不会保存明文密码**。")
    with st.form("login_gate", clear_on_submit=False):
        entered_user = st.text_input("用户名", autocomplete="username")
        entered_pwd = st.text_input("密码", type="password", autocomplete="current-password")
        ok = st.form_submit_button("登录", type="primary")
    if ok:
        users = _get_app_users()
        user = entered_user.strip()
        pwd = users.get(user, "")
        if pwd and secrets.compare_digest(entered_pwd, pwd):
            auth_token = _make_auth_token(user, pwd)
            st.session_state["_authed"] = True
            st.session_state["_authed_user"] = user
            st.query_params["auth"] = auth_token
            _persist_storage(_AUTH_COOKIE, auth_token)
            st.rerun()
        else:
            st.error("用户名或密码错误，请重试。")


def _restore_auth_session() -> bool:
    """尝试恢复登录态。未配置账户时视为已通过。"""
    users = _get_app_users()
    if not users:
        return True

    _storage_bridge(_AUTH_COOKIE, "auth")

    if st.session_state.get("_authed"):
        return True

    token = st.query_params.get("auth", "")
    authed_user = _verify_auth_token(token, users)
    if authed_user:
        st.session_state["_authed"] = True
        st.session_state["_authed_user"] = authed_user
        return True
    if token:
        _clear_auth_query()
        _clear_storage(_AUTH_COOKIE)
    return False


def main() -> None:
    _icon_arg = str(_DINO_ICON) if _DINO_ICON.exists() else "📘"
    st.set_page_config(
        page_title="VIPKID Dino 线下绘本教学",
        page_icon=_icon_arg,
        layout="wide",
        initial_sidebar_state="collapsed",
    )
    _inject_css()
    _init_main_nav()

    auth_ok = _restore_auth_session()
    _render_app_header_compact()
    if not auth_ok:
        _render_login_gate()
        _render_footer()
        return

    # Session 状态
    if "extracted" not in st.session_state:
        st.session_state.extracted = None
    if "outline" not in st.session_state:
        st.session_state.outline = None

    nav = st.session_state.get("main_nav", "overview")

    if nav == "overview":
        _render_overview_section()
        _render_faq_accordion()
        _render_footer()
        return
    if nav == "background":
        _render_background_section()
        _render_footer()
        return
    if nav == "features":
        _render_features_section()
        _render_footer()
        return
    if nav == "onboarding":
        _render_onboarding_section()
        _render_footer()
        return
    if nav == "metrics":
        _render_metrics_section()
        _render_footer()
        return
    if nav == "settings":
        _render_settings_section()
        _render_footer()
        return

    # ---------- nav == work：制作工作台 ----------
    st.markdown(
        """
        <section class='create-section' id='creation'>
          <div class='create-head'>
            <span class='create-icon'>✨</span>
            <div>
              <h2>生成绘本</h2>
            </div>
          </div>
        </section>
        """,
        unsafe_allow_html=True,
    )
    st.caption("在此填写书名、级别、故事 → AI 抽取 → 出图 → 下载 4 件套")
    _key_status_banner()

    nav_sel = _render_deliverable_nav()

    # ---------- 模式：单本 / 批量 ----------
    mode = st.radio(
        "制作模式",
        ["📖 单本制作", "📚 批量生产", "📤 上传成品绘本 → 教辅"],
        horizontal=True,
        key="produce_mode",
        help="单本：逐本精调出图 + 4 件套。批量：一次跑多本大纲。上传成品：已有绘本(PDF/PPT/散图)→只出 Worksheet+RR+TG。",
    )
    if mode == "📚 批量生产":
        _render_batch_mode()
        return
    if mode == "📤 上传成品绘本 → 教辅":
        _render_upload_mode()
        return

    # ---------- Section A：输入表单 ----------
    st.success(
        f"🎯 **必填 2 项**：① Book Title  ② Level  \n"
        "   **故事由 AI 根据书名自动生成**（约 30 秒），可在下方文本框人工微调后再抽取  \n"
        f"   其余（CEFR / 词表 / 语法 / 题目等）在 **AI 抽取** 时自动推断 → 产出 {_KIT_LABEL}"
    )

    if "raw_story_input" not in st.session_state:
        st.session_state.raw_story_input = ""

    with st.form("input_form"):
        # === 必填区 ===
        _sec_head("1", "必填基础信息")
        _chips(["① Book Title", "② Level", "③ 故事（AI 生成·可改）"], kind="ok")
        col1, col2 = st.columns([3, 1])
        with col1:
            title = st.text_input(
                "📕 Book Title *",
                value="",
                placeholder="例如：What Makes a Good Friend?",
                help="必填。系统会用它做文件名 / 封面 / 各文档大标题，并生成故事内容。",
            )
        with col2:
            level = st.selectbox(
                "🎚️ Level *", LEVEL_OPTIONS, index=0,
                help="必填。Smart / 0-2 = 双行词表，3-6 = 单行词表。决定 CEFR / Reader Type / 题数。",
            )

        st.markdown(
            "**📝 故事正文** — AI 根据书名自动生成；你可在此 **微调** 后再抽取"
            "（系统均分 7 页 · 封面另算共 8 页）"
        )
        st.caption(
            "流程：先填书名 + Level → 点「AI 生成故事草稿」→ 改满意了 → 点「AI 抽取」。"
            "若故事为空直接点抽取，也会自动生成。"
        )
        raw_story = st.text_area(
            "Raw story",
            label_visibility="collapsed",
            height=220,
            placeholder="填写书名后点下方「AI 生成故事草稿」，或留空在抽取时自动生成…",
            key="raw_story_input",
        )

        # === 选填区（折叠）===
        st.subheader("2️⃣  选填（不填 AI 全自动）")
        with st.expander("⚙️  选填字段（让交付物更精准）", expanded=False):
            col1, col2, col3 = st.columns(3)
            with col1:
                book_number = st.text_input(
                    "Book #", value="01",
                    help="书号，用于文件名 `Level X_BookXX_品类_标题.后缀`",
                )
                theme = st.text_input(
                    "Theme",
                    value="",
                    placeholder="留空 = AI 按书名推断",
                    help="主题，会用在 Writing 页 \"Write about ...\" 和 TG 的目标里",
                )
            with col2:
                cefr = st.text_input(
                    "CEFR (留空 = 按 Level 自动)",
                    value="",
                    help="Smart/L0/L1=Pre-A1, L2=A1, L3=A1+, L4=A2, L5=B1, L6=B1+",
                )
                # 年龄焊死：由 Level 自动决定（Smart/L0-2=8、L3-4=10、L5-6=12），
                # 默认锁定，只有勾选「高级手改」才允许覆盖，避免人物年龄跑偏。
                ip_age_default = resolve_ip_age(level)
                st.markdown(
                    f"**IP 年龄：{ip_age_default} 岁** 🔒（按 Level 自动焊死）"
                )
                _age_override = st.checkbox(
                    "高级：手动覆盖年龄", value=False,
                    help="一般不需要。Smart/L0-2=8 岁；L3-4=10 岁；L5-6=12 岁，由 Level 决定",
                )
                if _age_override:
                    ip_age = st.number_input(
                        "手动 IP 年龄", min_value=6, max_value=14,
                        value=ip_age_default,
                    )
                else:
                    ip_age = ip_age_default
            with col3:
                # L3-6 才让选 Fiction / Non-Fiction
                level_digits = "".join(ch for ch in level if ch.isdigit())
                try:
                    lvl_n = int(level_digits) if level_digits else 0
                except ValueError:
                    lvl_n = 0
                if lvl_n >= 3:
                    fiction_type = st.selectbox(
                        "Reader Type (L3-6)",
                        ["（AI 自动判断）", "fiction", "non-fiction"],
                        index=0,
                        help="L3-6 在 Reading Report 第一行显示为 'Fiction' 或 'Non-Fiction'。"
                             "默认让 AI 按故事内容判断；如需强制可手选。",
                    )
                else:
                    fiction_type = ""
                    rt_map = {
                        "smart": "Concept & Knowledge - Building Readers",
                        "0": "Concept & Knowledge - Building Readers",
                        "1": "Patterned Narrative & Informational Readers",
                        "2": "Early Independent Genre-Exposure Readers",
                    }
                    rt = rt_map.get(level.lower(), rt_map.get(level_digits, ""))
                    st.info(f"📖 Reader Type：**{rt}**（按 Level 固定，无需选）")

            st.divider()
            st.markdown("**👥 角色设定**（已注册角色系统自动识别，新人物才需在下方手填）")
            try:
                from character_registry import list_available
                chars = list_available()
                cols = st.columns(3)
                for i, ch in enumerate(chars):
                    with cols[i % 3]:
                        ages_str = "/".join(str(a) for a in ch["age_options"])
                        emoji = {"protagonist": "⭐", "supporting": "👥", "adult": "👩‍🏫",
                                 "pet": "🐱", "brand": "🦖", "family": "👨‍👩‍👧"}.get(ch["kind"], "•")
                        check = "✅" if ch["reference_exists"] else "⚠️"
                        disp_name = ch["key"].replace("_", " ").title()
                        st.markdown(f"{emoji} **{disp_name}** ({ages_str}) {check}")
                st.caption("✅ = 有官方参考图  ⚠️ = 缺参考图。出现这些名字时系统会自动加载形象。")
            except Exception as e:
                st.warning(f"角色注册表加载失败：{e}")
            custom_chars_text = st.text_area(
                "🆕 新人物注册（只在故事出现 \"全新\" 人物时填，每行：name | description）",
                value="",
                height=60,
                help="如：lucy | 8y GIRL, twin braids, red sweater, freckles —— 已在角色库的角色（mia/tommy/anna/teacher_kim/winnie 等）不用填",
            )

        btn_gen, btn_ext = st.columns(2)
        with btn_gen:
            gen_story_btn = st.form_submit_button(
                "✨ AI 生成故事草稿", width="stretch",
            )
        with btn_ext:
            submitted = st.form_submit_button(
                "🤖 AI 抽取 / 重新抽取", type="primary", width="stretch",
            )

    if gen_story_btn:
        if not (title or "").strip():
            st.error("❌ 请先填写 Book Title")
        else:
            draft = run_with_live_timer(
                "AI 生成故事草稿", generate_story_draft,
                title.strip(), level, theme.strip(),
                done_note="可在下方微调",
            )
            st.session_state.raw_story_input = draft
            st.success("✅ 故事草稿已生成，请在下方微调后点「AI 抽取」。")
            st.rerun()

    if submitted:
        errors = _validate_required_inputs(title, raw_story)
        if errors:
            for err in errors:
                st.error(f"❌ {err}")
        else:
            story_text = (raw_story or "").strip()
            if not story_text:
                story_text = run_with_live_timer(
                    "AI 生成故事草稿", generate_story_draft,
                    title.strip(), level, theme.strip(),
                )
                st.session_state.raw_story_input = story_text
            raw_story = story_text
            st.warning(
                "⏳ **本次 AI 抽取预计需要 2–3 分钟**，请稍候。"
                "抽取过程中 **请勿关闭或刷新** 当前页面，否则会中断任务。"
            )
            # 先做 AI 自动推断（角色识别 + CEFR/Lexile/Theme/Fiction-NF）—— 实时计时
            auto = run_with_live_timer(
                "AI 推断（角色 / CEFR / 主题）", auto_summary, level, raw_story, title,
            )
            # 用户没填 cefr / theme / fiction_type 时，用自动值兜底
            cefr_final = cefr.strip() or auto["cefr"]
            theme_final = theme.strip() or auto["theme"]
            # Reader Type：默认「（AI 自动判断）」→ 用 AI 判定值；否则用老师手选值
            _ft = (fiction_type or "").strip()
            fiction_final = auto["fiction_type"] if _ft.startswith("（AI") else (_ft or auto["fiction_type"])

            ec = run_with_live_timer(
                "AI 抽取（抽词 · 拆段 · 出题）", extract_all,
                raw_story=raw_story, title=title, level=level,
                cefr=cefr_final, theme=theme_final,
            )
            outline = _build_outline(
                title=title, level=level, book_number=book_number,
                cefr=cefr_final, theme=theme_final, ip_age=int(ip_age),
                raw_story=raw_story, custom_chars_text=custom_chars_text,
                fiction_type=fiction_final,
            )
            apply_extracted_to_outline(outline, ec)
            if enrich_from_syllabus(outline):
                st.session_state["_syllabus_hit"] = True
                _sync_ec_from_syllabus(ec, outline)   # 命中大纲→把官方词表/拼读逐字同步进工作台
            st.session_state.extracted = ec
            st.session_state.outline = outline
            st.session_state.auto = auto
            st.session_state.auto["_lexile"] = auto["lexile"]  # 留作生成阶段元信息

            # v1.9：预生成每页中文 prompt（按 Seedream 4.5 官方指南），存 session 供用户编辑
            ip_age_val = int(ip_age)
            cast_pool = st.session_state.get("story_cast_pool") or None
            generic_overrides = st.session_state.get("generic_overrides") or None

            def _gen_page_prompts() -> dict[int, dict]:
                pp: dict[int, dict] = {}
                for page in outline.pages:
                    built = build_cn_page_prompt(
                        page, outline, ip_age_val,
                        cast_pool=cast_pool, generic_overrides=generic_overrides,
                    )
                    pp[page.index] = {
                        "positive": built.positive,
                        "negative": built.negative,
                        "prompt": built.prompt,
                        "references": [str(r) for r in built.references],
                        "must_include": "",
                        "label": page.label,
                        "display_name": page_display_name(page.index),
                    }
                return pp

            page_prompts = run_with_live_timer(
                "生成每页画面提示词", _gen_page_prompts,
                done_note=f"{len(outline.pages)} 页",
            )
            st.session_state.page_prompts = page_prompts

            st.success(
                "✅ 抽取完成。请审核下方「AI 推断 + 主角识别 + 每页卡片（文本+场景+prompt+必须出现）」，"
                "再点最底部 Generate All。"
            )

    # ---------- v2.0：7 步严格解锁绘本组装工作流 ----------
    if st.session_state.extracted is not None:
        # 抽取完成 → 底层逻辑已前置为全局只读面板，绘本轨直接解锁到 Step 3（IP+画风）
        if st.session_state.get("book_unlocked_step", 1) < 3:
            st.session_state["book_unlocked_step"] = 3

        st.divider()
        st.subheader("🛠️ 4 交付物工作台（4 个角标 · 共享同一份 AI 抽取数据 · 一起产出）")

        # 共享抽取数据（词表/语法/拼读/读者类型）— 4 件套都用，放在角标上方
        _render_shared_extract_panel()

        # 交付物横向导航选中哪个 → 主区只渲染那一个
        _sel_label = dict(_DELIVERABLES).get(nav_sel, "")
        st.markdown(f"#### {_sel_label}")

        if nav_sel == "book":
            with st.expander("📋 绘本规格（点开查看）", expanded=False):
                st.markdown(render_deliverable_spec_md("book", st.session_state.outline.level))
            # 向导式：顶部步骤条 + 一次只显示一个阶段（最易上手）
            _render_book_wizard()
        elif nav_sel == "ws":
            st.caption("逐题打磨：①换题型/难度 ②🤖 AI 重出 ③手动改 ④配图来源。改完点「👀 生成初稿」预览。")
            _render_worksheet_editor()
        elif nav_sel == "rr":
            st.caption("阅读报告字段与阅读表达题预览/单题改；空白版/示例版双模式在组装时选。")
            _render_rr_editor()
        elif nav_sel == "tg":
            _render_tg_panel()


# =============================== 编辑器 ===============================
def _render_auto_summary_panel() -> None:
    """v1.8.2：AI 自动推断 + 主角识别透明化面板。

    展示两块内容，让老师在生图前能审核：
      1. AI 推断卡片 — CEFR / 蓝思 / 字数 / 故事类型 / 主题
      2. 主角识别面板 — 故事里识别到的官方 IP（含参考图 ✓/⚠️）+ 未命名 girl/boy 的默认形象建议
    """
    auto = st.session_state.auto

    st.divider()
    st.subheader("🤖 AI 自动推断（你不用填，系统已根据故事 + 级别算好）")

    outline = st.session_state.get("outline")
    _lex_disp, _lex_src = _lexile_display(outline)
    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("📏 CEFR 等级", auto["cefr"])
    c2.metric("📊 蓝思 Lexile", _lex_disp, help=f"来源：{_lex_src}")
    c3.metric("🔢 故事字数", auto["word_count"])
    c4.metric("📖 故事类型", "Fiction" if auto["fiction_type"] == "fiction" else "Non-Fiction")
    c5.metric("🏷️ 主题", auto["theme"] or "(空)")

    st.caption(
        "💡 以上字段 AI 自动生成，需要覆盖请去顶部「⚙️ 选填字段」展开后填入对应项 → 会覆盖自动值。"
    )

    # 块11：蓝思值半自动取值（禁止编造）—— 大纲没有官方 Lexile 时，用官方分析器取真实值后回填
    _render_lexile_panel(outline)

    # === 主角识别面板 ===
    st.subheader("🎭 主角识别（让你确认 AI 理解的故事人物是谁）")
    chars = auto.get("characters", [])
    generic = auto.get("generic_roles", [])

    if not chars and not generic:
        st.info("ℹ️ 未在故事里识别到已注册的 IP 角色。系统将按通用 girl/boy 生成（默认套 Mia/Tommy 形象）。")
        return

    if chars:
        st.markdown(f"**✅ 已匹配到 {len(chars)} 个官方 IP**（生图时会自动用对应参考图保证形象一致）：")
        ncols = min(len(chars), 4)
        cols = st.columns(ncols)
        for i, ch in enumerate(chars):
            with cols[i % ncols]:
                ref_badge = "✅ 有官方参考图" if ch["reference_exists"] else "⚠️ 缺参考图"
                kind_emoji = {
                    "protagonist": "⭐", "supporting": "👥", "adult": "👩‍🏫",
                    "pet": "🐱", "brand": "🦖", "family": "👨‍👩‍👧",
                }.get(ch["kind"], "•")
                st.markdown(
                    f"### {kind_emoji} {ch['name_in_story']}\n"
                    f"**Registry key**: `{ch['matched_key']}`  \n"
                    f"**类别**: {ch['kind']}  |  **性别**: {ch['gender']}  |  **年龄**: {ch['age']}y  \n"
                    f"**参考图**: {ref_badge}"
                )
                if ch["reference_exists"] and ch["reference_path"]:
                    try:
                        _zoom_image(ch["reference_path"],
                                    key=f"match{ch['matched_key']}",
                                    caption=ch["name_in_story"])
                    except Exception:
                        pass
                with st.expander(f"形象描述（生图 prompt 用）"):
                    st.caption(ch["description"])

    # === v2.1 新增：故事人物库 IP 选择器 + 无名角色映射 ===
    _render_ip_cast_panel(chars, generic)

    st.warning(
        "🔍 **请审核以上识别结果。** 如果识别错了（例如把 Mary 误识别为 Mia），"
        "请去顶部「⚙️ 选填字段」→「🆕 新人物注册」加一行 `mary | 12y GIRL, ...` 把新人物注册进来，"
        "再点 AI 重新抽取。"
    )


def _render_ip_cast_panel(detected_chars: list[dict], detected_generic: list[dict]) -> None:
    """v2.1: 故事人物库可视化选择器 + 无名角色映射。

    UI 结构：
      1. 「📚 这本绘本会出现的全部 IP」— 按 kind 分组，每个 IP 显示缩略图，
         可勾选。默认勾选已检测到的 IP，老师可加勾其他 IP。
      2. 「🎭 无名角色映射」— 故事文本里出现的 a girl/a boy 等，老师选用谁的形象。
      3. 结果写入 st.session_state['story_cast_pool'] (list of IP key) 和
         st.session_state['generic_overrides'] (dict role→key)，
         供 cn_prompt_builder 使用。
    """
    from ip_library import list_by_kind, get_ip
    from config import resolve_ip_age

    st.markdown("---")
    st.subheader("📚 这本绘本会出现的全部 IP（勾选 = 生图时会作为参考形象）")
    st.caption(
        "✅ 已自动勾选「故事里识别到的角色」。"
        "如果某页需要其他 IP 出场（如班里其他同学、家人探班），可加勾。"
        "**勾选项越多，生图时可挑选的参考图越多**。"
    )

    # 当前 outline 的 IP age（用于挑年龄档）
    outline = st.session_state.get("outline")
    level = outline.level if outline else "5"
    target_age = resolve_ip_age(level)

    # 自动默认勾选：识别到的 IP（按 age 档匹配） + 无名角色默认套的 IP
    auto_keys: set[str] = set()
    for ch in detected_chars:
        # ch["matched_key"] 可能是 base（"anna"），我们要带 age 后缀的（"anna_12"）
        from ip_library import resolve_name_to_ip
        ip = resolve_name_to_ip(ch["matched_key"], target_age)
        # 用户拍板（2026-06-04 / 06-06）：不自动拉「只有 8/10 岁档的支持角色」（在 12 岁书里会偏小）；
        # 但故事里【点名出现】且【正好有本级年龄档】的角色（如 Anna 12y）必须自动带上——
        # 它就是这本书的真实角色。其余无名 girl/boy 仍走 girl/boy→Mia/Tommy 映射，不新创造别人。
        if ip and (ip.kind != "supporting" or ip.age == target_age):
            auto_keys.add(ip.key)
    for g in detected_generic:
        from ip_library import resolve_generic_role
        ip = resolve_generic_role(g["role"], target_age)
        if ip:
            auto_keys.add(ip.key)

    # 上次手动选过的话，沿用；否则用 auto 默认
    if "story_cast_pool" not in st.session_state or not st.session_state["story_cast_pool"]:
        st.session_state["story_cast_pool"] = sorted(auto_keys)

    selected: set[str] = set(st.session_state.get("story_cast_pool", []))

    # 分组渲染（用户拍板 2026-06-06：精简为 4 段，固定顺序）：
    #   ① 主角（只 Mia/Tommy 三档）② 家人（爸妈·爷奶）
    #   ③ 同学·朋友·老师（含 Anna，supporting+adult+pet 合并）④ Dino
    sections = [
        ("⭐ 主角（Mia · Tommy）",        ["protagonist"],          True),
        ("👨‍👩‍👧 家人（爸妈 · 爷奶 · 宠物）", ["family", "pet"],        False),
        ("👥 同学 · 朋友 · 老师",          ["supporting", "adult"],  True),
        ("🦖 Dino",                        ["brand"],                False),
    ]
    grouped = list_by_kind()

    new_selected: set[str] = set()
    for label, kinds, expanded in sections:
        entries = [e for k in kinds for e in grouped.get(k, [])]
        if not entries:
            continue
        with st.expander(f"{label} ({len(entries)})", expanded=expanded):
            ncols = 4
            for row_start in range(0, len(entries), ncols):
                cols = st.columns(ncols)
                for j, e in enumerate(entries[row_start:row_start + ncols]):
                    with cols[j]:
                        try:
                            _zoom_image(e.image_path, key=f"ip{e.key}", caption=e.name)
                        except Exception:
                            pass
                        is_checked = st.checkbox(
                            e.name,
                            value=(e.key in selected),
                            key=f"ip_pick_{e.key}",
                            help=e.desc[:120],
                        )
                        if is_checked:
                            new_selected.add(e.key)

    st.session_state["story_cast_pool"] = sorted(new_selected)
    st.success(f"✅ 已选 **{len(new_selected)}** 个 IP 进入这本绘本的人物池")

    # === 无名角色映射器 ===
    if detected_generic:
        st.markdown("---")
        st.markdown("##### 🎭 无名角色映射（故事里 a girl / a boy 等 → 用哪个 IP 形象）")
        st.caption("默认按规则套（girl→Mia / boy→Tommy / woman→Mom 等）。可改成你想要的 IP。")

        overrides = st.session_state.get("generic_overrides", {}) or {}
        new_overrides: dict[str, str] = {}

        # 候选 IP = 当前选中的人物池（按 gender 过滤）
        pool_entries = [get_ip(k) for k in new_selected if get_ip(k)]
        # 性别过滤候选
        def _pick_options(role_gender: str) -> list[tuple[str, str]]:
            """返回 [(key, label), ...] 候选。"""
            cands = pool_entries
            if role_gender in ("girl", "woman"):
                cands = [e for e in pool_entries if e.gender in ("girl", "woman")]
            elif role_gender in ("boy", "man"):
                cands = [e for e in pool_entries if e.gender in ("boy", "man")]
            elif role_gender in ("cat",):
                cands = [e for e in pool_entries if e.kind == "pet"]
            opts = [("__skip__", "（跳过 / 不强制 IP）")]
            opts += [(e.key, e.name) for e in cands]
            return opts

        role_gender_map = {
            "girl": "girl", "boy": "boy",
            "woman": "woman", "man": "man",
            "old woman": "woman", "old man": "man",
            "cat": "cat", "kitty": "cat",
        }

        for g in detected_generic:
            role = g["role"]
            role_g = role_gender_map.get(role.lower(), "")
            opts = _pick_options(role_g)

            # 当前选择：先看 override，否则用 detected 的 default_key + age
            current = overrides.get(role, "")
            if not current:
                # 用 auto 推断
                from ip_library import resolve_generic_role
                ip = resolve_generic_role(role, target_age)
                current = ip.key if ip else "__skip__"

            keys_list = [o[0] for o in opts]
            labels_list = [o[1] for o in opts]
            default_idx = keys_list.index(current) if current in keys_list else 0

            sel = st.selectbox(
                f"故事里的 **{role}** → 套用谁？",
                options=range(len(opts)),
                index=default_idx,
                format_func=lambda i: labels_list[i],
                key=f"generic_pick_{role}",
                help=g.get("note", ""),
            )
            chosen_key = keys_list[sel]
            if chosen_key != "__skip__":
                new_overrides[role] = chosen_key

        st.session_state["generic_overrides"] = new_overrides


def _generate_worksheet_preview() -> None:
    """v1.9：基于当前编辑过的题目，立即生成一份 worksheet.pptx 初稿供老师下载预览。

    复用已生图的 page_xx.png 作为 Sentence 页/Match 页的图片来源；如果还没生图，用占位图。
    """
    ec = st.session_state.extracted
    outline: BookOutline = st.session_state.outline
    if not ec or not outline:
        st.error("请先点 AI 抽取按钮。")
        return
    apply_extracted_to_outline(outline, ec)
    rqc = int(st.session_state.get("ws_reading_q_count", 4))
    attach_worksheet_questions(outline, ec.worksheet_questions, reading_q_count=rqc)

    run_dir, img_dir, name_prefix = _ensure_run_dir()
    # 已生图就按 index 排序复用；没生图就传空列表（builder 用占位）
    image_results = st.session_state.get("image_results") or {}
    image_paths: list[Path] = []
    for idx in sorted(image_results.keys()):
        p = image_results[idx].get("path")
        if p and Path(p).exists():
            image_paths.append(Path(p))

    draft_path = run_dir / f"{name_prefix}_Worksheet_DRAFT.pptx"
    try:
        with st.spinner("正在生成 Worksheet 初稿..."):
            build_worksheet(
                outline,
                draft_path,
                image_paths=image_paths or None,
                sentence_image_mode=_ws_sentence_mode(),
                second_reading_mode=_ws_second_reading_mode(),
            )
    except Exception as e:
        st.error(f"Worksheet 初稿生成失败：{e}")
        return
    st.session_state.ws_draft_path = str(draft_path)
    st.toast("✅ Worksheet 初稿已生成，下方可下载", icon="✅")


def _generate_rr_preview() -> None:
    """基于当前编辑过的阅读表达题，立即生成一份 Reading Report 初稿供老师可视化预览。

    严格 1 页 A4（builder 内部实测页数 + 降档重排硬保证）。版本（空白/示例）跟随
    rr_answer_mode 选择。
    """
    ec = st.session_state.extracted
    outline: BookOutline = st.session_state.outline
    if not ec or not outline:
        st.error("请先点 AI 抽取按钮。")
        return
    apply_extracted_to_outline(outline, ec)
    attach_rr_questions(outline, ec.rr_questions)

    with_answers = str(st.session_state.get("rr_answer_mode", "")).startswith("示例")
    run_dir, _img_dir, name_prefix = _ensure_run_dir()
    draft_path = run_dir / f"{name_prefix}_Reading_Report_DRAFT.docx"
    try:
        run_with_live_timer(
            "生成 Reading Report 初稿（严格 1 页）",
            build_reading_report, outline, draft_path,
            done_note="含实测页数校验，确保单页",
            with_answers=with_answers,
        )
    except Exception as e:
        st.error(f"Reading Report 初稿生成失败：{e}")
        return
    st.session_state.rr_draft_path = str(draft_path)
    st.toast("✅ Reading Report 初稿已生成（严格 1 页），下方可预览/下载", icon="✅")


def _rebuild_all_cn_prompts() -> None:
    """v1.9：用最新 outline + ec.pages 重生所有页的中文 prompt（保留 must_include）。"""
    ec = st.session_state.extracted
    outline: BookOutline = st.session_state.outline
    if not ec or not outline:
        return
    # 把 UI 上改过的 text / scene 应用到 outline
    apply_extracted_to_outline(outline, ec)
    ip_age = outline.ip_age or resolve_ip_age(outline.level)
    page_prompts = st.session_state.get("page_prompts") or {}
    cast_pool = st.session_state.get("story_cast_pool") or None
    generic_overrides = st.session_state.get("generic_overrides") or None
    for page in outline.pages:
        built = build_cn_page_prompt(
            page, outline, ip_age,
            cast_pool=cast_pool, generic_overrides=generic_overrides,
        )
        prev_must = (page_prompts.get(page.index) or {}).get("must_include", "")
        page_prompts[page.index] = {
            "positive": built.positive,     # v3: 正向
            "negative": built.negative,     # v3: 反向
            "prompt": built.prompt,         # v3: 拼接后兜底（向后兼容）
            "references": [str(r) for r in built.references],
            "must_include": prev_must,
            "label": page.label,
            "display_name": page_display_name(page.index),
        }
    st.session_state.page_prompts = page_prompts


def _render_unified_page_panel() -> None:
    """v1.9：合并面板 — 每页一张卡片，包含：
      ① 故事文本（可改）
      ② AI 场景描述（可改）
      ③ 配图 prompt（中文，可改）
      ④ 必须出现（可加）
      ⑤ 参考图 + Shot + Expression

    页编号约定：Cover / Page 2 / Page 3 ... / Page 8（无 Page 1，Cover 即第 1 页印刷面）。
    """
    ec = st.session_state.extracted
    page_prompts = st.session_state.page_prompts

    with st.expander(
        f"🖼️ 每页详情（共 {len(page_prompts)} 页 — 强烈建议先审一遍再生图）",
        expanded=True,
    ):
        st.caption(
            "💡 一张卡片 = 一页绘本。改动顺序：先看 ① 故事原文 → 改 ② 场景描述 → 看 ③ 中文 prompt 是否到位 → "
            "在 ④ 必须出现 里补关键角色/道具（会追加到 prompt 末尾）。  \n"
            "📌 **页编号约定**：**Cover** 就是封面（印刷第 1 页）；**Page 2** 是故事第 1 句对应的图（印刷第 2 页），以此类推到 Page 8。**没有 Page 1**。"
        )

        ec_pages_by_idx = {p.get("index"): p for p in (ec.pages or [])}

        for idx in sorted(page_prompts.keys()):
            entry = page_prompts[idx]
            display_name = entry.get("display_name") or page_display_name(idx)
            refs = entry.get("references", [])
            ref_count = len(refs)
            ref_badge = f"🖼️ {ref_count} 张参考图" if ref_count else "⚠️ 无参考图"

            ec_page = ec_pages_by_idx.get(idx) or {}
            label_kind = "封面" if idx == 0 else "故事页"

            st.markdown(f"#### {display_name} · {label_kind} · {ref_badge}")

            # 第 1 行：故事文本 + 中文画面描述（scene_cn 是 prompt 的主体）
            c1, c2 = st.columns([1, 1])
            with c1:
                st.caption("① 故事原文（英文）")
                new_text = st.text_area(
                    f"text_{idx}",
                    value=ec_page.get("text", ""),
                    height=140,
                    key=f"uni_text_{idx}",
                    label_visibility="collapsed",
                )
                if new_text != ec_page.get("text", ""):
                    ec_page["text"] = new_text
            with c2:
                st.caption("② 中文画面描述 scene_cn（AI 生成 → 喂给画图模型，最关键）")
                new_scene_cn = st.text_area(
                    f"scene_cn_{idx}",
                    value=ec_page.get("scene_cn", ""),
                    height=140,
                    key=f"uni_scene_cn_{idx}",
                    label_visibility="collapsed",
                    placeholder=(
                        "AI 应该生成 120-220 字连贯描述：\n"
                        "主体（人物+具体外观）+ 动作（具体动词姿势）+ 环境（可见物品）+ 氛围（光照）\n"
                        "例：教室一角，Anna 坐在木课桌后，双手紧握放在桌面微微颤抖，眼神紧张（只写动作表情，不写外观）..."
                    ),
                )
                if new_scene_cn != ec_page.get("scene_cn", ""):
                    ec_page["scene_cn"] = new_scene_cn

            # 第 2 行：v3 正向 prompt + 反向 prompt + 必须出现
            c3a, c3b, c4 = st.columns([3, 2, 2])
            with c3a:
                st.caption("③✅ 正向 Prompt（火山风单段流畅自然语言，画什么）")
                # 兼容老 session：如果 entry 只有 "prompt" 没有 "positive"，把它当 positive
                cur_pos = entry.get("positive", entry.get("prompt", ""))
                new_pos = st.text_area(
                    f"positive_{idx}",
                    value=cur_pos,
                    height=240,
                    key=f"uni_pos_{idx}",
                    label_visibility="collapsed",
                )
                entry["positive"] = new_pos
            with c3b:
                st.caption("③❌ 反向 Prompt（不要什么）")
                cur_neg = entry.get("negative", "")
                new_neg = st.text_area(
                    f"negative_{idx}",
                    value=cur_neg,
                    height=240,
                    key=f"uni_neg_{idx}",
                    label_visibility="collapsed",
                    placeholder=(
                        "本页禁忌（每行一条），例如：\n"
                        "Tommy 戴眼镜\n"
                        "Mia 散发\n"
                        "除 Anna 外任何人穿绿色\n"
                        "画面出现宠物（本页不应有）"
                    ),
                )
                entry["negative"] = new_neg
            with c4:
                st.caption("④ 必须出现（追加到正向末尾，每行一条）")
                entry["must_include"] = st.text_area(
                    f"must_{idx}",
                    value=entry.get("must_include", ""),
                    height=180,
                    key=f"uni_must_{idx}",
                    label_visibility="collapsed",
                    placeholder=(
                        "（可选）写明必须画上的元素，例如：\n"
                        "Anna 头戴白色发箍\n"
                        "桌上有一摞 5 本课本\n"
                        "教室背景：暖米白空墙 + 单侧窗柔光 + 浅米瓷砖地"
                    ),
                )
                # Shot + 机位角度 + 表情
                bc1, bc2, bc3 = st.columns([1, 1, 1])
                with bc1:
                    cur_shot = ec_page.get("shot", "medium") or "medium"
                    new_shot = st.selectbox(
                        "镜头",
                        SHOT_OPTIONS,
                        index=max(0, SHOT_OPTIONS.index(cur_shot) if cur_shot in SHOT_OPTIONS else 1),
                        key=f"uni_shot_{idx}",
                    )
                    if new_shot != cur_shot:
                        ec_page["shot"] = new_shot
                with bc2:
                    cur_angle = ec_page.get("camera_angle", "") or "eye"
                    new_angle = st.selectbox(
                        "机位角度",
                        ANGLE_OPTIONS,
                        index=max(0, ANGLE_OPTIONS.index(cur_angle) if cur_angle in ANGLE_OPTIONS else 0),
                        format_func=lambda a: ANGLE_LABELS.get(a, a),
                        key=f"uni_angle_{idx}",
                    )
                    if new_angle != cur_angle:
                        ec_page["camera_angle"] = new_angle
                with bc3:
                    st.text_input(
                        "表情",
                        value=ec_page.get("expression", ""),
                        key=f"uni_expr_{idx}",
                        on_change=None,
                    )
                if refs:
                    st.caption(
                        "🎨 参考图: " + " | ".join(Path(r).name for r in refs[:4])
                    )

            st.divider()


def _render_editable_preview() -> None:
    """兼容旧调用：依次渲染 共享抽取 + RR + Worksheet（单页布局）。"""
    _render_shared_extract_panel()
    _render_rr_editor()
    _render_worksheet_editor()


def _render_shared_extract_panel() -> None:
    """共享：词汇/语法/拼读/读者类型（4 件套都用到，放在 Tab 上方或绘本轨内）。"""
    ec = st.session_state.extracted
    outline: BookOutline = st.session_state.outline

    # 词表
    with st.expander("📚 词汇表 / 语法 / 拼读 / 读者类型（4 件套共享）", expanded=True):
        col1, col2 = st.columns(2)
        is_dual = outline.is_dual_vocab_level
        with col1:
            if is_dual:
                ec.mastery = _csv_input("Mastery（3-4 词，lemma 原型小写）",
                                         ec.mastery, key="mastery")
                ec.exposure = _csv_input("Exposure（3-4 词，lemma 原型小写）",
                                          ec.exposure, key="exposure")
            else:
                ec.vocabulary = _csv_input("Vocabulary（4 词，lemma 原型小写）",
                                            ec.vocabulary, key="vocab")
        with col2:
            ec.grammar_focus = st.text_input(
                "Grammar Focus", value=ec.grammar_focus, key="grammar")
            ec.phonics = st.text_input(
                "Phonics", value=ec.phonics, key="phonics")
            _rt_opts = _reader_type_options(st.session_state.outline.level)
            if len(_rt_opts) == 1:
                ec.reader_type = _rt_opts[0]
                st.text_input("Reader Type", value=_rt_opts[0], disabled=True,
                              key="reader", help="L0-2 按 Level 固定体裁类，无需选择。")
            else:
                _norm = {"fiction": "Fiction", "non-fiction": "Non-Fiction",
                         "nonfiction": "Non-Fiction"}
                _cur = _norm.get((ec.reader_type or "").strip().lower(),
                                 (ec.reader_type or "").strip())
                if _cur not in _rt_opts:
                    _cur = _rt_opts[0]
                ec.reader_type = st.selectbox(
                    "Reader Type", _rt_opts, index=_rt_opts.index(_cur), key="reader",
                    help="L3-6 = Fiction / Non-Fiction。")
            ec.word_count = st.number_input(
                "Word Count", min_value=0, value=int(ec.word_count or 0), key="wc")

    # v2.0：unified_page_panel 已被 Step 4 + Step 5 取代，这里不再渲染


def _render_rr_editor() -> None:
    """📄 阅读报告 Tab：阅读表达题（按 Level 题量梯度）+ 预览/单题改。"""
    ec = st.session_state.extracted
    st.markdown(render_deliverable_spec_md("rr", st.session_state.outline.level))
    st.radio(
        "输出版本",
        ["空白版（教师手填）", "示例答案版（演示）"],
        horizontal=True,
        key="rr_answer_mode",
        help="示例答案版会在每题下方加灰色斜体示例答案；组装 4 件套时按此设置生成。",
    )
    # RR 题目
    with st.expander("📝 Reading Report 阅读表达题（按 Level 题量梯度）", expanded=True):
        for i, q in enumerate(ec.rr_questions):
            star = "⭐" * int(q.get("stars") or 1)
            cols = st.columns([10, 2])
            with cols[0]:
                q["q"] = st.text_input(
                    f"Q{i + 1} {star}", value=q.get("q", ""),
                    key=f"rr_q_{i}",
                )
            with cols[1]:
                if q.get("page") is None:
                    st.caption("(开放题)")
                    q["page"] = None
                else:
                    page_val = max(2, int(q.get("page") or 2))
                    q["page"] = st.number_input(
                        "P#", min_value=2, max_value=8, value=page_val,
                        key=f"rr_p_{i}", label_visibility="collapsed",
                        help="P1 = 封面，故事页是 P2-P8",
                    )

    # 即时预览：当前题目立刻生成一份 RR（严格 1 页 A4）并在网页内可视化预览
    st.markdown("---")
    col_pv1, col_pv2 = st.columns([3, 2])
    with col_pv1:
        st.caption(
            "👀 **生成前先看效果**：点右边按钮立即用当前阅读表达题生成一份 Reading Report，"
            "**严格锁定 1 张 A4**（系统会实测渲染页数，超一页自动降档重排）。"
        )
    with col_pv2:
        if st.button("👀 生成 Reading Report 预览", key="rr_preview_btn", width="stretch"):
            _generate_rr_preview()

    rr_draft = st.session_state.get("rr_draft_path")
    if rr_draft and Path(rr_draft).exists():
        st.success(f"✅ 初稿已生成：`{Path(rr_draft).name}`（约 {Path(rr_draft).stat().st_size//1024} KB）")
        _render_doc_preview(Path(rr_draft), "Reading Report", key="rr_draft", cols=1)
        with open(rr_draft, "rb") as f:
            st.download_button(
                "⬇️ 下载 Reading Report DOCX",
                data=f.read(),
                file_name=Path(rr_draft).name,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                key="rr_draft_dl",
            )


def _ws_sentence_mode() -> str:
    """读取 Sentence 页配图来源选择 → builder 参数。"""
    return "none" if st.session_state.get("ws_sentence_image_mode", "").startswith("不配图") else "reuse"


# 第 2 张 Reading 页内容：UI 文案 → builder 参数
_SECOND_READING_LABELS: dict[str, str] = {
    "自动（L0-2 思维导图 / L3-6 写作）": "auto",
    "思维导图 SWBST 复述": "mindmap",
    "写作脚手架（标题 Reading）": "writing",
    "写作 Writing（官方风格，标题 Writing）": "writing_official",
    "PBL 迷你项目": "pbl",
    "阅读理解延伸（再来一页阅读题）": "reading",
}


def _ws_second_reading_mode() -> str:
    """读取『第 2 张 Reading 页内容』选择 → builder 参数（默认 auto）。"""
    label = st.session_state.get("ws_second_reading_mode", "")
    return _SECOND_READING_LABELS.get(label, "auto")


# Reader Type 分级体裁（用户拍板）：L3-6 = Fiction/Non-Fiction；L0-2 = 固定体裁类
_L02_READER_TYPES = {
    "smart": "Concept & Knowledge-Building Readers",
    "0": "Concept & Knowledge-Building Readers",
    "1": "Patterned Narrative & Informational Readers",
    "2": "Early Independent Genre-Exposure Readers",
}


def _reader_type_options(level: str) -> list[str]:
    """按级别返回 Reader Type 候选。L3-6 给 Fiction/Non-Fiction 两选；L0-2 固定一类。"""
    digits = "".join(ch for ch in (level or "") if ch.isdigit())
    try:
        n = int(digits) if digits else 0
    except ValueError:
        n = 0
    if n >= 3:
        return ["Fiction", "Non-Fiction"]
    key = (level or "").lower()
    return [_L02_READER_TYPES.get(key,
            _L02_READER_TYPES.get(digits, "Early Independent Genre-Exposure Readers"))]


def _render_worksheet_editor() -> None:
    """📝 练习册 Tab：6 题(2 词汇+2 句型+2 阅读)逐题打磨 + 配图 + 预览。"""
    ec = st.session_state.extracted
    outline: BookOutline = st.session_state.outline
    st.markdown(render_deliverable_spec_md("worksheet", outline.level))
    st.radio(
        "Sentence 页配图来源",
        ["复用绘本图（默认）", "不配图"],
        horizontal=True,
        key="ws_sentence_image_mode",
        help="句子题默认复用绘本插画（page_02 起）；选『不配图』则选项满宽展开。",
    )
    st.selectbox(
        "第 2 张 Reading 页内容",
        list(_SECOND_READING_LABELS.keys()),
        index=0,
        key="ws_second_reading_mode",
        help="练习册固定 6 页：2 词汇 + 2 句型 + 2 阅读。第 1 张 Reading 是阅读理解，"
             "第 2 张可选思维导图/写作/PBL/阅读延伸（标题统一 Reading）。默认按级别。",
    )
    # Worksheet 题目
    with st.expander("📋 Worksheet 6 道题（题型 + 题项 JSON）", expanded=True):
        st.caption("题项是 JSON 格式（list of dict），如需自由编辑可直接改下面文本框。")

        # v2.0：Reading MC 页题数选择（4/6/8）
        cur_rqc = int(st.session_state.get("ws_reading_q_count", 4))
        st.session_state["ws_reading_q_count"] = st.select_slider(
            "📖 Reading MC 页题量（影响 Worksheet 第 4 页 reading 题量）",
            options=[4, 6, 8],
            value=cur_rqc,
            help=(
                "官方 L5-1 Worksheet 标准是 8 题。"
                "如果题目太多挤不下，可降到 6 或 4（默认 4）。"
            ),
            key="ws_rqc_slider",
        )

        # v1.9：初稿预览按钮（即时生成 worksheet.pptx 让老师看）
        st.markdown("---")
        col_pv1, col_pv2 = st.columns([3, 2])
        with col_pv1:
            st.caption(
                "👀 **生成前先看初稿**：点右边按钮会立即用当前编辑的题目生成一份 worksheet.pptx，"
                "你可以下载下来在 PowerPoint 里查看版式/字号/题目效果，确认无误再去步骤 C 组装。"
            )
        with col_pv2:
            if st.button(
                "👀 生成 Worksheet 初稿（预览）",
                key="ws_preview_btn",
                width="stretch",
            ):
                _generate_worksheet_preview()

        # 初稿下载链接
        ws_draft = st.session_state.get("ws_draft_path")
        if ws_draft and Path(ws_draft).exists():
            st.success(f"✅ 初稿已生成：`{Path(ws_draft).name}`（约 {Path(ws_draft).stat().st_size//1024} KB）")
            with open(ws_draft, "rb") as f:
                st.download_button(
                    "⬇️ 下载 Worksheet 初稿 PPTX 查看",
                    data=f.read(),
                    file_name=Path(ws_draft).name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="ws_draft_dl",
                )

        st.markdown("---")
        # v3.3 过程性出题：难度梯度说明（按 Level 自动匹配）
        _render_ws_difficulty_header(outline.level)
        st.caption(
            "💡 固定结构 **2 词汇 + 2 句子 + 2 阅读**。每道题可：① 换题型/难度 ② 🤖 AI 重出（紧扣故事+级别）"
            "③ 手动改 word / 句子 / 选项。改完点上方「👀 生成 Worksheet 初稿」看效果。"
        )
        for i, ws in enumerate(ec.worksheet_questions):
            qtype = (ws.get("type") or "").lower()
            section = _WS_SECTION_BY_SLOT[i] if i < len(_WS_SECTION_BY_SLOT) else "reading"
            sec_label = {"vocab": "📖 词汇", "sentence": "✍️ 句子", "reading": "📚 阅读"}[section]
            with st.container(border=True):
                st.markdown(f"**Activity {i + 1} · {sec_label} · {QUESTION_TITLES.get(qtype, qtype)}**")

                # ① 题型 / 难度切换（限定在本 section 内，保证版式不乱）
                tcol, bcol = st.columns([3, 2])
                with tcol:
                    _render_ws_type_switch(ws, i, section)
                with bcol:
                    st.markdown("&nbsp;", unsafe_allow_html=True)
                    if st.button(
                        "🤖 AI 重出这道题",
                        key=f"ws_regen_{i}",
                        width="stretch",
                        help="让 AI 按当前题型 + 故事内容 + 本级别难度，重新生成这道题的题项",
                    ):
                        _regenerate_one_worksheet_question(i)
                        st.rerun()

                # ② 结构化题项编辑
                _render_ws_item_editor(ws, i)


def _render_tg_panel() -> None:
    """👩‍🏫 教师指南 Tab：展示 8 模块规格 + 生成预览（DOCX）。"""
    outline: BookOutline = st.session_state.outline
    st.markdown(render_deliverable_spec_md("tg", outline.level))
    st.caption(
        "教师指南 100% 英文、按 8 固定模块顺序生成；Answer Key 与 Worksheet 同源一致。"
        "内容全部取自大纲，不编造。"
    )
    if st.button("👀 生成 Teacher's Guide 预览（DOCX）", key="tg_preview_btn"):
        try:
            run_dir, _, name_prefix = _ensure_run_dir()
            tg_path = run_dir / f"{name_prefix}_Teachers_Guide.docx"
            build_teacher_guide(outline, tg_path)
            st.session_state["tg_draft_path"] = str(tg_path)
        except Exception as e:
            st.error(f"生成失败：{e}")
    tg_draft = st.session_state.get("tg_draft_path")
    if tg_draft and Path(tg_draft).exists():
        st.success(f"✅ 已生成：`{Path(tg_draft).name}`（约 {Path(tg_draft).stat().st_size//1024} KB）")
        with open(tg_draft, "rb") as f:
            st.download_button(
                "⬇️ 下载 Teacher's Guide DOCX",
                data=f.read(),
                file_name=Path(tg_draft).name,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                key="tg_draft_dl",
            )


# ============================================================================
# v3.3 过程性出题：题型切换 / 单题 AI 重出 / 难度梯度说明
# ============================================================================

# 6 道题固定结构：2 词汇 + 2 句子 + 2 阅读
_WS_SECTION_BY_SLOT = ["vocab", "vocab", "sentence", "sentence", "reading", "reading"]

# 每个 section 允许切换的题型（限定在 builder 能正确排版的范围内）
# (legacy type id, 中文/英文短标签, 难度星级)
# v3.4：彻底删除涂色/纯涂圈/纯绘画类（无明确语言输出）；低龄改为看图类（有图、有词/句输出）。
_WS_TYPE_CHOICES: dict[str, list[tuple[str, str, int]]] = {
    "vocab": [
        ("word_to_pic", "看图选词 Look & Choose Word", 1),
        ("circle_match", "词↔图连线 Picture-Word Match", 1),
        ("fill_blank_simple", "看图填词 Fill the Word (simple)", 1),
        ("unscramble", "拼词 Unscramble Letters", 2),
        ("fill_blank", "填空 Fill Blanks", 2),
        ("emotion_fill", "选情绪词 Choose the Emotion", 3),
        ("fill_blank_advanced", "进阶填空 Fill Blanks (advanced)", 3),
        ("match_definition", "词↔义连线 Match Word ↔ Definition", 3),
    ],
    "sentence": [
        ("true_false_simple", "看图判断 True / False (simple)", 1),
        ("word_order_simple", "看图排词 Word Order (simple)", 1),
        ("true_false", "判断对错 True / False", 2),
        ("word_order", "句子排序 Sentence Order", 2),
        ("fill_blank", "看图填句 Complete the Sentence", 2),
        ("story_sequence", "故事排序 Story Sequence", 3),
        ("rewrite_tense", "句型改写 Rewrite (tense)", 3),
        ("rewrite_voice", "句型改写 Rewrite (voice/style)", 3),
    ],
    "reading": [
        ("personal_simple", "看图写我 About Me", 1),
        ("plot_chart", "故事图表 Story Chart", 2),
        ("inference", "阅读单选 Read & Infer (MCQ)", 3),
        ("plot_chart_pbl", "情节反思 Plot & Reflection", 3),
        ("compare_contrast", "对比 Compare & Contrast", 3),
        ("personal_write", "写你自己 Write About Yourself", 3),
        ("open_ended_pbl", "项目表达 Project Response", 3),
        ("essay_short", "短文写作 Short Essay", 3),
        ("research_pbl", "微型调研 Mini Research", 3),
    ],
}


def _render_ws_difficulty_header(level: str) -> None:
    """显示本 Level 的难度梯度（星级分布），让老师知道难度是按级别自动匹配的。"""
    try:
        from worksheet_question_types import _DIFFICULTY_PROFILE
        key = str(level or "5").strip().lower().lstrip("l").strip()
        if "smart" in str(level).lower():
            key = "smart"
        if key not in _DIFFICULTY_PROFILE:
            key = "5"
        n1, n2, n3 = _DIFFICULTY_PROFILE[key]
        st.info(
            f"🎚️ **Level {level} 难度梯度**：6 道题里约 "
            f"**{n1} 道 ⭐**（基础识记） · **{n2} 道 ⭐⭐**（理解运用） · **{n3} 道 ⭐⭐⭐**（推理/写作）。"
            "低级别更偏图/识记，高级别更偏推理/写作 —— 切题型时括号里的星级就是难度。"
        )
    except Exception:
        pass


def _render_ws_type_switch(ws: dict, idx: int, section: str) -> None:
    """本 section 内切换题型（带难度星级），换型时同步 title/instruction。"""
    choices = list(_WS_TYPE_CHOICES.get(section, []))
    cur_type = (ws.get("type") or "").lower()
    ids = [c[0] for c in choices]
    if cur_type and cur_type not in ids:
        # 当前题型不在候选里（兼容旧数据）→ 补进去当首选
        choices = [(cur_type, QUESTION_TITLES.get(cur_type, cur_type.replace("_", " ").title()), 2)] + choices
        ids = [c[0] for c in choices]

    labels = [f"{lbl} ({'⭐' * stars})" for _tid, lbl, stars in choices]
    cur_index = ids.index(cur_type) if cur_type in ids else 0
    sel = st.selectbox(
        "题型 / 难度",
        options=list(range(len(choices))),
        index=cur_index,
        format_func=lambda k: labels[k],
        key=f"ws_type_{idx}",
        help="只在本类（词汇/句子/阅读）里换，保证 Worksheet 版式不乱。括号内为难度。",
    )
    new_type = ids[sel]
    if new_type != cur_type:
        ws["type"] = new_type
        ws["title"] = QUESTION_TITLES.get(new_type, new_type.replace("_", " ").title())
        ws["instruction"] = QUESTION_INSTRUCTIONS.get(new_type, "")
        # 换题型后建议点「🤖 AI 重出这道题」重新出题项


def _regenerate_one_worksheet_question(idx: int) -> None:
    """逐题 AI 重出：用故事原文 + 当前题型 + 级别，重新生成这道题的题项。"""
    ec = st.session_state.extracted
    outline: BookOutline = st.session_state.outline
    if not ec or idx >= len(ec.worksheet_questions):
        return
    ws = ec.worksheet_questions[idx]
    qtype = (ws.get("type") or "").lower()
    raw_story = "\n".join((p.get("text") or "") for p in (ec.pages or []) if p.get("text"))
    title = getattr(outline, "title", "") or ""
    level = getattr(outline, "level", "") or "5"
    theme = getattr(outline, "theme", "") or ""
    cefr = getattr(outline, "cefr", "") or ""
    with st.spinner(f"🤖 正在按「{QUESTION_TITLES.get(qtype, qtype)}」+ 故事 + Level {level} 重出这道题..."):
        try:
            new_q = generate_one_worksheet_question(
                qtype, raw_story, title, level, cefr=cefr, theme=theme,
            )
            ec.worksheet_questions[idx] = new_q
            st.session_state.extracted = ec
            st.toast("✅ 已重出这道题，下面可继续微调。", icon="🤖")
        except Exception as e:  # noqa: BLE001
            st.error(f"重出失败：{e}")


def _render_ws_item_editor(ws: dict, idx: int) -> None:
    """v1.9：按 worksheet 题型展开成结构化输入（替代原 JSON 编辑器）。

    题型 → 结构：
      match_definition → 5 行 (word, def)
      fill_blank / fill_blank_advanced / emotion_fill → 5 行 (sentence, answer)
      true_false → 4 行 (statement, T/F)
      inference / reading_mc → 4 行 (q, options[3], correct)
      unscramble → 5 行 (scrambled, answer)
      其他 → 简化的 items JSON
    """
    qtype = (ws.get("type") or "").lower()
    # 健壮归一化：items 可能是 list / dict（{idx: item} 映射或单个 item dict）/ 其它
    _raw_items = ws.get("items")
    if isinstance(_raw_items, dict):
        _vals = list(_raw_items.values())
        # 全是子 dict → 当作 {idx: item} 映射取 values；否则视为单个 item dict
        items = _vals if (_vals and all(isinstance(v, dict) for v in _vals)) else [_raw_items]
    elif isinstance(_raw_items, list):
        items = _raw_items
    else:
        items = []

    if qtype == "match_definition":
        st.caption("每行：单词 → 词典定义（小学生能看懂的简单解释）")
        new_items = []
        for j, it in enumerate(items[:5]):
            c1, c2 = st.columns([1, 4])
            with c1:
                w = st.text_input(f"word_{idx}_{j}", value=it.get("word", ""),
                                   key=f"ws{idx}_w{j}", label_visibility="collapsed")
            with c2:
                d = st.text_input(f"def_{idx}_{j}",
                                   value=it.get("def") or it.get("definition", ""),
                                   key=f"ws{idx}_d{j}", label_visibility="collapsed",
                                   placeholder="feeling worried about something that will happen")
            if w.strip():
                new_items.append({"word": w.strip(), "def": d.strip()})
        ws["items"] = new_items

    elif qtype in ("fill_blank", "fill_blank_simple", "fill_blank_advanced", "emotion_fill"):
        st.caption("每行：句子（用 ____ 留空） → 答案")
        new_items = []
        for j, it in enumerate(items[:5]):
            c1, c2 = st.columns([4, 1])
            with c1:
                s = st.text_input(f"sent_{idx}_{j}", value=it.get("sentence", ""),
                                   key=f"ws{idx}_s{j}", label_visibility="collapsed",
                                   placeholder="Anna ____ on her first day.")
            with c2:
                a = st.text_input(f"ans_{idx}_{j}", value=it.get("answer", ""),
                                   key=f"ws{idx}_a{j}", label_visibility="collapsed",
                                   placeholder="felt nervous")
            if s.strip():
                new_items.append({"sentence": s.strip(), "answer": a.strip()})
        ws["items"] = new_items
        ws["extra"] = st.text_input(
            f"词库 extra_{idx}", value=ws.get("extra", ""), key=f"ws{idx}_e",
            placeholder="逗号分隔的词库，如：nervous, shared, helped, listened",
        )

    elif qtype in ("true_false", "true_false_simple"):
        st.caption("每行：陈述句 + T/F")
        new_items = []
        for j, it in enumerate(items[:4]):
            c1, c2 = st.columns([5, 1])
            with c1:
                s = st.text_input(f"stmt_{idx}_{j}", value=it.get("statement", ""),
                                   key=f"ws{idx}_st{j}", label_visibility="collapsed",
                                   placeholder="Anna shared pencils with a quiet boy.")
            with c2:
                a = st.selectbox(f"tf_{idx}_{j}", ["T", "F"],
                                  index=0 if (it.get("answer") or "T").upper() == "T" else 1,
                                  key=f"ws{idx}_tf{j}", label_visibility="collapsed")
            if s.strip():
                new_items.append({"statement": s.strip(), "answer": a})
        ws["items"] = new_items

    elif qtype in ("inference", "reading_mc"):
        st.caption("每行：问题 + 3 个选项 + 正确答案 (A/B/C)")
        new_items = []
        for j, it in enumerate(items[:4]):
            q = st.text_input(f"infq_{idx}_{j}", value=it.get("q") or it.get("question", ""),
                               key=f"ws{idx}_q{j}",
                               label_visibility="visible",
                               placeholder="Why did Anna help pick up the books?")
            opts = list(it.get("options") or ["", "", ""])
            while len(opts) < 3:
                opts.append("")
            cc1, cc2, cc3, cc4 = st.columns([3, 3, 3, 1])
            with cc1:
                opts[0] = st.text_input(f"o0_{idx}_{j}", value=opts[0],
                                          key=f"ws{idx}_o0{j}", label_visibility="collapsed",
                                          placeholder="A. She wanted to be kind.")
            with cc2:
                opts[1] = st.text_input(f"o1_{idx}_{j}", value=opts[1],
                                          key=f"ws{idx}_o1{j}", label_visibility="collapsed",
                                          placeholder="B. The teacher told her.")
            with cc3:
                opts[2] = st.text_input(f"o2_{idx}_{j}", value=opts[2],
                                          key=f"ws{idx}_o2{j}", label_visibility="collapsed",
                                          placeholder="C. She wanted the books.")
            with cc4:
                cur = int(it.get("correct", 0))
                cur = max(0, min(2, cur))
                ans = st.selectbox(f"co_{idx}_{j}", ["A", "B", "C"], index=cur,
                                    key=f"ws{idx}_co{j}", label_visibility="collapsed")
            if q.strip():
                new_items.append({
                    "q": q.strip(),
                    "options": [o.strip() for o in opts],
                    "correct": {"A": 0, "B": 1, "C": 2}[ans],
                })
        ws["items"] = new_items

    elif qtype == "unscramble":
        st.caption("每行：打散字母 → 正确答案")
        new_items = []
        for j, it in enumerate(items[:5]):
            c1, c2 = st.columns([2, 2])
            with c1:
                sc = st.text_input(f"sc_{idx}_{j}", value=it.get("scrambled", ""),
                                    key=f"ws{idx}_sc{j}", label_visibility="collapsed",
                                    placeholder="o c k l")
            with c2:
                a = st.text_input(f"sca_{idx}_{j}", value=it.get("answer", ""),
                                    key=f"ws{idx}_sca{j}", label_visibility="collapsed",
                                    placeholder="lock")
            if sc.strip():
                new_items.append({"scrambled": sc.strip(), "answer": a.strip()})
        ws["items"] = new_items

    else:
        # 兜底（compare_contrast、plot_chart、写作类等）：保留 JSON 但收起
        with st.expander("📝 items JSON（高级题型用 JSON 编辑）", expanded=False):
            import json
            items_str = st.text_area(
                f"items_json_{idx}",
                value=json.dumps(items, ensure_ascii=False, indent=2),
                height=120, key=f"ws_json_{idx}",
                label_visibility="collapsed",
            )
            try:
                ws["items"] = json.loads(items_str)
            except json.JSONDecodeError:
                st.warning("items JSON 解析失败，保留旧值")
        ws["extra"] = st.text_input(
            f"extra_{idx}", value=ws.get("extra", ""), key=f"ws_e_{idx}",
            placeholder="题目附加说明（如反思问题、写作 prompt）",
        )


def _csv_input(label: str, words: list[str], *, key: str) -> list[str]:
    raw = st.text_input(label, value=", ".join(words or []), key=key)
    return [w.strip().lower() for w in raw.split(",") if w.strip()]


def _name_prefix(outline: BookOutline) -> str:
    """新规范：Level X_BookXX_品类_标题  →  品类 + 后缀由调用方拼接。

    返回去除非法字符（/\\:*?"<>| 及空格替换为下划线）的前缀串。
    """
    level_s = outline.level or "1"
    if "smart" in level_s.lower():
        lvl_part = "Smart"
    else:
        digits = "".join(ch for ch in level_s if ch.isdigit())
        lvl_part = f"Level {digits or '1'}"
    book_s = (outline.book_number or "01").strip()
    if not book_s.lower().startswith("book"):
        book_s = f"Book{book_s}"
    title = re.sub(r'[\\/:*?"<>|]', "_", (outline.title or "Untitled"))
    title = re.sub(r"_+", "_", title).strip("_ ")
    return f"{lvl_part}_{book_s}_{title}"


# =============================== 生成器 ===============================
def _ensure_run_dir() -> tuple[Path, Path, str]:
    """复用同一次 run 目录（生图 + 组装文档共用），避免每次重新生图都换路径。"""
    outline: BookOutline = st.session_state.outline
    if not st.session_state.get("run_dir"):
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        run_dir = OUTPUTS_DIR / f"{outline.slug}_{timestamp}"
        st.session_state.run_dir = str(run_dir)
        st.session_state.run_name_prefix = _name_prefix(outline)
    run_dir = Path(st.session_state.run_dir)
    img_dir = run_dir / "images"
    run_dir.mkdir(parents=True, exist_ok=True)
    img_dir.mkdir(parents=True, exist_ok=True)
    return run_dir, img_dir, st.session_state.run_name_prefix


def _build_final_prompt_for_page(page, outline: BookOutline, ip_age: int) -> tuple[str, list[Path]]:
    """v3.1：组装最终 prompt =
       (Step3 风格设定正向 + page positive + 教师必须出现)
       + ==请勿出现==
       + (page negative + Step3 风格设定反向)

    优先用老师在 UI 上编辑过的正向/反向；否则用 build_cn_page_prompt 现场生成。
    Step3 的全局风格设定如果存在，会自动**前置**到正向 + **追加**到反向。
    """
    from cn_prompt_builder import BuiltPromptCN

    style_cfg = st.session_state.get("style_config") or {}
    style_blocks = _format_style_config_preview(style_cfg) if style_cfg else {
        "positive_block": "", "negative_block": ""
    }
    style_pos = style_blocks.get("positive_block", "").strip()
    style_neg = style_blocks.get("negative_block", "").strip()

    page_prompts = st.session_state.get("page_prompts") or {}
    edited = page_prompts.get(page.index)
    if edited:
        positive = edited.get("positive") or edited.get("prompt") or ""
        negative = edited.get("negative") or ""

        # v3.2 A 层：注入 page_constraints（姿势/情绪/视线/位置/互动/文字位置/焦点）
        pc_block = _format_page_constraints(edited.get("page_constraints"))
        if pc_block:
            positive = positive.rstrip() + "\n\n【A 层 · 本页人物状态约束】\n" + pc_block

        must_inc = (edited.get("must_include") or "").strip()
        if must_inc:
            positive = (
                positive.rstrip()
                + "\n\n【教师锁定 · 必须出现】\n"
                + must_inc
            )
        refs = [Path(r) for r in edited.get("references", []) if r]
    else:
        cast_pool = st.session_state.get("story_cast_pool") or None
        generic_overrides = st.session_state.get("generic_overrides") or None
        built = build_cn_page_prompt(
            page, outline, ip_age,
            cast_pool=cast_pool, generic_overrides=generic_overrides,
        )
        positive = built.positive if hasattr(built, "positive") else built.prompt
        negative = built.negative if hasattr(built, "negative") else ""
        refs = built.references

    # v3.1 注入 Step 3 风格设定
    if style_pos:
        positive = f"【全局风格设定】\n{style_pos}\n\n" + positive
    if style_neg:
        negative = (negative.rstrip() + "\n" if negative else "") + style_neg

    # v3.2 反馈学习的负向（来自每页生图后的问题反馈面板）
    if edited:
        learned_neg = (edited.get("learned_negatives") or "").strip()
        if learned_neg:
            negative = (negative.rstrip() + "\n" if negative else "") + learned_neg

    # v3.3：主角恒为「首位参考图」。gpt-image-2 只吃单张参考，必须保证它是主角，
    # 否则主角（尤其用 she/he 代称、本句没点名时）会跳帧。
    prot_ref = _book_protagonist_ref(outline, ip_age)
    if prot_ref is not None and page.page_type != "cover":
        refs = [Path(r) for r in refs]
        if prot_ref in refs:
            refs.remove(prot_ref)
        refs.insert(0, prot_ref)

    # v4 支柱一：多角色页「合成定妆参考图」——gpt-image-2 单参考图限制下，
    # 把本页所有角色定妆图拼成 1 张白底合集作唯一参考，锁住每个人的长相/发型/服装，
    # 解决「只发主角参考、Tommy/Mia 等配角崩形」的根因。
    refs = _maybe_build_reference_sheet(refs, page)
    if isinstance(refs, tuple):
        refs, sheet_note = refs
        if sheet_note:
            positive = positive.rstrip() + "\n\n" + sheet_note

    final_prompt = BuiltPromptCN.join(positive, negative)
    return final_prompt, refs


def _ip_name_for_ref(ref_path: Path) -> str:
    """把一张参考图路径反查成 IP 名字（用于定妆合集的标签）。"""
    try:
        from ip_library import load_library
        rp = str(Path(ref_path)).lower()
        for e in load_library():
            if str(e.image_path).lower() == rp:
                return e.name_base
    except Exception:
        pass
    return ""


def _maybe_build_reference_sheet(refs: list[Path], page):
    """v6 图生图 IP 锁定（用户拍板 2026-06-06：连续性绘本，形象永久锁定，用图不用文字）：
    本页出场的【每一个 IP】都必须由其官方定妆图驱动，绝不靠文字瞎编。

      · 单角色页 → 把那张定妆图贴身裁切成干净单锚图作唯一参考；
      · 多角色页 → 把本页所有角色的官方定妆图横向拼成一张「角色定妆表」作唯一参考，
        一张图同时锁住每个人的脸型/发型/发色/服装/配色。

    并配套强约束 note：照定妆表 1:1 还原每个人的形象，【只允许改变姿势与表情】，
    严禁改动形象，严禁照搬定妆表的白底/并排排版/站姿/多视图。

    返回 ([anchor_or_sheet], prompt_note)；无可用本地图时原样返回 refs（不带 note）。
    """
    local = [Path(r) for r in refs if r and not str(r).startswith(("http://", "https://")) and Path(r).exists()]
    if not local:
        return refs
    try:
        from seedream_client import crop_character_portrait, build_reference_sheet
        _, img_dir, _ = _ensure_run_dir()
        anchor_dir = img_dir / "_anchors"
        names = [(_ip_name_for_ref(r) or "") for r in local]

        if len(local) >= 2:
            # 多角色：拼定妆表（用图锁住每个人）
            dest = anchor_dir / f"_refsheets" / f"sheet_p{page.index:02d}.png"
            sheet = build_reference_sheet(local, dest, labels=names)
            if sheet is not None:
                disp = "、".join(n for n in names if n) or "本页角色"
                note = (
                    "【参考图＝角色定妆表（连续性绘本·形象永久锁定）】所附唯一参考图是一张白底"
                    f"「角色定妆表」，并排展示本页所有出场角色（{disp}，每个面板上方有英文名标签）。"
                    "这是一套【已出版的连续性绘本】的官方既定形象，必须严格沿用：请把定妆表里每个角色的"
                    "脸型、五官、发型、发色、肤色、服装款式与配色、鞋子等形象细节，1:1 精确还原到画面中，"
                    "做到与定妆表完全一致、与往期绘本同一个人。"
                    "【唯一允许改变的是：每个角色的姿势/动作 与 面部表情】，用来贴合本页剧情；"
                    "其余形象一律不得改动、不得重新设计、不得换装换发型换配色。"
                    "严禁照搬定妆表的白色背景、并排站姿、正/侧/背多视图排版——只取‘谁长什么样’，"
                    "把他们自然放进本页描述的真实场景里。"
                )
                return ([sheet], note)

        # 单角色（或拼图失败兜底）：贴身裁切单锚
        primary = local[0]
        dest = anchor_dir / f"anchor_p{page.index:02d}.png"
        anchor = crop_character_portrait(primary, dest) or primary
        primary_name = names[0] or "主角"
        note = (
            f"【参考图＝{primary_name} 官方定妆图（连续性绘本·形象永久锁定）】所附唯一参考图是 "
            f"{primary_name} 的单人定妆图（贴身裁切、干净背景），来自一套已出版的连续性绘本。"
            f"请把 {primary_name} 的脸型、五官、发型、发色、肤色、服装款式与配色 1:1 精确还原，"
            "做到与定妆图完全一致、与往期绘本同一个人。"
            "【唯一允许改变的是该角色的姿势/动作与面部表情】以贴合本页剧情，"
            "形象其余部分一律不得改动；不要照搬参考图的背景或姿势。"
        )
        return ([anchor], note)
    except Exception:
        return refs


@st.cache_data(show_spinner=False)
def _book_protagonist_ref_cached(cast_keys: tuple[str, ...], all_text: str, title: str, ip_age: int) -> str:
    """识别全书主角的参考图路径（字符串）。主角 = cast_pool 里 kind=protagonist
    且名字在标题/正文出现最多的那个。返回 '' 表示无。
    """
    from ip_library import get_ip
    best_key, best_score = "", -1
    text_low = (title + " " + all_text).lower()
    for key in cast_keys:
        ip = get_ip(key)
        if not ip or ip.kind != "protagonist":
            continue
        name = ip.name_base.lower()
        score = text_low.count(name)
        if score > best_score:
            best_key, best_score = key, score
    if not best_key:
        return ""
    ip = get_ip(best_key)
    return str(ip.image_path) if ip and ip.image_path else ""


def _book_protagonist_ref(outline: BookOutline, ip_age: int) -> Path | None:
    cast_pool = st.session_state.get("story_cast_pool") or []
    if not cast_pool:
        return None
    all_text = " ".join((p.text or "") for p in outline.pages)
    p = _book_protagonist_ref_cached(tuple(cast_pool), all_text, outline.title or "", ip_age)
    return Path(p) if p else None


def _run_image_generation_only(mock_images: bool) -> None:
    """v1.8.3 步骤 A：只生图，不组装文档。生成完把结果存到 session_state.image_results 供审核。"""
    ec = st.session_state.extracted
    outline: BookOutline = st.session_state.outline
    apply_extracted_to_outline(outline, ec)
    attach_rr_questions(outline, ec.rr_questions)
    rqc = int(st.session_state.get("ws_reading_q_count", 4))
    attach_worksheet_questions(outline, ec.worksheet_questions, reading_q_count=rqc)

    run_dir, img_dir, _ = _ensure_run_dir()
    ip_age = outline.ip_age or resolve_ip_age(outline.level)

    progress = st.progress(0, "Initializing image generation...")
    status = st.empty()

    n = len(outline.pages)
    image_results: dict[int, dict] = st.session_state.get("image_results") or {}

    # 并发生图：模型单张 20-60 秒去不掉，但多张同时跑可整体缩短到 ~原时间/并发数。
    # 主线程先算好每页 prompt/refs/dest（子线程不碰 st/session_state），再丢线程池并发。
    concurrency = max(1, int(os.getenv("IMAGE_CONCURRENCY", "2")))

    # 块2（用户拍板 2026-06-08）：同一地点连续页【图链锚定】——
    #   段内首页=锚点，先生成；其后页把锚点成图作为【场景参考图】传入，让走廊/房间真正一致。
    try:
        from cn_prompt_builder import _location_runs
        _runs = _location_runs(outline)
    except Exception:
        _runs = {}
    # dep_anchor[page_index] = 锚点页 index（该页布景应跟随锚点成图）
    dep_anchor: dict[int, int] = {}
    for _idx, _info in (_runs or {}).items():
        _lo = _info[2]
        if _idx != _lo:
            dep_anchor[_idx] = _lo

    pages_to_gen = []
    skipped_locked = 0
    for page in outline.pages:
        prev = image_results.get(page.index, {})
        # 块7（用户拍板 2026-06-08）：批量重生时跳过【已锁定 ✅】的图，只重生未锁定的，锁定图原样保留。
        if prev.get("locked") and prev.get("path") and Path(prev["path"]).exists():
            skipped_locked += 1
            continue
        pages_to_gen.append(page)

    if skipped_locked:
        status.text(f"🔒 跳过 {skipped_locked} 张已锁定的图，只重生其余 {len(pages_to_gen)} 张。")
    if not pages_to_gen:
        progress.progress(1.0, "全部图已锁定 ✅，无需重生")
        status.empty()
        st.info("🔒 所有图都已锁定，没有需要重生的页。如需重生请先在下方解锁对应图。")
        st.session_state.image_results = image_results
        return
    n = len(pages_to_gen)

    # 把待生成页分两批：第1批=锚点+非连续页（可并发）；第2批=同地点跟随页（等锚点成图后再生，带场景参考）
    gen_idx = {p.index for p in pages_to_gen}
    first_pages = [p for p in pages_to_gen if p.index not in dep_anchor or dep_anchor[p.index] not in gen_idx]
    dep_pages = [p for p in pages_to_gen if p.index in dep_anchor and dep_anchor[p.index] in gen_idx]

    def _scene_ref_for(page) -> list[Path]:
        """跟随页：取锚点页已生成的成图作为额外【场景参考图】。"""
        a = dep_anchor.get(page.index)
        if a is None:
            return []
        rec = image_results.get(a) or {}
        p = rec.get("path")
        if p and Path(p).exists():
            return [Path(p)]
        return []

    def _build_task(page):
        prev = image_results.get(page.index, {})
        final_prompt, refs = _build_final_prompt_for_page(page, outline, ip_age)
        # 场景参考图排在角色参考之后，整体不超过 3 张（角色一致性优先，场景锚定其次）
        scene_refs = _scene_ref_for(page)
        if scene_refs:
            refs = (list(refs) + scene_refs)[:3]
        version = int(prev.get("version", 0)) + 1
        dest = img_dir / f"page_{page.index:02d}_v{version}.png"
        return (page, final_prompt, refs, version, dest)

    errors: dict[int, str] = {}
    done = 0
    start_ts = time.time()

    def _work(t):
        page, final_prompt, refs, version, dest = t
        generate_image_for_level(
            outline.level,
            prompt=final_prompt, dest=dest,
            references=refs, mock=mock_images, label=page.label,
            deliver_print=False,  # 审图阶段不放大；定稿组装时只对锁定图放大
            review_meta={
                "page_text": getattr(page, "text", "") or "",
                "ip_age": getattr(outline, "ip_age", None),
            },
        )
        return t

    def _run_pass(tasks, pass_label: str):
        nonlocal done
        if not tasks:
            return
        workers = min(concurrency, max(1, len(tasks)))
        status.text(f"⏳ {pass_label}：提交 {len(tasks)} 张，{workers} 张并发…（单张约 30–60 秒）")
        with ThreadPoolExecutor(max_workers=workers) as ex:
            fut_map = {ex.submit(_work, t): t for t in tasks}
            pending = set(fut_map)
            while pending:
                finished, pending = wait(pending, timeout=3, return_when=FIRST_COMPLETED)
                for fut in finished:
                    page, final_prompt, refs, version, dest = fut_map[fut]
                    done += 1
                    display_name = page_display_name(page.index)
                    try:
                        fut.result()
                    except Exception as e:
                        errors[page.index] = str(e)
                        try:
                            with (run_dir / "image_errors.log").open("a", encoding="utf-8") as _lf:
                                _lf.write(f"[{display_name} / {page.label}] {type(e).__name__}: {e}\n")
                        except Exception:
                            pass
                        continue
                    prev = image_results.get(page.index, {})
                    image_results[page.index] = {
                        "path": str(dest), "prompt": final_prompt, "label": page.label,
                        "version": version, "locked": prev.get("locked", False),
                    }
                elapsed = int(time.time() - start_ts)
                progress.progress(min(done, n) / n, f"已完成 {min(done, n)}/{n}")
                status.text(
                    f"⏳ {pass_label} · 已用 {elapsed}s · 完成 {min(done, n)}/{n}"
                    + (f"（{len(errors)} 张失败）" if errors else "")
                    + f" · {len(pending)} 张进行中…"
                )

    progress.progress(0.0, "生图中... 0/%d" % n)
    # 第1批：锚点 + 非连续页
    first_tasks = [_build_task(p) for p in first_pages]
    _run_pass(first_tasks, "第1批·锚点/独立页")
    # 第2批：同地点跟随页（此时锚点已成图 → 作为场景参考注入）
    dep_tasks = [_build_task(p) for p in dep_pages]
    if dep_tasks:
        _run_pass(dep_tasks, "第2批·同场景跟随页(锚定布景)")

    # ★ 失败页自动补跑：并发突发常因限流/图床抖动丢几张，这里串行（单张、退避）重试。
    all_tasks = first_tasks + dep_tasks
    task_by_idx = {t[0].index: t for t in all_tasks}
    for _retry_round in range(2):
        if not errors:
            break
        failed_idx = list(errors.keys())
        for ri, page_index in enumerate(failed_idx):
            t = task_by_idx.get(page_index)
            if not t:
                continue
            page, final_prompt, refs, version, dest = t
            display_name = page_display_name(page.index)
            status.text(
                f"🔁 补跑失败页（第 {_retry_round + 1} 轮）：{display_name} "
                f"· 剩 {len(failed_idx) - ri}"
            )
            try:
                generate_image_for_level(
                    outline.level,
                    prompt=final_prompt, dest=dest,
                    references=refs, mock=mock_images, label=page.label,
                    deliver_print=False,
                )
            except Exception as e:
                try:
                    with (run_dir / "image_errors.log").open("a", encoding="utf-8") as _lf:
                        _lf.write(f"[RETRY{_retry_round+1} {display_name}] {type(e).__name__}: {e}\n")
                except Exception:
                    pass
                continue
            prev = image_results.get(page.index, {})
            image_results[page.index] = {
                "path": str(dest), "prompt": final_prompt, "label": page.label,
                "version": version, "locked": prev.get("locked", False),
            }
            errors.pop(page_index, None)

    st.session_state.image_results = image_results
    status.empty()
    if errors:
        msg = "；".join(f"{page_display_name(k)}: {v}" for k, v in errors.items())
        progress.progress(1.0, f"完成 {n - len(errors)}/{n}（{len(errors)} 张失败）")
        st.error(
            f"已自动补跑仍有 {len(errors)} 张失败（可单张 🔁 重生）：{msg}。"
            f"若反复失败多为接口限流，可把并发调小：环境变量 IMAGE_CONCURRENCY=2。"
        )
    else:
        progress.progress(1.0, f"✅ {n} 张图全部生成")
        st.success(
            f"✅ 已生成 {n} 张图到 `{img_dir}`。"
            f"下方按页审核 —— 不满意的可点 🔁 单张重生；满意的勾 ✅ 锁定。"
        )


def _regenerate_single_image(page_index: int, mock_images: bool) -> None:
    """单张重生（用最新的 prompt + must_include）。"""
    outline: BookOutline = st.session_state.outline
    ip_age = outline.ip_age or resolve_ip_age(outline.level)
    _, img_dir, _ = _ensure_run_dir()

    page = next((p for p in outline.pages if p.index == page_index), None)
    if not page:
        st.error(f"找不到页 index={page_index}")
        return
    display_name = page_display_name(page_index)

    final_prompt, refs = _build_final_prompt_for_page(page, outline, ip_age)
    image_results = st.session_state.image_results
    prev = image_results.get(page_index, {})
    version = int(prev.get("version", 0)) + 1
    dest = img_dir / f"page_{page_index:02d}_v{version}.png"
    try:
        run_with_live_timer(
            f"重新生成 {display_name} (v{version})", generate_image,
            prompt=final_prompt, dest=dest,
            references=refs, mock=mock_images, label=page.label,
        )
    except Exception as e:
        st.error(f"{display_name} 重生失败：{e}")
        return
    image_results[page_index] = {
        "path": str(dest),
        "prompt": final_prompt,
        "label": page.label,
        "version": version,
        "locked": False,
    }
    st.session_state.image_results = image_results
    st.success(f"✅ {display_name} 已重新生成（v{version}）")


def _review_fix_single(page_index: int, mock_images: bool) -> None:
    """块3：对单张图跑 GPT 视觉自审；审出硬伤则图生图定向修，存为新版本。结果写入 session 供面板展示。"""
    reviews = st.session_state.get("image_reviews") or {}
    image_results = st.session_state.get("image_results") or {}
    entry = image_results.get(page_index) or {}
    img_path = entry.get("path")
    if not img_path or not Path(img_path).exists():
        st.error("图片不存在，无法审图")
        return
    outline: BookOutline = st.session_state.outline
    page = next((p for p in outline.pages if p.index == page_index), None)
    ip_age = outline.ip_age or resolve_ip_age(outline.level)
    if mock_images:
        st.info("占位图模式不审图。")
        return
    try:
        from seedream_client import _review_image, _fix_prompt, generate_image as _gen, host_image_to_url
    except Exception as e:
        st.error(f"审图模块不可用：{e}")
        return
    with st.spinner(f"GPT 审图中… {page_display_name(page_index)}"):
        verdict = _review_image(
            Path(img_path),
            page_text=(getattr(page, "text", "") or "") if page else "",
            scene_cn=(getattr(page, "scene_cn", "") or "") if page else "",
            ip_age=ip_age,
        )
    verdict = dict(verdict or {})
    verdict["fixed"] = False
    if not verdict.get("ok", True) and (verdict.get("issues") or verdict.get("fix")):
        # 定向修图（图生图，保留画风/构图）→ 新版本
        _, img_dir, _ = _ensure_run_dir()
        version = int(entry.get("version", 0)) + 1
        dest = img_dir / f"page_{page_index:02d}_v{version}.png"
        try:
            ref_url = host_image_to_url(Path(img_path))
            with st.spinner("按问题定向修图中…"):
                _gen(prompt=_fix_prompt(verdict.get("issues") or [], verdict.get("fix") or ""),
                     dest=dest, reference_url=ref_url, mock=False,
                     label=f"{page_display_name(page_index)} 定向修图", deliver_print=False)
            image_results[page_index] = {
                "path": str(dest), "prompt": entry.get("prompt", ""),
                "label": entry.get("label", ""), "version": version,
                "locked": False,
            }
            st.session_state.image_results = image_results
            verdict["fixed"] = True
        except Exception as e:
            st.error(f"定向修图失败：{e}")
    reviews[page_index] = verdict
    st.session_state.image_reviews = reviews


def _render_image_review_panel(mock_imgs: bool) -> None:
    """生图后展示 8 张缩略图，每张可：🔁 重生 / ✏️ 改 prompt 重生 / ✅ 锁定。"""
    image_results = st.session_state.image_results
    n_total = len(image_results)
    n_locked = sum(1 for r in image_results.values() if r.get("locked"))
    st.caption(
        f"已生成 {n_total} 张图，已锁定 ✅ {n_locked} 张。"
        "建议每张都点 🔁 重生几次直到满意，再勾 ✅ 锁定，所有图锁定后再组装文档。"
    )

    # v3.3 图在前：每页一行，左图（大）右信息（故事原文 / 剧情要点 / 操作 / 提示词折叠）
    outline = st.session_state.get("outline")
    pages_by_idx = {p.index: p for p in (outline.pages if outline else [])}

    indices = sorted(image_results.keys())
    for idx in indices:
        entry = image_results[idx]
        page = pages_by_idx.get(idx)
        display_name = page_display_name(idx)
        lock_emoji = "🔒 已锁定" if entry.get("locked") else ""
        st.markdown(
            f"#### {display_name} · {entry.get('label', '')} · v{entry.get('version', 1)}　{lock_emoji}"
        )

        c_img, c_info = st.columns([3, 2])
        with c_img:
            img_path = entry.get("path")
            if img_path and Path(img_path).exists():
                _zoom_image(
                    img_path,
                    key=f"pg{idx}v{entry.get('version', 1)}",
                    caption=f"{display_name} · v{entry.get('version', 1)}",
                )
            else:
                st.caption("⚠️ 图片不存在")

        with c_info:
            # 图在前，文字在旁：该页故事原文 + 剧情要点
            story_text = (page.text or "").strip() if page else ""
            plot = (getattr(page, "scene_cn", "") or "").strip() if page else ""
            if story_text:
                st.markdown(f"**📖 该页故事**\n\n{story_text}")
            else:
                st.markdown("**📖 该页**：封面")
            if plot:
                st.caption(f"🎬 剧情要点：{plot}")

            bcol1, bcol2 = st.columns(2)
            with bcol1:
                if st.button(
                    "🔁 重生",
                    key=f"regen_btn_{idx}",
                    width="stretch",
                    help="用当前 prompt 重新生成一次（保留参考图）",
                ):
                    _regenerate_single_image(idx, mock_imgs)
                    st.rerun()
            with bcol2:
                new_lock = st.checkbox(
                    "✅ 锁定",
                    value=bool(entry.get("locked")),
                    key=f"lock_chk_{idx}",
                    help="勾上表示这张图满意，组装时用它",
                )
                if new_lock != bool(entry.get("locked")):
                    entry["locked"] = new_lock
                    st.session_state.image_results[idx] = entry

            # 块3：GPT 审图面板 —— 一键审 + 按需定向修（展示问题 + 一键修图）
            if st.button("🩺 GPT 审图并按需修复", key=f"review_btn_{idx}", width="stretch",
                         help="让 GPT 看这张图，挑 IP不一致/多指畸形/图文不符/分身/家具比例等硬伤；有问题就图生图定向修，只修瑕疵不改画风"):
                _review_fix_single(idx, mock_imgs)
                st.rerun()
            _rv = (st.session_state.get("image_reviews") or {}).get(idx)
            if _rv is not None:
                if _rv.get("ok", True) and not _rv.get("issues"):
                    st.success("✅ 审图通过：未发现硬伤")
                else:
                    st.warning("🔍 审出问题：" + "；".join(str(i) for i in (_rv.get("issues") or [])))
                    if _rv.get("fixed"):
                        st.caption("🔧 已按问题定向修图，生成新版本。")

            # 提示词折叠（默认收起，不刷屏；像即梦：图为主，提示词点开才看）
            with st.expander("✏️ 提示词（点开编辑 + 重生）", expanded=False):
                page_prompts = st.session_state.get("page_prompts") or {}
                pp = page_prompts.get(idx, {})
                cur_prompt = pp.get("prompt", entry.get("prompt", ""))
                cur_must = pp.get("must_include", "")
                new_prompt = st.text_area(
                    "中文 Prompt",
                    value=cur_prompt,
                    height=180,
                    key=f"reedit_prompt_{idx}",
                )
                new_must = st.text_area(
                    "必须出现（追加到 prompt 末尾）",
                    value=cur_must,
                    height=70,
                    key=f"reedit_must_{idx}",
                    placeholder="例如：\nAnna 必须头戴白色发箍\nTommy 必须微笑看向 Anna",
                )
                if st.button(f"💾 保存并重生 {display_name}", key=f"reedit_save_{idx}"):
                    pp["prompt"] = new_prompt
                    pp["must_include"] = new_must
                    page_prompts[idx] = pp
                    st.session_state.page_prompts = page_prompts
                    _regenerate_single_image(idx, mock_imgs)
                    st.rerun()

            # v3.2 问题反馈面板：勾选问题 → 自动加进 learned_negatives，下次重生立即生效
            _render_image_feedback_panel(idx, mock_imgs)

        st.divider()


# v3.2 问题反馈面板的固定问题清单（问题 → 自动注入到反向 prompt 的文本）
_IMAGE_FEEDBACK_ITEMS = [
    ("frame_jump", "🔀 人物跳帧（外观和前后页不一致）",
     "人物外观与故事其他页不一致；服装/发型/配饰跳变；同一角色出现不同面孔"),
    ("multi_finger", "✋ 手指/手部畸形（多手指、缺手指、融合）",
     "多手指、缺手指、融合手指、第六根手指、手部畸形、手腕扭曲"),
    ("char_too_small", "🔍 人物太小（淹没在背景里）",
     "人物在画面中过小、淹没在背景、远景小人；主角必须占画面 30-40% 以上"),
    ("char_too_big", "🔭 人物太大（占满画面没有环境）",
     "人物占满画面、没有环境上下文、特写过头；必须保留环境留白"),
    ("background_messy", "🌪️ 背景杂乱（元素太碎太多）",
     "背景元素过多过碎、视觉混乱、注意力被背景抢走；背景要简洁"),
    ("ratio_wrong", "📐 角色间比例不对（成人/儿童/宠物大小关系错）",
     "角色之间比例失调；成年人比儿童矮、儿童比宠物小、宠物大如人；比例必须真实"),
    ("gaze_wrong", "👀 视线混乱（角色没看对地方）",
     "角色视线方向错误；对话中角色没看对方；眼神涣散；应当根据互动调整视线"),
    ("clothing_mismatch", "👕 服装与故事/IP 不符",
     "角色服装与 IP 锁定不符；同一本书内服装跳变；服装颜色错误"),
    ("emotion_wrong", "😐 表情和故事情绪不匹配",
     "角色表情与故事情境不符；该笑的没笑、该紧张的太放松"),
    ("extra_chars", "👥 多了不该出现的人物 / 少了关键角色",
     "画面里出现故事没有提到的人物；或者缺失故事里关键的角色"),
]


def _render_image_feedback_panel(idx: int, mock_imgs: bool) -> None:
    """v3.2 问题反馈面板：勾选这张图的问题 → 自动加进 learned_negatives → 一键重生。"""
    page_prompts = st.session_state.setdefault("page_prompts", {})
    entry = page_prompts.setdefault(idx, {})
    saved_codes: list[str] = entry.get("feedback_codes") or []

    with st.expander("🚨 这张图有问题？勾选 → 下次自动避免", expanded=False):
        st.caption("勾选的问题会作为反向 prompt 写进下次生图。你也可以补一句自由文本。")
        new_codes: list[str] = []
        for code, label, _ in _IMAGE_FEEDBACK_ITEMS:
            checked = st.checkbox(
                label,
                value=(code in saved_codes),
                key=f"fb_{idx}_{code}",
            )
            if checked:
                new_codes.append(code)

        free_text = st.text_input(
            "其他问题（自由写一句，自动加进反向）",
            value=entry.get("feedback_free", ""),
            key=f"fb_free_{idx}",
            placeholder="例：Anna 的发箍丢了，应该头戴白色发箍",
        )

        # 实时拼出"会注入的负向 prompt"预览
        neg_bits: list[str] = []
        for code, _, neg in _IMAGE_FEEDBACK_ITEMS:
            if code in new_codes:
                neg_bits.append(neg)
        if free_text.strip():
            neg_bits.append(free_text.strip())
        learned_neg = "；".join(neg_bits)

        if learned_neg:
            st.markdown("**🔬 将注入的反向 prompt：**")
            st.code(learned_neg, language="text")

        bb1, bb2 = st.columns(2)
        with bb1:
            if st.button("💾 保存反馈（不重生）", key=f"fb_save_{idx}", width="stretch"):
                entry["feedback_codes"] = new_codes
                entry["feedback_free"] = free_text
                entry["learned_negatives"] = learned_neg
                page_prompts[idx] = entry
                st.session_state.page_prompts = page_prompts
                st.success("已保存。下次任何形式的重生都会带上这些反向约束。")
        with bb2:
            if st.button("🔁 保存并重生本页", key=f"fb_regen_{idx}", width="stretch", type="primary"):
                entry["feedback_codes"] = new_codes
                entry["feedback_free"] = free_text
                entry["learned_negatives"] = learned_neg
                page_prompts[idx] = entry
                st.session_state.page_prompts = page_prompts
                _regenerate_single_image(idx, mock_imgs)
                st.rerun()


def _run_docs_assembly() -> None:
    """v1.8.3 步骤 C：用 session 里的图组装 4 件套 + 打包 ZIP。"""
    _assemble_t0 = time.time()
    ec = st.session_state.extracted
    outline: BookOutline = st.session_state.outline
    apply_extracted_to_outline(outline, ec)
    attach_rr_questions(outline, ec.rr_questions)
    rqc = int(st.session_state.get("ws_reading_q_count", 4))
    attach_worksheet_questions(outline, ec.worksheet_questions, reading_q_count=rqc)

    image_results = st.session_state.get("image_results") or {}
    if not image_results:
        st.error("还没生图。请先点上面的「🎨 生成所有图」。")
        return

    run_dir, _, name_prefix = _ensure_run_dir()

    # 按 page.index 顺序取图（用最新的 path）
    image_paths: list[Path] = []
    for page in outline.pages:
        entry = image_results.get(page.index)
        if not entry or not entry.get("path") or not Path(entry["path"]).exists():
            st.error(f"P{page.index} 缺图，请先点 🎨 生成。")
            return
        image_paths.append(Path(entry["path"]))

    # ★ 定稿放大（质量优先）：审图阶段为了快没做 4K 放大，这里只对最终选用的图
    #   统一做「居中裁 4:3 → 升印刷分辨率」。放大不改画面内容，质量无损、专供印刷 DPI。
    from seedream_client import postprocess_4k
    up = st.progress(0, "印刷放大（4K）...")
    for _k, _p in enumerate(image_paths):
        up.progress(_k / max(1, len(image_paths)), f"印刷放大 {_k + 1}/{len(image_paths)}...")
        try:
            postprocess_4k(_p)
        except Exception:
            pass
    up.progress(1.0, "✅ 印刷放大完成")

    progress = st.progress(0, "组装文档中...")

    # 1. Picture Book PPT
    progress.progress(1 / 5, "组装 Picture Book PPT...")
    pb_path = run_dir / f"{name_prefix}_Reader.pptx"
    build_picturebook_pptx(outline, image_paths, pb_path)

    # 2. Worksheet
    progress.progress(2 / 5, "生成 Worksheet PPTX...")
    ws_path = run_dir / f"{name_prefix}_Worksheet.pptx"
    build_worksheet(outline, ws_path, image_paths=image_paths,
                    sentence_image_mode=_ws_sentence_mode(),
                    second_reading_mode=_ws_second_reading_mode())

    # 3. Reading Report
    progress.progress(3 / 5, "生成 Reading Report DOCX...")
    rr_path = run_dir / f"{name_prefix}_Reading_Report.docx"
    rr_with_answers = st.session_state.get("rr_answer_mode", "").startswith("示例")
    build_reading_report(outline, rr_path, with_answers=rr_with_answers)

    # 4. Teacher's Guide
    progress.progress(4 / 5, "生成 Teacher's Guide DOCX...")
    tg_path = run_dir / f"{name_prefix}_Teachers_Guide.docx"
    build_teacher_guide(outline, tg_path)

    # 5. ZIP
    zip_path = run_dir / f"{name_prefix}_Full_Set.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:
        for path in [pb_path, ws_path, rr_path, tg_path]:
            z.write(path, arcname=path.name)
        # 也带上图片（屏幕用 PNG）
        for img in image_paths:
            z.write(img, arcname=f"images/{img.name}")

    progress.progress(1.0, "完成 ✅")

    _elapsed = time.time() - _assemble_t0
    st.success(f"4 件套已生成（组装用时 {_elapsed:.1f}s）。点击下面按钮下载：")
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        _download_button(pb_path, "📘 绘本 PPT")
    with col2:
        _download_button(ws_path, "📋 Worksheet")
    with col3:
        _download_button(rr_path, "📝 Reading Report")
    with col4:
        _download_button(tg_path, "📖 Teacher Guide")
    with col5:
        _download_button(zip_path, "📦 全套 ZIP", primary=True)

    # 4 交付物视觉预览（文档→PDF→PNG，按需渲染，带缓存）
    _render_deliverable_previews(pb_path, ws_path, rr_path, tg_path)

    # 预览图片（瀑布流 / Masonry 图片墙，错落铺满无大块空白）
    with st.expander("🖼️ 预览生成的图片", expanded=False):
        _masonry_gallery(
            [str(img) for img in image_paths],
            captions=[f"第 {i} 页" if i else "封面" for i in range(len(image_paths))],
        )


def _render_doc_preview(path: Path, label: str, *, key: str, cols: int = 2) -> None:
    """交付物视觉预览：首次点按钮渲染（文档→PDF→PNG），之后走缓存秒开。

    soffice/fitz 不可用时退回文字预览。
    """
    path = Path(path)
    if not path.exists():
        st.error(f"{label}: 文件不存在")
        return
    sk = f"_docprev::{key}::{path.name}"
    if st.button(f"🔍 预览 {label}", key=f"btn_{sk}", width="stretch"):
        if has_visual_preview():
            try:
                imgs = run_with_live_timer(
                    f"渲染预览 · {label}", render_to_images, path,
                    done_note="首次约 10 秒，之后秒开",
                )
            except Exception as e:
                imgs = []
                st.warning(f"视觉预览失败：{e}")
            st.session_state[sk] = [str(p) for p in imgs]
        else:
            st.session_state[sk] = []  # 标记“已尝试”，走文字兜底

    if sk not in st.session_state:
        st.caption("👆 点按钮生成预览（首次约 10 秒渲染，之后秒开）。")
        return

    imgs = st.session_state.get(sk) or []
    if not imgs:
        st.info("未生成视觉预览，改用文字预览：")
        st.text(extract_text(path))
        return
    _masonry_gallery(
        imgs,
        captions=[f"{label} · 第 {i + 1} 页" for i in range(len(imgs))],
        cols=cols,
    )


def _render_deliverable_previews(
    pb_path: Path, ws_path: Path, rr_path: Path, tg_path: Path,
) -> None:
    """组装完成后：4 交付物视觉预览（分页签，按需渲染）。"""
    with st.expander("👁️ 预览 4 交付物（绘本 / Worksheet / RR / TG）", expanded=False):
        if not has_visual_preview():
            st.caption("⚠️ 未检测到 LibreOffice，预览将以文字形式呈现。")
        t1, t2, t3, t4 = st.tabs(
            ["📘 绘本 Reader", "📋 Worksheet", "📝 Reading Report", "📖 Teacher Guide"]
        )
        with t1:
            _render_doc_preview(pb_path, "绘本 Reader", key="pb", cols=2)
        with t2:
            _render_doc_preview(ws_path, "Worksheet", key="ws", cols=2)
        with t3:
            _render_doc_preview(rr_path, "Reading Report", key="rr", cols=1)
        with t4:
            _render_doc_preview(tg_path, "Teacher Guide", key="tg", cols=2)


def _download_button(path: Path, label: str, *, primary: bool = False) -> None:
    if not path.exists():
        st.error(f"{label}: 文件不存在")
        return
    with open(path, "rb") as f:
        st.download_button(
            label=label,
            data=f.read(),
            file_name=path.name,
            mime="application/octet-stream",
            type="primary" if primary else "secondary",
            width="stretch",
        )


# =============================== 工具 ===============================
def _sync_ec_from_syllabus(ec, outline) -> None:
    """命中大纲时，把官方【词表 / 拼读 / 语法】逐字同步进 ec（工作台单一真值源）。

    工作台 UI 读的是 ec，而 enrich_from_syllabus 只改 outline；不同步会出现「明明命中大纲，
    工作台仍显示 AI 抽的词」。这里以大纲 verbatim 覆盖 ec，让显示/编辑/产出一致。
    """
    syl = getattr(outline, "syllabus", None)
    if syl is None:
        return
    try:
        if outline.is_dual_vocab_level:
            if getattr(syl, "vocab_mastery", None):
                ec.mastery = list(syl.vocab_mastery)
            if getattr(syl, "vocab_exposure", None):
                ec.exposure = list(syl.vocab_exposure)
        else:
            words = syl.vocab_words()
            if words:
                ec.vocabulary = list(words)
        if getattr(syl, "phonics_rule", ""):
            ph = syl.phonics_rule
            if getattr(syl, "phonics_examples", ""):
                ph = f"{ph}: {syl.phonics_examples}"
            ec.phonics = ph
        # 语法：大纲若有句型/句法焦点，用作 Grammar Focus（否则保留 AI 值）
        gram = getattr(syl, "syntax_focus", "") or getattr(syl, "sentence_pattern", "")
        if gram:
            ec.grammar_focus = gram
    except Exception:
        pass


def _build_outline(
    *, title: str, level: str, book_number: str, cefr: str, theme: str,
    ip_age: int, raw_story: str, custom_chars_text: str,
    fiction_type: str = "",
) -> BookOutline:
    """从表单数据搭一个 BookOutline 骨架（pages 留空，等 AI 抽取后填）。"""
    pages: list[PageSpec] = [PageSpec(index=0, page_type="cover", text=title, scene="")]
    for i in range(1, 8):
        pages.append(PageSpec(index=i, page_type="story", text="", scene=""))

    custom_chars: dict[str, str] = {}
    for line in (custom_chars_text or "").splitlines():
        if "|" in line:
            name, desc = line.split("|", 1)
            name = name.strip().lower()
            desc = desc.strip()
            if name and desc:
                custom_chars[name] = desc

    return BookOutline(
        title=title.strip() or "Picture Book",
        pages=pages,
        level=level,
        book_number=book_number,
        cefr=cefr,
        theme=theme,
        ip_age=ip_age,
        custom_characters=custom_chars,
        fiction_type=fiction_type,
    )


def _key_status_banner() -> None:
    ai_ok = bool(DOUBAO_API_KEY) and not MOCK_AI_EXTRACT
    img_ok = bool(JIMENG_API_KEY) and not MOCK_IMAGES

    msgs = []
    if ai_ok:
        msgs.append("🟢 Claude 文本可用")
    else:
        msgs.append("🟡 Claude 文本走 mock（无 key 或强制 mock）")
    if img_ok:
        msgs.append("🟢 gpt-image-2 图生成可用")
    else:
        msgs.append("🟡 gpt-image-2 走 mock 占位图")
    st.caption("&nbsp;&nbsp;|&nbsp;&nbsp;".join(msgs), unsafe_allow_html=True)


def _masonry_gallery(paths, captions=None, cols: int | None = None) -> None:
    """瀑布流（Masonry）图片墙：图片按列流式排布、按内容高度自适应、列间无大块空白。

    用 base64 内联 <img>，包进 .masonry 容器（CSS column 多列）→ Pinterest 式错落瀑布。
    每张图悬浮有 3D 抬升 + 标题浮层（微交互反馈）。
    cols：固定列数（不传则用 CSS 响应式 3/2/1 列）。"""
    paths = [str(p) for p in (paths or [])]
    if not paths:
        return
    cards, boxes = [], []
    for i, p in enumerate(paths):
        uri = _img_data_uri(p)
        if not uri:
            continue
        zid = "zi-" + re.sub(r"[^a-zA-Z0-9_-]", "", f"{i}-{Path(p).stem}")
        cap_txt = captions[i] if (captions and i < len(captions) and captions[i]) else ""
        cap = f'<figcaption class="mz-cap">{cap_txt}</figcaption>' if cap_txt else ""
        # 卡片本体即放大触发器（点击 → 全屏灯箱）
        cards.append(
            f'<a class="mz-card" href="#{zid}"><img loading="lazy" src="{uri}"/>{cap}</a>'
        )
        lbcap = f'<div class="lb-cap">{cap_txt}</div>' if cap_txt else ""
        boxes.append(
            f'<a class="lightbox" id="{zid}" href="#_close"><img src="{uri}"/>{lbcap}</a>'
        )
    if not cards:
        return
    style = f' style="column-count:{cols}"' if cols else ""
    st.markdown(f'<div class="masonry"{style}>{"".join(cards)}</div>{"".join(boxes)}',
                unsafe_allow_html=True)


def _img_data_uri(path) -> str | None:
    """读图 → base64 data URI（失败返回 None）。"""
    ext_mime = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                ".webp": "image/webp", ".gif": "image/gif"}
    try:
        data = Path(path).read_bytes()
    except Exception:
        return None
    mime = ext_mime.get(Path(path).suffix.lower(), "image/png")
    return f"data:{mime};base64,{base64.b64encode(data).decode('ascii')}"


def _zoom_image(path, key: str, caption: str = "") -> None:
    """单图点击放大（纯 CSS 灯箱）：点缩略图 → 全屏看大图；点任意处关闭。无需任何按钮。"""
    uri = _img_data_uri(path)
    if not uri:
        st.caption("⚠️ 图片不存在")
        return
    zid = "zi-" + re.sub(r"[^a-zA-Z0-9_-]", "", str(key))
    cap = f'<div class="lb-cap">{caption}</div>' if caption else ""
    st.markdown(
        f'<a class="zoomable" href="#{zid}"><img loading="lazy" src="{uri}"/></a>'
        f'<a class="lightbox" id="{zid}" href="#_close"><img src="{uri}"/>{cap}</a>',
        unsafe_allow_html=True,
    )


def _inject_css() -> None:
    st.markdown(
        """<style>
        @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@400;500;600;700&family=Plus+Jakarta+Sans:wght@500;600;700;800&display=swap');
        :root{
          --brand:#006b58; --brand-dark:#004a3c; --brand-tint:#e8faf5;
          --brand-2:#26c2a3; --primary-fixed:#6ff9d8;
          --ink:#0b1c30; --muted:#3c4a45; --faint:#6c7a75;
          --line:#dce9ff; --line-soft:#eff4ff; --card:#ffffff;
          --bg:#f8f9ff;
          --radius:14px; --radius-lg:20px;
          --shadow-sm:0 1px 2px rgba(11,28,48,.04),0 1px 3px rgba(11,28,48,.05);
          --shadow-md:0 4px 12px rgba(11,28,48,.06),0 2px 4px rgba(11,28,48,.04);
          --shadow-lg:0 18px 40px -12px rgba(0,107,88,.15);
          --ring:0 0 0 3px rgba(0,107,88,.16);
          --fs-base:15.5px; --fs-sm:13.5px; --fs-lg:17px;
        }
        html, body, [class*="css"]{
          -webkit-font-smoothing:antialiased; text-rendering:optimizeLegibility;
          font-size:var(--fs-base);
        }
        .stApp{ font-size:var(--fs-base); }
        .stApp{
          background:
            radial-gradient(900px 480px at 88% -8%, rgba(77,221,188,.12), transparent 60%),
            radial-gradient(820px 460px at -6% 4%, rgba(0,107,88,.05), transparent 55%),
            var(--bg);
        }
        .block-container { max-width: 1340px; padding-top: 1.2rem; }

        /* ---------- 隐藏侧边栏（交付物导航已移至主区） ---------- */
        [data-testid="stSidebar"],
        [data-testid="stSidebarCollapsedControl"],
        [data-testid="collapsedControl"]{
          display:none !important;
        }

        /* ---------- 顶栏 + 主导航 pill（PRD 风格） ---------- */
        #app-header-anchor + div[data-testid="stVerticalBlockBorderWrapper"],
        #app-header-anchor + div[data-testid="stVerticalBlock"] > div[data-testid="stVerticalBlockBorderWrapper"]{
          border:1px solid var(--line) !important;
          border-radius:var(--radius-lg) !important;
          padding:10px 16px 6px !important;
          margin-bottom:18px !important;
          background:linear-gradient(120deg, rgba(255,255,255,.96), rgba(232,250,245,.88)) !important;
          box-shadow:var(--shadow-md) !important;
        }
        .app-topbar-brand{ display:flex; align-items:center; gap:12px; min-width:0; }
        .app-topbar-title{
          font-family:'Plus Jakarta Sans','Noto Sans SC',sans-serif; font-size:20px; font-weight:800;
          color:var(--brand); line-height:1.2;
        }
        .nav-dino-logo{ width:40px; height:40px; object-fit:contain; }
        .nav-dino-fallback{ font-size:32px; line-height:1; }
        .app-topbar-sub{ font-size:12px; color:var(--muted); margin-top:2px; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
        .main-nav-wrap [data-testid="stRadio"]{ width:100%; }
        .main-nav-wrap [role="radiogroup"]{
          display:flex; flex-wrap:wrap; justify-content:flex-end; align-items:center;
          gap:2px; padding:4px; border-radius:999px;
          background:rgba(18,24,38,.04); border:1px solid var(--line);
        }
        .main-nav-wrap input[type="radio"]{ position:absolute !important; opacity:0 !important; width:0 !important; height:0 !important; pointer-events:none !important; }
        .main-nav-wrap [role="radiogroup"] > label{
          margin:0 !important; padding:6px 12px !important; border-radius:999px !important;
          border:none !important; background:transparent !important; box-shadow:none !important;
          font-weight:600 !important; font-size:12.5px !important; color:var(--muted) !important;
          transition:all .18s ease !important; cursor:pointer !important;
          min-height:0 !important; gap:0 !important;
        }
        .main-nav-wrap [role="radiogroup"] > label[data-baseweb="radio"] > div:first-child,
        .main-nav-wrap [role="radiogroup"] > label > div:first-child{
          display:none !important; width:0 !important; min-width:0 !important; margin:0 !important; padding:0 !important;
        }
        .main-nav-wrap [role="radiogroup"] > label:hover{
          color:var(--brand) !important; background:rgba(0,107,88,.06) !important;
        }
        .main-nav-wrap [role="radiogroup"] > label:has(input:checked){
          color:var(--brand) !important;
          background:linear-gradient(135deg, rgba(0,107,88,.14), rgba(38,194,163,.10)) !important;
          box-shadow:0 1px 6px rgba(0,107,88,.16), inset 0 0 0 1px rgba(0,107,88,.20) !important;
        }
        .main-nav-wrap [role="radiogroup"] > label > div:last-child,
        .main-nav-wrap [role="radiogroup"] > label p{
          padding:0 !important; margin:0 !important; font-size:12.5px !important;
        }
        .page-lead h2{ font-size:28px; font-weight:800; margin:0 0 8px; color:var(--ink); }
        .page-lead p{ color:var(--muted); font-size:var(--fs-lg); margin:0 0 16px; line-height:1.6; }
        @media (max-width: 1100px){
          .main-nav-wrap [role="radiogroup"] > label{ padding:5px 9px !important; font-size:11.5px !important; }
        }
        @media (max-width: 900px){
          .main-nav-wrap [role="radiogroup"]{ justify-content:flex-start; }
          .app-topbar-sub{ white-space:normal; }
        }
        .header-actions .stButton > button{
          white-space:nowrap !important;
          min-width:4.2rem;
          font-size:12.5px !important;
          padding:0.38rem 0.55rem !important;
          line-height:1.2 !important;
        }
        @media (max-width: 1100px){
          .header-actions .stButton > button{
            font-size:11.5px !important;
            padding:0.34rem 0.45rem !important;
            min-width:3.6rem;
          }
        }
        .nav-title-inline{
          font-weight:700; font-size:13px; letter-spacing:.4px; color:var(--faint); margin:.4rem 0 .5rem;
        }
        [data-testid="stRadio"] [role="radiogroup"][aria-label="交付物导航"]{
          gap:8px; flex-wrap:wrap;
        }
        [data-testid="stRadio"] [role="radiogroup"][aria-label="交付物导航"] label{
          padding:8px 16px; border-radius:999px; border:1px solid var(--line);
          font-weight:600; background:var(--card); transition:all .15s ease;
        }
        [data-testid="stRadio"] [role="radiogroup"][aria-label="交付物导航"] label:has(input:checked){
          color:var(--brand); background:var(--brand-tint); border-color:rgba(0,107,88,.35);
          box-shadow:inset 0 0 0 1px rgba(0,107,88,.12);
        }
        [data-testid="stRadio"] [role="radiogroup"][aria-label="交付物导航"] label > div:first-child{ display:none; }

        body, p, span, div, label, textarea, input, button, select{
          font-family:'Plus Jakarta Sans','Noto Sans SC','PingFang SC','Microsoft YaHei',sans-serif;
          font-size:var(--fs-base);
        }
        [data-testid="stCaption"], .stCaption, small{ font-size:var(--fs-sm) !important; }
        [data-testid="stMetricValue"]{ font-size:1.35rem !important; }
        [data-testid="stMetricLabel"]{ font-size:var(--fs-sm) !important; }

        /* ---------- Hero（mockup：mint 渐变 + Kidde + 特性 pill） ---------- */
        .hero-section{
          border-radius:var(--radius-lg); margin:0 0 24px; overflow:hidden;
          padding:32px 28px;
        }
        .hero-gradient{
          background:
            radial-gradient(circle at top right, rgba(77,221,188,.15), transparent),
            radial-gradient(circle at bottom left, rgba(0,107,88,.05), transparent);
        }
        .hero-grid{
          display:grid; grid-template-columns:1.05fr .95fr; gap:28px; align-items:center;
        }
        @media (max-width:900px){ .hero-grid{ grid-template-columns:1fr; } }
        .hero-pill{
          display:inline-flex; align-items:center; gap:6px;
          padding:4px 12px; border-radius:999px; margin-bottom:14px;
          background:var(--primary-fixed); color:#002019;
          font-size:12px; font-weight:600;
        }
        .hero-headline{
          font-family:'Plus Jakarta Sans','Noto Sans SC',sans-serif;
          font-size:clamp(32px,4vw,48px); font-weight:800; line-height:1.15;
          color:var(--ink); margin:0 0 16px;
        }
        .hero-accent{ color:var(--brand); }
        .hero-lead{
          font-size:18px; line-height:1.55; color:var(--muted);
          max-width:540px; margin:0 0 20px;
        }
        .hero-features{ display:flex; flex-wrap:wrap; gap:12px; }
        .hero-feat{
          display:flex; align-items:center; gap:10px;
          padding:12px 14px; border-radius:12px;
          background:#fff; border:1px solid #bbcac4; box-shadow:var(--shadow-sm);
          min-width:160px;
        }
        .feat-ic{ font-size:20px; }
        .feat-t{ font-size:14px; font-weight:600; color:var(--ink); margin:0; }
        .feat-s{ font-size:13px; color:var(--muted); margin:0; }
        .hero-visual{ position:relative; }
        .hero-glow{
          position:absolute; top:-40px; right:-40px; width:220px; height:220px;
          border-radius:50%; background:#4dddbc; filter:blur(60px); opacity:.2;
        }
        .hero-glass{
          position:relative; padding:8px; border-radius:14px;
          background:rgba(255,255,255,.8); backdrop-filter:blur(12px);
          border:1px solid #bbcac4; box-shadow:var(--shadow-lg); overflow:hidden;
        }
        .hero-kidde{
          width:100%; max-height:360px; object-fit:contain;
          border-radius:10px; display:block; background:#fff;
        }
        .hero-kidde-fallback{ font-size:80px; display:block; text-align:center; padding:40px; }
        .hero-glass-cap{
          position:absolute; left:20px; right:20px; bottom:20px;
          display:flex; justify-content:space-between; align-items:center;
          padding:12px 14px; border-radius:10px;
          background:rgba(255,255,255,.85); backdrop-filter:blur(10px);
          border:1px solid rgba(255,255,255,.5);
        }
        .cap-t{ font-size:14px; font-weight:600; color:var(--brand); margin:0; }
        .cap-s{ font-size:13px; color:var(--ink); margin:0; }
        .cap-play{ color:var(--brand); font-size:22px; }

        /* ---------- 生成绘本区块头 ---------- */
        .create-section{ margin:0 0 8px; }
        .create-head{
          display:flex; align-items:center; gap:14px; margin-bottom:16px;
        }
        .create-icon{
          width:48px; height:48px; border-radius:12px;
          background:var(--brand); color:#fff; display:flex;
          align-items:center; justify-content:center; font-size:22px;
        }
        .create-head h2{
          font-size:32px; font-weight:700; color:var(--ink); margin:0 0 4px;
        }
        .create-head p{ font-size:16px; color:var(--muted); margin:0; }

        /* ---------- FAQ 手风琴 ---------- */
        .faq-section{
          margin:28px 0 20px; padding:28px 0;
          background:var(--bg); border-radius:var(--radius-lg);
        }
        .faq-head{ text-align:center; margin-bottom:22px; }
        .faq-head h2{
          font-size:32px; font-weight:700; color:var(--ink); margin:0 0 6px;
        }
        .faq-head p{ color:var(--muted); margin:0; }
        .faq-list{ max-width:880px; margin:0 auto; display:flex; flex-direction:column; gap:12px; }
        .faq-item{
          background:#fff; border:1px solid #bbcac4; border-radius:12px; overflow:hidden;
        }
        .faq-item summary{
          display:flex; justify-content:space-between; align-items:center;
          padding:18px 20px; cursor:pointer; list-style:none;
        }
        .faq-item summary::-webkit-details-marker{ display:none; }
        .faq-q{ font-size:14px; font-weight:600; color:var(--ink); }
        .faq-chevron{ color:var(--muted); transition:transform .2s; font-size:12px; }
        .faq-item[open] .faq-chevron{ transform:rotate(180deg); }
        .faq-a{
          padding:0 20px 18px; font-size:15px; line-height:1.6;
          color:var(--muted); border-top:1px solid #bbcac4; padding-top:14px;
        }

        /* ---------- 页脚 ---------- */
        .site-footer{
          display:flex; flex-wrap:wrap; justify-content:space-between; align-items:center;
          gap:16px; margin-top:32px; padding:28px 0 12px;
          border-top:1px solid #bbcac4;
        }
        .footer-logo{ font-size:14px; font-weight:700; color:var(--brand); }
        .site-footer p{ font-size:13px; color:var(--muted); margin:4px 0 0; }
        .footer-links{ display:flex; flex-wrap:wrap; gap:20px; }
        .footer-links span{
          font-size:12px; font-weight:500; color:var(--muted); cursor:default;
        }

        /* ---------- legacy hero 兼容（登录页等） ---------- */
        .hero{
          position:relative; display:flex; align-items:center; gap:20px;
          background:linear-gradient(120deg, rgba(255,255,255,.86), rgba(232,250,245,.78));
          backdrop-filter:blur(14px) saturate(1.2); -webkit-backdrop-filter:blur(14px) saturate(1.2);
          border:1px solid rgba(255,255,255,.7); border-radius:var(--radius-lg);
          padding:22px 26px; margin:4px 0 20px;
          box-shadow:var(--shadow-lg);
          overflow:hidden;
        }
        .hero::before{
          content:""; position:absolute; inset:0; z-index:0; pointer-events:none;
          background:
            radial-gradient(420px 220px at 90% -40%, rgba(0,107,88,.22), transparent 60%),
            radial-gradient(360px 200px at 8% 130%, rgba(38,194,163,.18), transparent 60%);
        }
        .hero::after{
          content:""; position:absolute; left:0; right:0; bottom:0; height:3px; z-index:1;
          background:linear-gradient(90deg, transparent, var(--brand), var(--brand-2), transparent);
          opacity:.8;
        }
        .hero > *{ position:relative; z-index:2; }
        .hero-icon{ flex:0 0 auto; display:flex; }
        .hero-dino{
          width:64px; height:64px; object-fit:contain;
          filter:drop-shadow(0 6px 14px rgba(0,107,88,.34));
        }
        .hero-text{ flex:1 1 auto; min-width:0; }
        .hero-title{
          font-family:'Poppins','Inter',sans-serif;
          font-size:28px; font-weight:800; letter-spacing:.2px;
          background:linear-gradient(92deg,var(--ink) 0%, var(--ink) 38%, var(--brand) 100%);
          -webkit-background-clip:text; background-clip:text; color:transparent;
          line-height:1.18;
        }
        .hero-sub{ color:var(--muted); font-size:13.5px; margin-top:5px; font-weight:500; }
        .hero-badge{
          flex:0 0 auto; align-self:flex-start;
          background:linear-gradient(135deg,var(--brand),var(--brand-2)); color:#fff;
          font-weight:700; font-size:12px; letter-spacing:.3px;
          padding:6px 14px; border-radius:999px;
          box-shadow:0 6px 16px rgba(0,107,88,.32); border:1px solid rgba(255,255,255,.35);
        }

        /* ---------- 侧边栏（玻璃质感 + 细描边） ---------- */
        [data-testid="stSidebar"]{
          background:linear-gradient(180deg,#ffffff, #fbfcfe);
          border-right:1px solid var(--line);
        }
        .side-brand{
          display:flex; align-items:center; gap:11px;
          padding:8px 6px 14px; margin-bottom:8px; border-bottom:1px solid var(--line);
        }
        .side-brand .side-dino{
          width:42px; height:42px; object-fit:contain;
          filter:drop-shadow(0 3px 7px rgba(0,107,88,.28));
        }
        .side-brand span{ font-weight:800; color:var(--ink); font-size:16px; line-height:1.15; }
        .side-brand small{ font-weight:600; color:var(--muted); font-size:11px; }

        /* ---------- 标题层级 ---------- */
        h1,h2,h3,h4{ color:var(--ink); font-weight:700; letter-spacing:.1px; font-family:'Plus Jakarta Sans','Noto Sans SC',sans-serif; }

        /* ---------- 按钮（高级实心 + 柔和描边 + 微动效） ---------- */
        .stButton > button{
          border-radius:11px; font-weight:600; transition:transform .16s ease, box-shadow .16s ease, background .16s ease, border-color .16s ease;
          border:1px solid var(--line); box-shadow:var(--shadow-sm);
        }
        .stButton > button:hover{ transform:translateY(-1px); box-shadow:var(--shadow-md); }
        .stButton > button:active{ transform:translateY(0); }
        .stButton > button:focus-visible{ box-shadow:var(--ring) !important; }
        /* primary = 品牌橙渐变实心 */
        .stButton > button[kind="primary"],
        [data-testid="stBaseButton-primary"]{
          background:linear-gradient(135deg,var(--brand),var(--brand-2)) !important;
          border:1px solid rgba(0,107,88,.4) !important;
          color:#fff !important; box-shadow:0 8px 20px -6px rgba(0,107,88,.5);
        }
        .stButton > button[kind="primary"]:hover,
        [data-testid="stBaseButton-primary"]:hover{
          filter:brightness(1.03);
          box-shadow:0 12px 26px -6px rgba(0,107,88,.58) !important;
        }
        /* secondary = 描边，hover 染橙 */
        .stButton > button[kind="secondary"]:hover,
        [data-testid="stBaseButton-secondary"]:hover{
          color:var(--brand) !important; border-color:var(--brand) !important;
          background:var(--brand-tint) !important;
        }

        /* ---------- 输入控件（聚焦发光环，更通透） ---------- */
        textarea, input, [data-baseweb="select"] > div, [data-baseweb="input"] > div{
          border-radius:11px !important;
        }
        textarea:focus, input:focus,
        [data-baseweb="input"]:focus-within, [data-baseweb="select"]:focus-within{
          box-shadow:var(--ring) !important; border-color:var(--brand) !important;
        }

        /* ---------- Tabs（卡片 + 橙色高亮条） ---------- */
        .stTabs [data-baseweb="tab-list"]{
          gap:6px; border-bottom:2px solid var(--line); padding-bottom:0;
        }
        .stTabs [data-baseweb="tab"]{
          border-radius:11px 11px 0 0; padding:9px 20px; font-weight:600;
          color:var(--muted); background:transparent; transition:all .16s ease;
        }
        .stTabs [data-baseweb="tab"]:hover{ color:var(--brand); background:var(--brand-tint); }
        .stTabs [aria-selected="true"]{
          color:var(--brand) !important; background:var(--brand-tint) !important;
        }
        .stTabs [data-baseweb="tab-highlight"]{ background:linear-gradient(90deg,var(--brand),var(--brand-2)) !important; height:3px !important; border-radius:3px; }
        .stTabs [data-baseweb="tab-border"]{ display:none; }

        /* ---------- 卡片容器（st.container(border=True)）：玻璃白 + 悬浮抬升 ---------- */
        [data-testid="stVerticalBlockBorderWrapper"]{
          border-radius:var(--radius) !important; border:1px solid var(--line) !important;
          box-shadow:var(--shadow-sm); background:var(--card);
          transition:box-shadow .2s ease, transform .2s ease, border-color .2s ease;
        }
        [data-testid="stVerticalBlockBorderWrapper"]:hover{
          box-shadow:var(--shadow-md); border-color:#e6e8ee !important;
        }

        /* ---------- expander 卡片化 ---------- */
        [data-testid="stExpander"]{
          border-radius:var(--radius); border:1px solid var(--line);
          box-shadow:var(--shadow-sm); background:var(--card);
        }
        [data-testid="stExpander"] summary:hover{ color:var(--brand); }

        /* ---------- 指标 / 提示框圆角统一 ---------- */
        [data-testid="stMetric"]{
          background:var(--card); border:1px solid var(--line); border-radius:var(--radius);
          padding:14px 16px; box-shadow:var(--shadow-sm);
        }
        [data-testid="stAlert"]{ border-radius:var(--radius); }

        /* ---------- 向导顶部进度条 ---------- */
        .wiz-track{
          position:relative; height:7px; border-radius:7px; background:var(--line-soft);
          margin:2px 0 14px; overflow:hidden;
        }
        .wiz-fill{
          position:absolute; left:0; top:0; height:100%; border-radius:7px;
          background:linear-gradient(90deg,var(--brand),var(--brand-2));
          box-shadow:0 0 12px rgba(0,107,88,.45);
          transition:width .4s cubic-bezier(.4,0,.2,1);
        }

        /* ---------- 侧边栏交付物导航（左侧竖向 tabs） ---------- */
        [data-testid="stSidebar"] .nav-title{
          font-weight:700; font-size:13px; letter-spacing:.6px; text-transform:uppercase;
          color:var(--faint); margin:.2rem 0 .6rem;
        }
        [data-testid="stSidebar"] [role="radiogroup"]{ gap:6px; }
        [data-testid="stSidebar"] [role="radiogroup"] label{
          display:flex; align-items:center; width:100%;
          padding:11px 14px; border-radius:11px; border:1px solid transparent;
          font-weight:600; color:var(--muted); cursor:pointer; transition:all .16s ease;
        }
        [data-testid="stSidebar"] [role="radiogroup"] label:hover{
          background:var(--brand-tint); color:var(--brand); transform:translateX(2px);
        }
        [data-testid="stSidebar"] [role="radiogroup"] label:has(input:checked){
          background:linear-gradient(90deg,var(--brand-tint),rgba(232,250,245,.4));
          color:var(--brand); border-color:rgba(0,107,88,.28);
          box-shadow:inset 3px 0 0 var(--brand);
        }
        [data-testid="stSidebar"] [role="radiogroup"] label > div:first-child{ display:none; }

        /* ---------- 实时计时药丸 ---------- */
        .timer-pill{
          display:inline-flex; align-items:center; gap:9px;
          padding:7px 15px; border-radius:999px; font-size:13.5px; font-weight:600;
          color:var(--brand-dark); background:var(--brand-tint);
          border:1px solid rgba(0,107,88,.28); box-shadow:var(--shadow-sm);
          margin:4px 0;
        }
        .timer-pill.done{
          color:#0a7d4d; background:#eafaf2; border-color:rgba(16,185,129,.3);
        }
        .timer-pill b{ font-variant-numeric:tabular-nums; }
        .timer-spin{
          width:13px; height:13px; border-radius:50%;
          border:2px solid rgba(0,107,88,.3); border-top-color:var(--brand);
          animation:tmr-spin .7s linear infinite;
        }
        @keyframes tmr-spin{ to{ transform:rotate(360deg); } }

        /* ---------- 分隔线更轻 ---------- */
        hr{ margin:.8rem 0; border:none; border-top:1px solid var(--line); }

        /* ============================================================
           v2 精修层（2026-06-06）：字体层级 / 呼吸感 / 立体分层 /
           状态反馈 / 关键词标签 / 区块标题 / 瀑布流 / 微动效
           ============================================================ */

        /* —— 1) 字体层级（清晰的标题阶梯 + 行高呼吸）—— */
        .block-container h1{ font-size:30px; line-height:1.2; margin:.2rem 0 .6rem; }
        .block-container h2{ font-size:23px; line-height:1.25; margin:1.3rem 0 .55rem; }
        .block-container h3{ font-size:18.5px; line-height:1.3; margin:1.05rem 0 .45rem; }
        .block-container h4{ font-size:15.5px; line-height:1.35; margin:.85rem 0 .35rem; color:var(--muted); letter-spacing:.2px; }
        .block-container p, .block-container li{ line-height:1.62; }
        [data-testid="stCaptionContainer"], .stCaption, small{ color:var(--faint) !important; }
        [data-testid="stWidgetLabel"] label p{ font-weight:600 !important; color:var(--ink) !important; letter-spacing:.1px; }

        /* —— 2) 呼吸感（统一竖向节奏 + 更宽松留白）—— */
        .block-container{ padding-bottom:4rem; }
        [data-testid="stVerticalBlock"]{ gap:.85rem; }
        [data-testid="stHorizontalBlock"]{ gap:1rem; }
        [data-testid="stVerticalBlockBorderWrapper"] > div{ padding:2px; }

        /* —— 3) 立体分层（卡片更有层次的悬浮抬升 = 三维实用感）—— */
        [data-testid="stVerticalBlockBorderWrapper"]:hover{ transform:translateY(-2px); }
        [data-testid="stMetric"]{ transition:transform .18s ease, box-shadow .18s ease; }
        [data-testid="stMetric"]:hover{ transform:translateY(-2px); box-shadow:var(--shadow-md); }
        [data-testid="stMetricValue"]{ font-weight:800; color:var(--ink); }
        [data-testid="stMetricLabel"]{ color:var(--muted); font-weight:600; }

        /* —— 4) 状态反馈（提示框左侧强调条 + 语义色，反馈更明确）—— */
        [data-testid="stAlert"]{ border:1px solid var(--line); box-shadow:var(--shadow-sm); position:relative; overflow:hidden; }
        [data-testid="stAlert"]::before{ content:""; position:absolute; left:0; top:0; bottom:0; width:4px; }
        [data-testid="stAlert"]:has([data-testid="stAlertContentSuccess"])::before{ background:#16a34a; }
        [data-testid="stAlert"]:has([data-testid="stAlertContentInfo"])::before{ background:#2563eb; }
        [data-testid="stAlert"]:has([data-testid="stAlertContentWarning"])::before{ background:#f59e0b; }
        [data-testid="stAlert"]:has([data-testid="stAlertContentError"])::before{ background:#ef4444; }

        /* —— 5) 关键词标签 chip（“战格/标签”可复用组件）—— */
        .chip-row{ display:flex; flex-wrap:wrap; gap:8px; margin:6px 0 4px; }
        .chip{
          display:inline-flex; align-items:center; gap:6px;
          padding:5px 12px; border-radius:999px; font-size:12.5px; font-weight:600;
          color:var(--brand-dark); background:var(--brand-tint);
          border:1px solid rgba(0,107,88,.22); line-height:1;
          transition:all .14s ease;
        }
        .chip:hover{ box-shadow:var(--shadow-sm); transform:translateY(-1px); }
        .chip.gray{ color:var(--muted); background:var(--line-soft); border-color:var(--line); }
        .chip.ok{ color:#0a7d4d; background:#eafaf2; border-color:rgba(16,185,129,.28); }

        /* —— 6) 统一区块标题（图标 + 标题 + 渐隐细线，对齐有序）—— */
        .sec-head{
          display:flex; align-items:center; gap:10px; margin:.4rem 0 .7rem;
          font-family:'Poppins','Inter',sans-serif; font-weight:700; color:var(--ink); font-size:16px;
        }
        .sec-head .ic{
          display:inline-flex; align-items:center; justify-content:center;
          width:30px; height:30px; border-radius:9px;
          background:linear-gradient(135deg,var(--brand),var(--brand-2)); color:#fff;
          box-shadow:0 5px 12px -3px rgba(0,107,88,.5); font-size:15px;
        }
        .sec-head .line{ flex:1 1 auto; height:1px; background:linear-gradient(90deg,var(--line),transparent); }

        /* —— 7) 状态药丸 badge（运行中 / 完成 / 待办）—— */
        .badge{ display:inline-flex; align-items:center; gap:6px; padding:4px 11px; border-radius:999px;
          font-size:12px; font-weight:700; letter-spacing:.2px; border:1px solid transparent; }
        .badge.run{ color:#b45309; background:#fff7ed; border-color:#fed7aa; }
        .badge.ok{ color:#0a7d4d; background:#eafaf2; border-color:rgba(16,185,129,.3); }
        .badge.todo{ color:var(--muted); background:var(--line-soft); border-color:var(--line); }

        /* —— 8) 瀑布流图片容器（HTML 画廊用 .masonry 包裹 <img>）—— */
        .masonry{ column-count:3; column-gap:14px; }
        .masonry > *{ break-inside:avoid; margin:0 0 14px; border-radius:var(--radius);
          overflow:hidden; box-shadow:var(--shadow-sm); display:block; }
        @media (max-width:1100px){ .masonry{ column-count:2; } }
        @media (max-width:680px){ .masonry{ column-count:1; } }

        /* —— 9) 图片圆角 + 轻边框（统一观感，立体）—— */
        [data-testid="stImage"] img{ border-radius:12px; box-shadow:var(--shadow-sm); }

        /* —— 10) 滚动条 / 选区 / 入场微动效（视觉呼吸）—— */
        ::-webkit-scrollbar{ width:11px; height:11px; }
        ::-webkit-scrollbar-thumb{ background:#d7dbe3; border-radius:8px; border:3px solid var(--bg); }
        ::-webkit-scrollbar-thumb:hover{ background:#c2c8d2; }
        ::selection{ background:rgba(0,107,88,.22); }
        .hero, [data-testid="stVerticalBlockBorderWrapper"]{ animation:rise .42s cubic-bezier(.2,.7,.2,1) both; }
        @keyframes rise{ from{ opacity:0; transform:translateY(8px); } to{ opacity:1; transform:translateY(0); } }
        @media (prefers-reduced-motion: reduce){ *{ animation:none !important; } }

        /* ============================================================
           v3 沉浸层（2026-06-06）：瀑布流 / 玻璃拟态 / 粘性顶栏+侧栏 /
           分层视差 / 3D 空间 / Logo·标题环绕动效 / 微交互反馈 / 慢滚顺滑
           ============================================================ */

        /* —— 0) 屏宽与定位：拉宽内容、居中、慢滚顺滑、铺满不留大块空白 —— */
        html{ scroll-behavior:smooth; }
        .block-container{ max-width:1480px; padding-left:2.4rem; padding-right:2.4rem; margin:0 auto; }

        /* —— 1) 分层视差背景：背景固定（慢）、前景内容滚动（快）→ 纵深分层感 —— */
        .stApp{ background-attachment:fixed; }
        .stApp::before{
          content:""; position:fixed; inset:0; z-index:0; pointer-events:none;
          background:
            radial-gradient(1100px 560px at 84% -12%, rgba(0,107,88,.10), transparent 62%),
            radial-gradient(900px 520px at -10% 8%, rgba(120,120,255,.07), transparent 58%),
            radial-gradient(700px 700px at 50% 120%, rgba(38,194,163,.06), transparent 60%);
          animation:bgdrift 26s ease-in-out infinite alternate;
        }
        @keyframes bgdrift{
          from{ transform:translate3d(0,0,0) scale(1); }
          to{ transform:translate3d(0,-18px,0) scale(1.04); }
        }
        [data-testid="stAppViewContainer"] > .main{ position:relative; z-index:1; }

        /* —— 2) 粘性顶栏：hero 吸顶 + 玻璃半透明（滚动时悬浮在内容之上）—— */
        .hero{
          position:sticky; top:.5rem; z-index:50;
          background:linear-gradient(120deg, rgba(255,255,255,.72), rgba(232,250,245,.6));
          backdrop-filter:blur(18px) saturate(1.35); -webkit-backdrop-filter:blur(18px) saturate(1.35);
          border:1px solid rgba(255,255,255,.55);
        }

        /* —— 3) 粘性侧栏（窗口内）：品牌头吸顶 + 内容随窗滚动 + 玻璃质感 —— */
        [data-testid="stSidebar"]{
          background:linear-gradient(180deg, rgba(255,255,255,.78), rgba(251,252,254,.7));
          backdrop-filter:blur(16px) saturate(1.2); -webkit-backdrop-filter:blur(16px) saturate(1.2);
        }
        [data-testid="stSidebar"] [data-testid="stSidebarUserContent"]{
          position:sticky; top:0; max-height:100vh; overflow-y:auto;
        }
        .side-brand{
          position:sticky; top:0; z-index:5;
          background:linear-gradient(180deg, rgba(255,255,255,.92), rgba(255,255,255,.7));
          backdrop-filter:blur(10px); -webkit-backdrop-filter:blur(10px);
        }

        /* —— 4) Logo · 标题环绕动效：恐龙做轻盈环绕浮游，标题流光扫过 —— */
        .hero-dino{ animation:orbit 6.5s ease-in-out infinite; transform-origin:center; }
        @keyframes orbit{
          0%{ transform:translate(0,0) rotate(-4deg); }
          25%{ transform:translate(3px,-5px) rotate(2deg); }
          50%{ transform:translate(0,-8px) rotate(4deg); }
          75%{ transform:translate(-3px,-5px) rotate(2deg); }
          100%{ transform:translate(0,0) rotate(-4deg); }
        }
        .hero-icon{ position:relative; }
        .hero-icon::after{   /* 环绕光环 */
          content:""; position:absolute; inset:-8px; border-radius:50%;
          border:1.5px dashed rgba(0,107,88,.35); animation:ring-spin 14s linear infinite;
        }
        @keyframes ring-spin{ to{ transform:rotate(360deg); } }
        .hero-title{
          background:linear-gradient(92deg,var(--ink) 0%, var(--ink) 30%, var(--brand) 50%, var(--ink) 70%, var(--ink) 100%);
          background-size:220% 100%;
          -webkit-background-clip:text; background-clip:text; color:transparent;
          animation:sheen 7s ease-in-out infinite;
        }
        @keyframes sheen{ 0%,100%{ background-position:140% 0; } 50%{ background-position:-40% 0; } }

        /* —— 5) 3D 空间：卡片悬浮抬升（轻量，不裁剪内容、不产生内部滚动条）—— */
        [data-testid="stVerticalBlockBorderWrapper"]:hover{
          transform:translateY(-3px); box-shadow:var(--shadow-lg);
        }
        /* 图片容器绝不内部滚动：自然高度铺满，不出现小框 + 滚动轴 */
        [data-testid="stImage"]{ overflow:visible !important; }
        [data-testid="stImage"] img{ height:auto !important; max-height:none !important; }

        /* —— 6) 扁平边框：去厚重、统一 1px 干净描边（flat frame）—— */
        [data-testid="stExpander"], [data-testid="stMetric"], [data-testid="stAlert"]{
          box-shadow:none !important; border:1px solid var(--line) !important;
        }
        [data-testid="stExpander"]:hover, [data-testid="stMetric"]:hover{
          box-shadow:var(--shadow-sm) !important;
        }

        /* —— 7) 瀑布流图片墙 Masonry：错落铺满 + 悬浮抬升 + 标题浮层（点击放大）—— */
        .masonry .mz-card{
          position:relative; display:block; margin:0 0 14px; border-radius:var(--radius);
          overflow:hidden; box-shadow:var(--shadow-sm); break-inside:avoid; cursor:zoom-in;
          text-decoration:none;
          transition:transform .28s cubic-bezier(.2,.7,.2,1), box-shadow .28s ease;
        }
        .masonry .mz-card img{ display:block; width:100%; height:auto; border-radius:var(--radius); }
        .masonry .mz-card:hover{
          transform:translateY(-4px) scale(1.015);
          box-shadow:var(--shadow-lg); z-index:2;
        }
        .masonry .mz-cap{
          position:absolute; left:0; right:0; bottom:0;
          padding:18px 12px 9px; font-size:12.5px; font-weight:600; color:#fff;
          background:linear-gradient(transparent, rgba(16,24,40,.62));
          opacity:0; transform:translateY(6px); transition:all .26s ease;
        }
        .masonry .mz-card:hover .mz-cap{ opacity:1; transform:translateY(0); }

        /* —— 7b) 点击放大灯箱（纯 CSS :target，无需 JS / 按钮）—— */
        .zoomable{ display:block; cursor:zoom-in; text-decoration:none; }
        .zoomable img{
          display:block; width:100%; height:auto; border-radius:12px;
          box-shadow:var(--shadow-sm); transition:transform .25s ease, box-shadow .25s ease;
        }
        .zoomable:hover img{ transform:scale(1.012); box-shadow:var(--shadow-md); }
        .lightbox{ display:none; }
        .lightbox:target{
          display:flex; position:fixed; inset:0; z-index:9999; cursor:zoom-out;
          align-items:center; justify-content:center; padding:3vh 3vw;
          background:rgba(16,24,40,.84);
          backdrop-filter:blur(8px); -webkit-backdrop-filter:blur(8px);
          animation:lb-in .22s ease;
        }
        .lightbox:target img{
          max-width:96vw; max-height:92vh; width:auto; height:auto;
          border-radius:14px; box-shadow:0 32px 90px rgba(0,0,0,.55);
        }
        .lightbox .lb-cap{
          position:fixed; left:0; right:0; bottom:18px; text-align:center;
          color:#fff; font-size:14px; font-weight:600; letter-spacing:.2px;
          text-shadow:0 2px 8px rgba(0,0,0,.6);
        }
        @keyframes lb-in{ from{ opacity:0; } to{ opacity:1; } }

        /* —— 8) 微交互反馈：按钮按压涟漪、图片/控件即时反馈 —— */
        .stButton > button{ position:relative; overflow:hidden; }
        .stButton > button::after{
          content:""; position:absolute; left:50%; top:50%; width:0; height:0;
          border-radius:50%; background:rgba(255,255,255,.5);
          transform:translate(-50%,-50%); transition:width .45s ease, height .45s ease, opacity .6s ease;
          opacity:0;
        }
        .stButton > button:active::after{ width:240px; height:240px; opacity:1; transition:0s; }
        [data-testid="stImage"] img{ transition:transform .25s ease, box-shadow .25s ease; }
        [data-testid="stImage"] img:hover{ transform:scale(1.012); box-shadow:var(--shadow-md); }
        [data-baseweb="select"] > div:hover, [data-baseweb="input"] > div:hover{ border-color:var(--brand-2) !important; }
        [data-testid="stSlider"] [role="slider"]{ transition:transform .14s ease; }
        [data-testid="stSlider"] [role="slider"]:hover{ transform:scale(1.18); }

        @media (prefers-reduced-motion: reduce){
          .hero-dino, .hero-icon::after, .hero-title, .stApp::before{ animation:none !important; }
        }
        </style>""",
        unsafe_allow_html=True,
    )


if __name__ == "__main__":
    main()
