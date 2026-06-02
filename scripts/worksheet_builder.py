"""Worksheet PPTX 生成器 v2.0（6 页固定模板，对齐真实 L5-1 样本）。

页面结构（强制 6 页）：
    Page 1  Vocabulary  - Match the words to their definitions   (5 对连线)
    Page 2  Vocabulary  - Use the words / phrase to fill blanks  (5 题填空 + 词库条)
    Page 3  Sentence    - Choose the correct sentence            (4 题二选一 + 绘本图)
    Page 4  Reading     - Choose the correct answer              (全文 + 8 道 3 选)
    Page 5  Writing     - Write about [theme]                    (5 步骨架 + 写作区)
    Page 6  Reading     - Filling the mind map                   (3 列表 5 行)

字体/字号（v1.6 真实样本）：
    大标题 Poppins Bold 20pt #333333  / 副标题 Poppins Regular 12pt #666666
    题号  Poppins Bold 16pt 圆形粉底白字
    题干  Poppins Regular 16pt 黑色  /  Reading 长文 12pt 黑色

品牌外框（统一 6 页）：
    粉色外背景 (BRAND_COLORS[level]) + 内白圆角
    左上 VIPKID Dino Reading Club logo（Dino 头像 + 白色文字）
    右上 Name 五角形角标 (粉底白字)
    右下 footer "Level X - <Title>" 白字
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as _RawInches, Pt, Emu

from config import BRAND_DIR, brand_color_rgb
from parser import BookOutline
from text_format import (
    _to_us_spelling,
    format_word_answer,
    format_sentence_answer,
    smart_format_answer,
    is_sentence_like,
)


# ---------- 几何尺寸 ----------
# 历史版式在 10.83 x 7.50 in（275x190mm）设计稿空间里布局；
# 输出时整体等比放大到「真·A4 横向」297x210mm = 11.69 x 8.27 in。
# 缩放按高度填满（系数≈1.1027）；粉色满版背景略向右溢出被裁切无碍，
# 真实内容最大右边界 10.53in×1.1027≈11.61in < 11.69in，不会被裁。
DESIGN_W = 10.83
DESIGN_H = 7.50
A4_LAND_W = 11.69   # 297 mm
A4_LAND_H = 8.27    # 210 mm
_WS_SCALE = A4_LAND_H / DESIGN_H

# 版式代码继续以设计稿坐标书写，下面的 Inches() 会自动放大。
SLIDE_W = DESIGN_W
SLIDE_H = DESIGN_H


def Inches(value):
    """把设计稿英寸坐标整体等比放大到 A4 横向后再转 EMU。"""
    return _RawInches(value * _WS_SCALE)

# 内容白底（粉色外框内的圆角白区）
CONTENT_X = 0.30
CONTENT_Y = 0.70
CONTENT_W = 10.23
CONTENT_H = 6.55
CONTENT_ROUND = 0.06  # 圆角调整比例（python-pptx adjustments[0]）

# 顶部 logo 区（露在粉色背景上）
LOGO_X = 0.40
LOGO_Y = 0.10
LOGO_ICON_W = 0.55
LOGO_ICON_H = 0.55

# Name 角标（v1.9：放成矩形标签，避免旋转 PENTAGON 跑位）
NAME_X = 8.85
NAME_Y = 0.18
NAME_W = 1.50
NAME_H = 0.42

# Footer（v1.9：往上挪到内容白底底边附近，避免被裁切）
FOOTER_X = 0.30
FOOTER_Y = 6.85
FOOTER_W = 10.23
FOOTER_H = 0.30


# ---------- 字号/颜色 ----------
FONT = "Poppins"
FONT_BOLD = "Poppins"
# underscore 字符在 Poppins 下被压扁，改用 Arial 才显示得清晰粗实
FONT_BLANK = "Arial"

# 按用户偏好：大标题 40pt + 副标题 22pt（更大气）
TITLE_PT = 40
SUBTITLE_PT = 22
BODY_PT = 18
READING_PT = 13
QNUM_PT = 18
HEADER_PT = 20  # mind map 表头
LOGO_TEXT_PT = 20
NAME_PT = 16
FOOTER_PT = 14.5

TITLE_RGB = RGBColor(0x33, 0x33, 0x33)  # 深炭灰
SUB_RGB = RGBColor(0x66, 0x66, 0x66)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAME_FILL = RGBColor(0xF8, 0xC8, 0xDC)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

# Mind Map 三列表头颜色（v1.6 紫粉 / 黄 / 绿）
MM_PURPLE = RGBColor(0xE6, 0xD8, 0xF2)
MM_YELLOW = RGBColor(0xFA, 0xEB, 0xC6)
MM_GREEN = RGBColor(0xCE, 0xE7, 0xCD)
MM_PURPLE_DARK = RGBColor(0xC9, 0xB3, 0xE2)
MM_YELLOW_DARK = RGBColor(0xF5, 0xDC, 0x95)
MM_GREEN_DARK = RGBColor(0xA8, 0xD3, 0xA6)

# Reading 红框
READ_BORDER = RGBColor(0xE9, 0x52, 0x83)


# ============================================================
#  对外入口
# ============================================================

def build_worksheet(
    outline: BookOutline,
    out_path: Path,
    *,
    image_paths: Optional[Iterable[Path]] = None,
) -> Path:
    """生成 6 页 worksheet。image_paths 为绘本图（page_00.png .. page_07.png）。

    Sentence 页（Page 3）会复用 image_paths[2..5]（即故事 P2..P5）作为 4 题配图。
    若不传 image_paths，Page 3 用占位灰块。
    """
    data = _resolve_worksheet_data(outline)
    brand_rgb = brand_color_rgb(outline.level)
    level_label = _level_label(outline.level)
    footer_text = f"{level_label} - {outline.title}"
    # v2.1: 用专门的 Dino 头 icon（不是设定卡 dino_logo.png）
    logo_icon = BRAND_DIR / "dino_head_icon.png"
    if not logo_icon.exists():
        logo_icon = BRAND_DIR / "dino_logo.png"  # 兜底
    images = list(image_paths or [])

    prs = Presentation()
    prs.slide_width = _RawInches(A4_LAND_W)
    prs.slide_height = _RawInches(A4_LAND_H)
    blank = prs.slide_layouts[6]

    # Page 1: Vocabulary - Match
    s = prs.slides.add_slide(blank)
    _draw_brand_frame(s, brand_rgb, footer_text, logo_icon)
    _build_p1_match(s, brand_rgb, data["match_pairs"], images)

    # Page 2: Vocabulary - Fill blanks
    s = prs.slides.add_slide(blank)
    _draw_brand_frame(s, brand_rgb, footer_text, logo_icon)
    _build_p2_fill(s, brand_rgb, data["fill_blanks"], data["word_bank"], images)

    # Page 3: Sentence - MC
    s = prs.slides.add_slide(blank)
    _draw_brand_frame(s, brand_rgb, footer_text, logo_icon)
    _build_p3_sentence(s, brand_rgb, data["sentence_mcs"], images)

    # Page 4: Reading - MC
    s = prs.slides.add_slide(blank)
    _draw_brand_frame(s, brand_rgb, footer_text, logo_icon)
    _build_p4_reading(s, brand_rgb, data["reading_text"], data["reading_mcs"])

    # Page 5: Writing
    s = prs.slides.add_slide(blank)
    _draw_brand_frame(s, brand_rgb, footer_text, logo_icon)
    _build_p5_writing(s, brand_rgb, data["writing"])

    # Page 6: Mind Map
    s = prs.slides.add_slide(blank)
    _draw_brand_frame(s, brand_rgb, footer_text, logo_icon)
    _build_p6_mindmap(s, data["mind_map"])

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def attach_worksheet_questions(
    outline: BookOutline, data, *, reading_q_count: int = 4
) -> None:
    """挂载 AI 抽取的 worksheet 内容到 outline。

    v1.9：兼容两种输入格式：
      - dict：直接用（match_pairs/word_bank/fill_blanks/sentence_mcs/reading_text/reading_mcs/writing/mind_map）
      - list[dict]（AI 抽出来的 6 道题 list，每条 {type, items, ...}）→ 自动跑 adapter

    v2.0: 加 reading_q_count 参数（4/6/8），控制 Reading MC 页题量。
    """
    if isinstance(data, list):
        data = _questions_list_to_template_data(data, outline)
    if isinstance(data, dict):
        data = dict(data)
        data["_reading_q_count"] = reading_q_count
    setattr(outline, "_worksheet_data", data)


def _questions_list_to_template_data(qlist: list[dict], outline: BookOutline) -> dict:
    """把 AI 抽取的"按 level 池子的 6 道题"映射到 worksheet PPTX 模板的 7 个固定字段。

    Worksheet 模板有 6 页固定结构：
      P1 Match (word ↔ definition)
      P2 Fill blanks (word_bank + 5 句 fill_blanks)
      P3 Sentence MC (4 题 2 选 1，可配图)
      P4 Reading MC (4 题，每题 3 options)
      P5 Writing (scaffold + 写作区)
      P6 Mind Map (character / problem / solution)

    AI 抽取的题型五花八门（match_definition / fill_blank / true_false / inference / unscramble / …）
    本函数把它们分发到对应模板字段，并用 outline 数据兜底空字段。
    """
    out: dict = {
        "match_pairs": [],
        "word_bank": [],
        "fill_blanks": [],
        "sentence_mcs": [],
        "reading_text": "",
        "reading_mcs": [],
        "writing": {},
        "mind_map": [],
    }

    for q in qlist or []:
        if not isinstance(q, dict):
            continue
        qtype = (q.get("type") or "").lower()
        # AI 抽取偶尔把 items 给成字符串列表（而非 dict 列表），统一规整成 dict，
        # 字符串元素塞进 _str，供 match/fill 等分支兜底使用，避免 .get 崩溃。
        items = [
            it if isinstance(it, dict) else {"_str": str(it)}
            for it in (q.get("items") or [])
            if it is not None
        ]
        extra = q.get("extra") or ""

        if qtype == "match_definition":
            out["match_pairs"] = [
                {"word": it.get("word") or it.get("_str", ""),
                 "def": it.get("def") or it.get("definition", "")}
                for it in items if (it.get("word") or it.get("_str"))
            ][:5]
            if not out["word_bank"]:
                out["word_bank"] = [
                    it.get("word") or it.get("_str", "") for it in items
                ][:5]

        elif qtype in ("fill_blank", "fill_blank_simple", "fill_blank_advanced", "emotion_fill"):
            out["fill_blanks"] = [
                {"sentence": it.get("sentence") or it.get("_str", ""),
                 "answer": it.get("answer", "")}
                for it in items if (it.get("sentence") or it.get("_str"))
            ][:5]
            if not out["word_bank"]:
                if extra:
                    bank = [w.strip() for w in extra.split(",") if w.strip()]
                else:
                    bank = [it.get("answer", "") for it in items]
                out["word_bank"] = bank[:5]

        elif qtype in ("true_false", "true_false_simple"):
            for it in items[:4]:
                stmt = it.get("statement", "")
                if not stmt:
                    continue
                # AI 给的 T/F 转成 sentence MC：正确句 vs 反义句
                ans = (it.get("answer") or "T").upper()
                opt_true = stmt
                opt_false = "Not " + stmt[0].lower() + stmt[1:] if stmt else ""
                out["sentence_mcs"].append({
                    "options": [opt_true, opt_false],
                    "correct": 0 if ans == "T" else 1,
                })

        elif qtype in ("inference", "reading_mc"):
            for it in items[:4]:
                stem = it.get("q") or it.get("question", "")
                if not stem:
                    continue
                opts = it.get("options") or []
                if len(opts) < 2:
                    continue
                out["reading_mcs"].append({
                    "q": stem,
                    "options": list(opts)[:3],
                    "correct": int(it.get("correct", 0)),
                })

        elif qtype in ("plot_chart", "plot_chart_pbl"):
            # AI 给 {label: Setting/Problem/Solution/..., answer: 内容}
            buf: dict[str, str] = {}
            for it in items:
                buf[(it.get("label") or "").lower()] = it.get("answer", "")
            character = buf.get("characters") or buf.get("character") or "Main character"
            problem = buf.get("problem") or buf.get("conflict") or ""
            solution = buf.get("solution") or buf.get("resolution") or buf.get("ending") or ""
            out["mind_map"].append({
                "character": character,
                "problem": problem,
                "solution": solution,
            })

        elif qtype == "compare_contrast":
            for it in items[:3]:
                out["mind_map"].append({
                    "character": it.get("topic", "Comparison"),
                    "problem": it.get("side_a", ""),
                    "solution": it.get("side_b", ""),
                })

        elif qtype in ("essay_short", "personal_write", "personal_simple",
                       "draw_favorite", "open_ended_pbl", "research_pbl"):
            out["writing"] = {
                "theme": outline.theme or "the story",
                "title": (q.get("title") or extra or f"Write about {outline.title}").strip(),
                "steps": ["", "", "", "", ""],
                "step_labels": [
                    "Beginning:", "First event:", "Second event:",
                    "Funny event:", "Ending:",
                ],
                "min_words": 50, "max_words": 80,
            }

        elif qtype == "unscramble":
            # 转化为 fill_blanks（"Unscramble: o c k l → ____" 的形式）
            if not out["fill_blanks"]:
                for it in items[:5]:
                    scr = it.get("scrambled", "")
                    ans = it.get("answer", "")
                    if scr and ans:
                        out["fill_blanks"].append({
                            "sentence": f"Unscramble: {scr} → ____",
                            "answer": ans,
                        })

        elif qtype in ("word_order", "word_order_simple", "story_sequence"):
            # 当 sentence_mcs（按 order 排）— 不太贴合，跳过
            pass

        elif qtype in ("rewrite_tense", "rewrite_voice"):
            if not out["fill_blanks"]:
                for it in items[:5]:
                    prm = it.get("prompt", "")
                    ans = it.get("answer", "")
                    if prm and ans:
                        out["fill_blanks"].append({
                            "sentence": f"Rewrite: {prm} → ____",
                            "answer": ans,
                        })

        elif qtype in ("color_match", "circle_match", "word_to_pic"):
            # 简单 vocab cue，用作 word_bank 兜底
            if not out["word_bank"]:
                out["word_bank"] = [
                    it.get("word", "") for it in items if it.get("word")
                ][:5]

    # ----- 兜底字段：用 outline 数据补全 -----
    pages = outline.pages or []
    story_text = " ".join(
        (p.text or "").strip() for p in pages if (p.text or "").strip()
    ).strip()

    if not out["reading_text"]:
        out["reading_text"] = story_text or "Story text goes here."

    # Match 兜底：用 vocab
    if not out["match_pairs"]:
        words = (outline.vocabulary_mastery or outline.vocabulary_simple or [])[:5]
        out["match_pairs"] = [
            {"word": w, "def": f"meaning of {w}"} for w in words
        ]
        if not out["word_bank"]:
            out["word_bank"] = words[:]

    # word_bank 兜底
    if not out["word_bank"]:
        out["word_bank"] = (outline.vocabulary_mastery or outline.vocabulary_simple or [])[:5]

    # Fill 兜底
    if not out["fill_blanks"] and out["word_bank"]:
        out["fill_blanks"] = [
            {"sentence": f"I feel ____ when I see this.", "answer": w}
            for w in out["word_bank"][:5]
        ]

    story_sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", story_text) if s.strip()]

    # Sentence MC 兜底：用故事原文 + 编一句反义
    if not out["sentence_mcs"]:
        for s in story_sentences[:4]:
            out["sentence_mcs"].append({
                "options": [s, "Not " + s[0].lower() + s[1:] if s else ""],
                "correct": 0,
            })

    # Reading MC 兜底：从故事生成简单 5W（按上限取，不固定 4）
    if not out["reading_mcs"] and story_sentences:
        for i, s in enumerate(story_sentences[:8]):
            out["reading_mcs"].append({
                "q": f"What does the passage say in sentence {i + 1}?",
                "options": [
                    s[:60] + ("..." if len(s) > 60 else ""),
                    "Something opposite happens.",
                    "The story does not mention it.",
                ],
                "correct": 0,
            })

    # Writing 兜底
    if not out["writing"]:
        out["writing"] = {
            "theme": outline.theme or "the story",
            "title": f"My Story About {outline.title}",
            "steps": ["", "", "", "", ""],
            "step_labels": [
                "Beginning:", "First event:", "Second event:",
                "Funny event:", "Ending:",
            ],
            "min_words": 50, "max_words": 80,
        }

    # Mind Map 兜底
    if not out["mind_map"]:
        out["mind_map"] = [{
            "character": "Main character",
            "problem": "What is the problem in the story?",
            "solution": "How is it solved?",
        }]

    return out


# ============================================================
#  品牌外框
# ============================================================

def _draw_brand_frame(slide, brand_rgb: tuple, footer_text: str, logo_icon: Path) -> None:
    """所有 6 页统一的外框：粉背景 + 圆角白底 + 左上 logo + 右上 Name + 右下 footer。"""
    # 1. 外背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*brand_rgb)
    bg.line.fill.background()
    bg.shadow.inherit = False

    # 2. 内白底圆角
    content = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X), Inches(CONTENT_Y),
        Inches(CONTENT_W), Inches(CONTENT_H),
    )
    content.adjustments[0] = CONTENT_ROUND
    content.fill.solid()
    content.fill.fore_color.rgb = WHITE
    content.line.fill.background()
    content.shadow.inherit = False

    # 3. 左上 logo (icon + 文字)
    if logo_icon and logo_icon.exists():
        try:
            slide.shapes.add_picture(
                str(logo_icon),
                Inches(LOGO_X), Inches(LOGO_Y),
                width=Inches(LOGO_ICON_W), height=Inches(LOGO_ICON_H),
            )
        except Exception:
            pass
    text_x = LOGO_X + LOGO_ICON_W + 0.10
    tb = slide.shapes.add_textbox(
        Inches(text_x), Inches(LOGO_Y), Inches(4.5), Inches(LOGO_ICON_H),
    )
    tb.text_frame.margin_left = tb.text_frame.margin_right = 0
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p.add_run()
    r.text = "VIPKID Dino Reading Club"
    r.font.name = FONT
    r.font.size = Pt(LOGO_TEXT_PT)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # 4. 右上 Name 角标（v2.0 改回盾形：上方矩形 + 下方三角尖角，对齐官方模板）
    # Name 颜色用 brand_rgb 加深 30%（更接近模板的酒红/暗粉效果）
    name_dark = RGBColor(
        max(0, brand_rgb[0] - 50),
        max(0, brand_rgb[1] - 30),
        max(0, brand_rgb[2] - 30),
    )

    # 矩形部分（上半部）
    name_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(NAME_X), Inches(NAME_Y),
        Inches(NAME_W), Inches(NAME_H),
    )
    name_rect.fill.solid()
    name_rect.fill.fore_color.rgb = name_dark
    name_rect.line.fill.background()
    name_rect.shadow.inherit = False
    tf = name_rect.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Name"
    r.font.name = FONT
    r.font.size = Pt(NAME_PT)
    r.font.color.rgb = WHITE
    r.font.bold = True

    # 三角形尖角部分（朝下，宽度比矩形窄，居中下方）
    tri_w = NAME_W * 0.55
    tri_h = 0.22
    tri_x = NAME_X + (NAME_W - tri_w) / 2
    tri_y = NAME_Y + NAME_H
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(tri_x), Inches(tri_y),
        Inches(tri_w), Inches(tri_h),
    )
    tri.rotation = 180  # 翻转让尖朝下
    tri.fill.solid()
    tri.fill.fore_color.rgb = name_dark
    tri.line.fill.background()
    tri.shadow.inherit = False

    # 5. 底部 footer（v1.9：全宽 + 右对齐，往上挪到白底里面避免被裁切）
    fo = slide.shapes.add_textbox(
        Inches(FOOTER_X), Inches(FOOTER_Y),
        Inches(FOOTER_W), Inches(FOOTER_H),
    )
    fo.text_frame.margin_left = fo.text_frame.margin_right = Inches(0.10)
    p = fo.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = footer_text
    r.font.name = FONT
    r.font.size = Pt(FOOTER_PT)
    r.font.color.rgb = RGBColor(*brand_rgb)
    r.font.bold = True


def _get_subtitle(qtype_id: str, fallback: str) -> str:
    """v2.0 从题型库取标准英文 instruction；找不到就用兜底文本。"""
    try:
        from worksheet_question_types import get_type
        t = get_type(qtype_id)
        if t and t.en_instr:
            return t.en_instr
    except Exception:
        pass
    return fallback


def _add_title(slide, title: str, subtitle: str) -> None:
    """大标题 + 副标题（居中，位于内容白底顶部）。"""
    # Title
    tb = slide.shapes.add_textbox(
        Inches(CONTENT_X), Inches(CONTENT_Y + 0.20),
        Inches(CONTENT_W), Inches(0.55),
    )
    tb.text_frame.margin_left = tb.text_frame.margin_right = 0
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(TITLE_PT)
    r.font.color.rgb = TITLE_RGB

    # Subtitle
    sb = slide.shapes.add_textbox(
        Inches(CONTENT_X), Inches(CONTENT_Y + 0.78),
        Inches(CONTENT_W), Inches(0.30),
    )
    sb.text_frame.margin_left = sb.text_frame.margin_right = 0
    p2 = sb.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = FONT
    r2.font.size = Pt(SUBTITLE_PT)
    r2.font.color.rgb = SUB_RGB


# ============================================================
#  Page 1 — Vocabulary 连线
# ============================================================

def _build_p1_match(slide, brand_rgb: tuple, pairs: list[dict], images: list[Path]) -> None:
    """5 对连线：左列小图（绘本图） + 中列粉色实心词卡 ↔ 右列白底粉边定义卡。

    v1.8 新增：每个 vocab 配一张绘本小图（用 page_02..page_06.png 循环）当 visual cue。
    """
    _add_title(slide, "Vocabulary", _get_subtitle("vocab_match_definition", "Match the words to their definitions."))

    n = min(len(pairs), 5)
    if n == 0:
        return

    # 区域
    area_top = CONTENT_Y + 1.30
    area_bottom = CONTENT_Y + CONTENT_H - 0.30
    area_h = area_bottom - area_top

    # 三列布局：小图 0.90 + 词卡 2.10 + 定义卡 4.20
    img_x = CONTENT_X + 0.40
    img_w = 0.90
    word_x = img_x + img_w + 0.20  # 1.50
    word_w = 2.10
    def_x = word_x + word_w + 0.30  # 3.90 起
    def_w = CONTENT_W - (def_x - CONTENT_X) - 0.40
    row_gap = 0.18
    row_h = (area_h - row_gap * (n - 1)) / n

    for i, pair in enumerate(pairs[:n]):
        y = area_top + i * (row_h + row_gap)

        # 小图（左，绘本插画，正方形等比缩放）
        img_idx = (i % max(1, len(images) - 2)) + 2  # 从 page_02 起
        img_path = images[img_idx] if (img_idx < len(images) and images[img_idx]) else None
        if img_path and img_path.exists():
            try:
                from PIL import Image as _PILImg
                with _PILImg.open(str(img_path)) as _pim:
                    iw, ih = _pim.size
                aspect = iw / ih if ih else 1.0
                fit_w = img_w
                fit_h = fit_w / aspect
                if fit_h > row_h:
                    fit_h = row_h
                    fit_w = fit_h * aspect
                off_x = img_x + (img_w - fit_w) / 2
                off_y = y + (row_h - fit_h) / 2
                slide.shapes.add_picture(
                    str(img_path),
                    Inches(off_x), Inches(off_y),
                    width=Inches(fit_w), height=Inches(fit_h),
                )
            except Exception:
                _draw_image_placeholder(slide, img_x, y, img_w, row_h)
        else:
            _draw_image_placeholder(slide, img_x, y, img_w, row_h)

        # 词卡（中，粉色实心，白字）
        wc = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(word_x), Inches(y), Inches(word_w), Inches(row_h),
        )
        wc.adjustments[0] = 0.28
        wc.fill.solid()
        wc.fill.fore_color.rgb = RGBColor(*brand_rgb)
        wc.line.fill.background()
        wc.shadow.inherit = False
        tf = wc.text_frame
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(pair.get("word", "")).strip()
        r.font.name = FONT
        r.font.size = Pt(BODY_PT)
        r.font.color.rgb = WHITE
        r.font.bold = False

        # 定义卡（右，白底粉边，黑字）
        dc = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(def_x), Inches(y), Inches(def_w), Inches(row_h),
        )
        dc.adjustments[0] = 0.28
        dc.fill.solid()
        dc.fill.fore_color.rgb = WHITE
        dc.line.color.rgb = RGBColor(*brand_rgb)
        dc.line.width = Pt(1.8)
        dc.shadow.inherit = False
        tf = dc.text_frame
        tf.margin_left = tf.margin_right = Inches(0.15)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT  # v1.8: 定义文本左对齐更易读
        r = p.add_run()
        r.text = str(pair.get("def", "")).strip()
        r.font.name = FONT
        r.font.size = Pt(BODY_PT)
        r.font.color.rgb = BLACK


# ============================================================
#  Page 2 — Vocabulary 填空
# ============================================================

def _build_p2_fill(slide, brand_rgb: tuple, fills: list[dict], word_bank: list[str],
                    images: Optional[list[Path]] = None) -> None:  # noqa: ARG001
    """顶部粉色词库条 + 5 道填空题（用 ____ 表示空）。images 参数留作未来扩展。"""
    _add_title(slide, "Vocabulary", _get_subtitle("vocab_fill_blank", "Use the words to fill each blank."))

    # 词库条（粉色实心圆角，水平排列词）
    bank_top = CONTENT_Y + 1.80
    bank_h = 0.55
    bank_x = CONTENT_X + 1.50
    bank_w = CONTENT_W - 3.00
    if word_bank:
        bk = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bank_x), Inches(bank_top), Inches(bank_w), Inches(bank_h),
        )
        bk.adjustments[0] = 0.4
        bk.fill.solid()
        bk.fill.fore_color.rgb = RGBColor(*brand_rgb)
        bk.line.fill.background()
        bk.shadow.inherit = False
        tf = bk.text_frame
        tf.margin_left = tf.margin_right = Inches(0.15)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        # 水平排列：词之间用空格
        joined = "    ".join(word_bank)
        r = p.add_run()
        r.text = joined
        r.font.name = FONT
        r.font.size = Pt(BODY_PT)
        r.font.color.rgb = WHITE
        r.font.bold = False

    # 5 道填空（垂直排列）
    n = min(len(fills), 5)
    if n == 0:
        return
    qa_top = bank_top + bank_h + 0.45
    qa_bottom = CONTENT_Y + CONTENT_H - 0.30
    qa_h = qa_bottom - qa_top
    row_h = qa_h / n

    for i, qa in enumerate(fills[:n]):
        y = qa_top + i * row_h
        tb = slide.shapes.add_textbox(
            Inches(CONTENT_X + 0.80), Inches(y),
            Inches(CONTENT_W - 1.60), Inches(row_h),
        )
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        sentence = _ensure_blank(str(qa.get("sentence", "")))
        _emit_with_underscore_lock(p, f"{i + 1}.  {sentence}", BODY_PT, BLACK)


def _ensure_blank(text: str) -> str:
    """若 sentence 里没有 ____ 占位，自动加 underscores。"""
    if "____" in text or "_____" in text or "___" in text:
        return text
    # 把第一个 [blank] / [BLANK] / [BLK] 替换为 ____
    for token in ("[blank]", "[BLANK]", "[BLK]"):
        if token in text:
            return text.replace(token, "________", 1)
    return text


def _emit_with_underscore_lock(paragraph, text: str, size_pt: int, color: RGBColor) -> None:
    """把文本里连续的下划线段（____）切出来用 Arial 字体渲染（Poppins 下 _ 太扁），
    其余部分仍用 Poppins。"""
    import re as _re
    parts = _re.split(r"(_{3,})", text)  # 至少 3 个 _ 才视作占位
    if len(parts) == 1:
        # 无下划线段，直接整段 Poppins
        r = paragraph.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size_pt)
        r.font.color.rgb = color
        return
    for chunk in parts:
        if not chunk:
            continue
        r = paragraph.add_run()
        r.text = chunk
        if chunk.startswith("___"):
            r.font.name = FONT_BLANK  # Arial，下划线粗实
        else:
            r.font.name = FONT
        r.font.size = Pt(size_pt)
        r.font.color.rgb = color


# ============================================================
#  Page 3 — Sentence MC (二选一 + 绘本配图)
# ============================================================

def _build_p3_sentence(slide, brand_rgb: tuple, mcs: list[dict], images: list[Path]) -> None:
    """4 题二选一，每题左侧绘本图 + 右侧 A/B 选项 + 行首圆圈题号。"""
    _add_title(slide, "Sentence", _get_subtitle("sent_tick_sentence", "Look at the picture and tick the correct sentence."))

    n = min(len(mcs), 4)
    if n == 0:
        return

    area_top = CONTENT_Y + 1.40
    area_bottom = CONTENT_Y + CONTENT_H - 0.30
    area_h = area_bottom - area_top
    row_gap = 0.10
    row_h = (area_h - row_gap * (n - 1)) / n

    qnum_size = 0.45
    qnum_x = CONTENT_X + 0.30
    img_x = qnum_x + qnum_size + 0.20
    img_w = 1.80
    opt_x = img_x + img_w + 0.30
    opt_w = CONTENT_W - (opt_x - CONTENT_X) - 0.40

    for i, mc in enumerate(mcs[:n]):
        y = area_top + i * (row_h + row_gap)
        cy = y + (row_h - qnum_size) / 2

        # 圆形题号
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(qnum_x), Inches(cy),
            Inches(qnum_size), Inches(qnum_size),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*brand_rgb)
        circle.line.fill.background()
        circle.shadow.inherit = False
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(QNUM_PT)
        r.font.color.rgb = WHITE

        # 图（复用绘本 page_02..page_05），保持长宽比，居中放入 img_w x (row_h-0.10) 框
        img_idx = i + 2  # P3 题 1 用绘本 page_02
        img_path = images[img_idx] if (img_idx < len(images) and images[img_idx]) else None
        max_box_w = img_w
        max_box_h = row_h - 0.10
        if img_path and img_path.exists():
            try:
                from PIL import Image as _PILImg
                with _PILImg.open(str(img_path)) as _pim:
                    iw, ih = _pim.size
                aspect = iw / ih if ih else 1.0
                # 尝试按宽适配
                fit_w = max_box_w
                fit_h = fit_w / aspect
                if fit_h > max_box_h:
                    fit_h = max_box_h
                    fit_w = fit_h * aspect
                off_x = img_x + (max_box_w - fit_w) / 2
                off_y = y + 0.05 + (max_box_h - fit_h) / 2
                slide.shapes.add_picture(
                    str(img_path),
                    Inches(off_x), Inches(off_y),
                    width=Inches(fit_w), height=Inches(fit_h),
                )
            except Exception:
                _draw_image_placeholder(slide, img_x, y + 0.05, max_box_w, max_box_h)
        else:
            _draw_image_placeholder(slide, img_x, y + 0.05, max_box_w, max_box_h)

        # 选项 A/B（垂直堆叠 + 复选框）
        options = mc.get("options") or []
        opt_h = (row_h - 0.05) / max(len(options), 1)
        for j, opt in enumerate(options[:2]):
            oy = y + j * opt_h
            # 复选框
            cb_size = 0.22
            cb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(opt_x), Inches(oy + (opt_h - cb_size) / 2),
                Inches(cb_size), Inches(cb_size),
            )
            cb.fill.solid()
            cb.fill.fore_color.rgb = WHITE
            cb.line.color.rgb = BLACK
            cb.line.width = Pt(1.0)
            cb.shadow.inherit = False
            # 选项文字
            tb = slide.shapes.add_textbox(
                Inches(opt_x + cb_size + 0.15), Inches(oy),
                Inches(opt_w - cb_size - 0.20), Inches(opt_h),
            )
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            letter = chr(ord("A") + j)
            r.text = f"{letter}. {opt}"
            r.font.name = FONT
            r.font.size = Pt(BODY_PT)
            r.font.color.rgb = BLACK


def _draw_image_placeholder(slide, x: float, y: float, w: float, h: float) -> None:
    ph = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    ph.adjustments[0] = 0.05
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    ph.line.color.rgb = LIGHT_GRAY
    ph.line.width = Pt(1.0)
    ph.shadow.inherit = False


# ============================================================
#  Page 4 — Reading 全文 + 8 道 3 选
# ============================================================

def _build_p4_reading(slide, brand_rgb: tuple, text: str, mcs: list[dict]) -> None:
    _add_title(slide, "Reading", _get_subtitle("read_mc_questions", "Choose the correct answer for each question."))

    text = (text or "").strip()

    # 顶部红框 = Reading 全文。短文字数差异很大（长文会撑破固定框），
    # 按字符数自适应字号 + 行距，确保始终落在框内不压字。
    tlen = len(text)
    if tlen > 780:
        read_pt, read_ls = 10.5, 1.12
    elif tlen > 560:
        read_pt, read_ls = 11.5, 1.15
    elif tlen > 380:
        read_pt, read_ls = 12.5, 1.20
    else:
        read_pt, read_ls = float(READING_PT), 1.25

    text_top = CONTENT_Y + 1.35
    text_h = 2.15
    text_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 0.40), Inches(text_top),
        Inches(CONTENT_W - 0.80), Inches(text_h),
    )
    text_box.adjustments[0] = 0.03
    text_box.fill.solid()
    text_box.fill.fore_color.rgb = WHITE
    text_box.line.color.rgb = READ_BORDER
    text_box.line.width = Pt(1.5)
    text_box.shadow.inherit = False
    # python-pptx 不直接支持 dashed line per shape — 这里用实线代替

    tf = text_box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(read_pt)
    r.font.color.rgb = BLACK
    p.line_spacing = read_ls

    # 下方 4 道 MC（2x2 排版）
    n = min(len(mcs), 4)
    if n == 0:
        return
    mc_top = text_top + text_h + 0.20
    mc_bottom = CONTENT_Y + CONTENT_H - 0.25
    mc_h_total = mc_bottom - mc_top
    rows_per_col = (n + 1) // 2  # 4 题 → 2/2
    row_h = mc_h_total / max(rows_per_col, 1)

    col_w = (CONTENT_W - 0.80) / 2
    col_x = [CONTENT_X + 0.40, CONTENT_X + 0.40 + col_w]

    # 题干+选项字号按"最长一题"的总字符数自适应，避免换行后相互重叠/压页脚
    def _q_chars(mc: dict) -> int:
        return len(str(mc.get("q", ""))) + sum(
            len(str(o)) for o in (mc.get("options") or [])[:3]
        )

    max_chars = max((_q_chars(m) for m in mcs[:n]), default=0)
    if max_chars > 150:
        q_pt = 10.5
    elif max_chars > 115:
        q_pt = 11.5
    elif max_chars > 85:
        q_pt = 12.5
    else:
        q_pt = 14.0

    for i, mc in enumerate(mcs[:n]):
        col_idx = 0 if i < rows_per_col else 1
        row_idx = i if col_idx == 0 else (i - rows_per_col)
        x = col_x[col_idx] + 0.10
        y = mc_top + row_idx * row_h
        tb = slide.shapes.add_textbox(
            Inches(x), Inches(y),
            Inches(col_w - 0.20), Inches(row_h - 0.05),
        )
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        # 第一段：题号 + 题干
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = f"{i + 1}. {mc.get('q', '')}"
        r.font.name = FONT
        r.font.size = Pt(q_pt)
        r.font.color.rgb = BLACK
        r.font.bold = False
        # 后续段：选项 A/B/C 每题占一段（避免选项过长被截断）
        opts = (mc.get("options") or [])[:3]
        for j, opt in enumerate(opts):
            po = tf.add_paragraph()
            po.alignment = PP_ALIGN.LEFT
            po.space_after = Pt(1)
            po.line_spacing = 1.05
            ro = po.add_run()
            ro.text = f"    {chr(ord('A') + j)}. {opt}"
            ro.font.name = FONT
            ro.font.size = Pt(q_pt)
            ro.font.color.rgb = BLACK


# ============================================================
#  Page 5 — Writing 脚手架
# ============================================================

def _build_p5_writing(slide, brand_rgb: tuple, writing: dict) -> None:
    theme = writing.get("theme", "the story")
    _add_title(slide, "Writing", f"Write about {theme}.")

    # 中部黄虚线框 = 5 步骨架
    scaff_top = CONTENT_Y + 1.40
    scaff_h = 2.40
    scaff = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 1.20), Inches(scaff_top),
        Inches(CONTENT_W - 2.40), Inches(scaff_h),
    )
    scaff.adjustments[0] = 0.04
    scaff.fill.solid()
    scaff.fill.fore_color.rgb = RGBColor(0xFD, 0xF5, 0xE0)  # 浅黄
    scaff.line.color.rgb = RGBColor(0xE0, 0xC8, 0x80)
    scaff.line.width = Pt(1.5)
    scaff.shadow.inherit = False

    tf = scaff.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.15)
    tf.word_wrap = True

    # Title 行
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = "Title: "
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0x6B, 0x4D, 0xA8)  # 紫
    r2 = p.add_run()
    r2.text = writing.get("title", "_______________")
    r2.font.name = FONT
    r2.font.bold = True
    r2.font.size = Pt(14)
    r2.font.color.rgb = BLACK

    # 5 步
    step_colors = [
        RGBColor(0x4C, 0xA1, 0x65),  # Beginning - 绿
        RGBColor(0x4C, 0x8F, 0xD8),  # First event - 蓝
        RGBColor(0xE3, 0x76, 0x35),  # Second event - 橙
        RGBColor(0xC0, 0x39, 0x6F),  # Funny event - 红
        RGBColor(0x6B, 0x4D, 0xA8),  # Ending - 紫
    ]
    step_labels = writing.get("step_labels") or [
        "Beginning:", "First event:", "Second event:", "Funny event:", "Ending:",
    ]
    step_contents = writing.get("steps") or [""] * 5

    for i in range(5):
        pi = tf.add_paragraph()
        pi.alignment = PP_ALIGN.LEFT
        pi.space_after = Pt(2)
        rn = pi.add_run()
        rn.text = f"{i + 1}. "
        rn.font.name = FONT
        rn.font.bold = True
        rn.font.size = Pt(13)
        rn.font.color.rgb = BLACK
        rl = pi.add_run()
        rl.text = step_labels[i] + "  "
        rl.font.name = FONT
        rl.font.bold = True
        rl.font.size = Pt(13)
        rl.font.color.rgb = step_colors[i % len(step_colors)]
        # step content 里含 ________ 时切出来用 Arial 字体
        _emit_with_underscore_lock(
            pi, step_contents[i] if i < len(step_contents) else "", 13, BLACK
        )

    # 字数提示
    hint_y = scaff_top + scaff_h + 0.10
    hb = slide.shapes.add_textbox(
        Inches(CONTENT_X), Inches(hint_y),
        Inches(CONTENT_W), Inches(0.25),
    )
    p = hb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"Write {writing.get('min_words', 50)}-{writing.get('max_words', 80)} words."
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0xE6, 0xA8, 0x2B)

    # 蓝色横线写作区
    write_top = hint_y + 0.35
    write_bottom = CONTENT_Y + CONTENT_H - 0.30
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 0.40), Inches(write_top),
        Inches(CONTENT_W - 0.80), Inches(write_bottom - write_top),
    )
    box.adjustments[0] = 0.03
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(0x8C, 0xB6, 0xE6)
    box.line.width = Pt(1.2)
    box.shadow.inherit = False

    # 内部横线（5-6 条）
    line_left = CONTENT_X + 0.70
    line_right = CONTENT_X + CONTENT_W - 0.70
    n_lines = 5
    line_step = (write_bottom - write_top - 0.30) / n_lines
    for k in range(1, n_lines + 1):
        ly = write_top + k * line_step
        ln = slide.shapes.add_connector(
            1,
            Inches(line_left), Inches(ly),
            Inches(line_right), Inches(ly),
        )
        ln.line.color.rgb = RGBColor(0xB7, 0xCD, 0xEC)
        ln.line.width = Pt(0.75)


# ============================================================
#  Page 6 — Mind Map
# ============================================================

def _build_p6_mindmap(slide, rows: list[dict]) -> None:
    _add_title(slide, "Reading", _get_subtitle("read_extended_qa", "Fill in the mind map to organize the story."))

    n_rows = min(len(rows), 5) + 1  # +1 表头
    table_top = CONTENT_Y + 1.40
    table_bottom = CONTENT_Y + CONTENT_H - 0.50
    table_h = table_bottom - table_top
    row_h = table_h / n_rows

    # 3 列宽度
    table_left = CONTENT_X + 0.40
    table_w = CONTENT_W - 0.80
    col1_w = table_w * 0.30
    col2_w = table_w * 0.35
    col3_w = table_w * 0.35
    col_x = [table_left, table_left + col1_w, table_left + col1_w + col2_w]
    col_w = [col1_w, col2_w, col3_w]

    header_fills = [MM_PURPLE, MM_YELLOW, MM_GREEN]
    body_fills = [MM_PURPLE, MM_YELLOW, MM_GREEN]
    headers = ["Character", "Problem", "Solution"]

    # 表头行
    for c in range(3):
        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(col_x[c]), Inches(table_top),
            Inches(col_w[c]), Inches(row_h),
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fills[c]
        cell.line.color.rgb = WHITE
        cell.line.width = Pt(2.0)
        cell.shadow.inherit = False
        tf = cell.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = headers[c]
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(HEADER_PT)
        r.font.color.rgb = BLACK

    # 内容行（5 行）
    keys = ["character", "problem", "solution"]
    body_cell_fills = [
        RGBColor(0xF5, 0xEF, 0xFB),  # 极浅紫
        RGBColor(0xFE, 0xF9, 0xE6),  # 极浅黄
        RGBColor(0xEF, 0xF7, 0xEE),  # 极浅绿
    ]

    for i, row in enumerate(rows[:5]):
        ry = table_top + (i + 1) * row_h
        for c in range(3):
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(col_x[c]), Inches(ry),
                Inches(col_w[c]), Inches(row_h),
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = body_cell_fills[c]
            cell.line.color.rgb = WHITE
            cell.line.width = Pt(2.0)
            cell.shadow.inherit = False
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.10)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c > 0 else PP_ALIGN.LEFT
            text = str(row.get(keys[c], "")).strip()
            if c == 0:
                # 角色列加序号
                text = f"{i + 1}. {text}"
            # 含 ________ 时切出来用 Arial 字体
            _emit_with_underscore_lock(p, text, 13, BLACK)


# ============================================================
#  默认数据 fallback（AI 失败时也能出像样的 worksheet）
# ============================================================

def _resolve_worksheet_data(outline: BookOutline) -> dict:
    """优先用 outline._worksheet_data；否则用 outline 字段构造默认。
    在返回前统一做 v1.8 文本规整：美式拼写、答案格式、难度排序、剔除无效题型。"""
    data = getattr(outline, "_worksheet_data", None)
    if isinstance(data, dict) and data.get("match_pairs"):
        return _normalize_worksheet_data(data)
    return _normalize_worksheet_data(_build_default_data(outline))


# ----- v1.8 文本规整：所有英文走美式 + 答案格式 + 难度排序 -----

# 不允许的活动类型关键词（color / colour 只是涂色，不算 reading 输出）
_BANNED_PROMPT_PATTERNS = [
    "color the ", "colour the ",
    "color in ", "colour in ",
    "circle the picture",  # 仅涂圈没有输出
]


# v2.0 vocab 兜底真定义词典 — 防止 AI 输出 "meaning of X" 占位
# 数据源：常见低龄启蒙词汇 + L5-1 What Makes a Good Friend 词汇 + 高频感觉/动作/状态词
_KID_DICT: dict[str, str] = {
    # L5-1 What Makes a Good Friend
    "nervous":   "feeling worried and not calm",
    "shake":     "to move quickly from side to side or up and down",
    "recess":    "a short break at school for students to play outside",
    "a pile of": "many things lying one on top of another",
    "wooden":    "made of wood",
    # 高频形容词
    "excited":   "feeling very happy and full of energy",
    "amazed":    "feeling very surprised in a good way",
    "worried":   "feeling that something bad might happen",
    "happy":     "feeling glad and full of joy",
    "sad":       "feeling unhappy",
    "scared":    "feeling afraid",
    "tired":     "feeling like you need to rest or sleep",
    "kind":      "nice and caring to others",
    "quiet":     "making very little sound",
    "loud":      "making a lot of sound",
    "friendly":  "kind and nice to other people",
    "smart":     "able to think and learn quickly",
    "brave":     "not afraid of dangerous or difficult things",
    "shy":       "feeling not comfortable talking to new people",
    "gentle":    "kind and soft, not rough",
    # 高频动作
    "share":     "to give part of what you have to someone else",
    "help":      "to do something nice for someone who needs it",
    "listen":    "to pay attention to a sound or a person",
    "smile":     "to make a happy face with your mouth turned up",
    "laugh":     "to make a happy sound when something is funny",
    "look":      "to use your eyes to see something",
    "run":       "to move very fast on your feet",
    "walk":      "to move on your feet at a normal speed",
    "jump":      "to push yourself up into the air with your legs",
    "grab":      "to take hold of something quickly",
    "drop":      "to let something fall to the ground",
    "pick up":   "to lift something from the ground or a surface",
    "give":      "to let someone have something",
    # 高频名词
    "friend":    "someone you like and spend time with",
    "friendship": "the feeling of being friends with someone",
    "classmate": "a person in the same class as you at school",
    "kindness":  "the quality of being nice and caring",
    "teacher":   "a person whose job is to teach in a school",
    "student":   "a person who is learning at a school",
    "class":     "a group of students who learn together",
    "school":    "a place where children go to learn",
    "desk":      "a kind of table you sit at when you read or write",
    "chair":     "something you sit on",
    "book":      "pages with words and pictures joined together to read",
    "pencil":    "a thin stick used to write or draw",
    "eraser":    "a small piece of rubber used to remove pencil marks",
    "hamster":   "a small soft animal with short legs that people keep as a pet",
    # L1-L4 主题词
    "culture":   "the customs, beliefs, and ways of life of a group of people",
    "castle":    "a large old strong building where kings or queens used to live",
    "bagpipes":  "a musical instrument with a bag and pipes, played in Scotland",
    "journey":   "a trip from one place to another",
}


def _fix_vocab_def(word: str, def_text: str) -> str:
    """如果 def 是占位（meaning of X / definition of X / 空），用 _KID_DICT 兜底。"""
    bad = (
        not def_text
        or def_text.strip().lower().startswith("meaning of ")
        or def_text.strip().lower().startswith("definition of ")
        or def_text.strip().lower() == word.strip().lower()
    )
    if not bad:
        return def_text
    return _KID_DICT.get(word.strip().lower(), def_text or f"see story for the meaning of {word}")


def _normalize_worksheet_data(data: dict) -> dict:
    out = dict(data)
    # v2.0：reading_q_count 在 data 上可配置（4/6/8），默认 4
    reading_q_count = max(4, min(8, int(data.get("_reading_q_count") or 4)))

    # 1) match_pairs：word 小写无标点；def 美式 + 首字母小写（更像词典定义）+ v2.0 兜底真定义
    pairs = []
    for p in (data.get("match_pairs") or []):
        word = format_word_answer(p.get("word", ""))
        definition = _to_us_spelling(str(p.get("def", "")).strip().rstrip("."))
        definition = _fix_vocab_def(word, definition)
        if word:
            pairs.append({"word": word, "def": definition, "_len": len(word)})
    # v2.0：题数上限改为 4（按官方模板，每页 3-4 题最佳）
    pairs.sort(key=lambda x: (x["_len"], x["word"]))
    for p in pairs:
        p.pop("_len", None)
    out["match_pairs"] = pairs[:4]

    # 2) word_bank：小写美式
    out["word_bank"] = [format_word_answer(w) for w in (data.get("word_bank") or []) if w]

    # 3) fill_blanks：句子格式 + 答案小写
    fills = []
    for q in (data.get("fill_blanks") or []):
        sent = format_sentence_answer(q.get("sentence", ""))
        ans = format_word_answer(q.get("answer", ""))
        # 难度启发：句长升序
        fills.append({"sentence": sent, "answer": ans, "_len": len(sent)})
    # v2.0：上限 4 题（对齐官方每页 3-4 题）
    fills.sort(key=lambda x: x["_len"])
    for x in fills:
        x.pop("_len", None)
    out["fill_blanks"] = fills[:4]

    # 4) sentence_mcs：每个 option 走句子格式；过滤 color-only 等无效题
    mcs = []
    for q in (data.get("sentence_mcs") or []):
        opts = [format_sentence_answer(o) for o in (q.get("options") or []) if o]
        # 过滤无效
        if any(any(bad in o.lower() for bad in _BANNED_PROMPT_PATTERNS) for o in opts):
            continue
        if len(opts) < 2:
            continue
        mcs.append({
            "options": opts[:2],
            "correct": q.get("correct", 0),
            "_len": sum(len(o) for o in opts),
        })
    mcs.sort(key=lambda x: x["_len"])
    for x in mcs:
        x.pop("_len", None)
    out["sentence_mcs"] = mcs[:4]

    # 5) reading_text：美式
    out["reading_text"] = _to_us_spelling(str(data.get("reading_text", "")).strip())

    # 6) reading_mcs：题干补问号 + 美式 + 首字母大写；选项 smart_format
    rmcs = []
    for q in (data.get("reading_mcs") or []):
        stem = _to_us_spelling(str(q.get("q", "")).strip().rstrip(".!?"))
        # 独立 i 大写
        import re as _re
        stem = _re.sub(r"\bi\b", "I", stem)
        stem = _re.sub(r"\bi(['\u2019])", r"I\1", stem)
        # 首字母大写
        if stem:
            stem = stem[0].upper() + stem[1:]
        if stem and stem[-1] not in "?？":
            stem += "?"
        opts = [smart_format_answer(o) for o in (q.get("options") or []) if o]
        # 过滤 color-only 活动
        stem_low = stem.lower()
        if any(bad in stem_low for bad in _BANNED_PROMPT_PATTERNS):
            continue
        if not stem or len(opts) < 2:
            continue
        rmcs.append({
            "q": stem, "options": opts[:3],
            "correct": q.get("correct", 0),
            "_len": len(stem),
        })
    # 题干越短越简单，放顶上；上限按 _reading_q_count 配置（4/6/8）
    rmcs.sort(key=lambda x: x["_len"])
    for x in rmcs:
        x.pop("_len", None)
    out["reading_mcs"] = rmcs[:reading_q_count]

    # 7) writing：theme/title/steps 美式；步骤句子格式
    writing = dict(data.get("writing") or {})
    writing["theme"] = _to_us_spelling(str(writing.get("theme", "the story")))
    writing["title"] = _to_us_spelling(str(writing.get("title", "")))
    steps_raw = writing.get("steps") or [""] * 5
    writing["steps"] = [
        format_sentence_answer(s) if s.strip() else "" for s in steps_raw
    ]
    out["writing"] = writing

    # 8) mind_map：character 短语形式（首字母大写、无句号）；problem/solution 句子形式
    mm = []
    for r in (data.get("mind_map") or [])[:5]:
        ch_raw = str(r.get("character", "")).strip().rstrip(".")
        ch = _to_us_spelling(ch_raw) if ch_raw else ""
        if ch:
            ch = ch[0].upper() + ch[1:]
        mm.append({
            "character": ch,
            "problem": format_sentence_answer(r.get("problem", "")),
            "solution": format_sentence_answer(r.get("solution", "")),
        })
    out["mind_map"] = mm

    return out


def _build_default_data(outline: BookOutline) -> dict:
    words = (
        outline.vocabulary_mastery
        or outline.vocabulary_simple
        or ["nervous", "shake", "recess", "wooden", "a pile of"]
    )[:5]
    words = list(words) + ["word"] * max(0, 5 - len(words))
    words = words[:5]

    pages = outline.pages or []
    story_text = " ".join(
        (p.text or "").strip() for p in pages if (p.text or "").strip()
    ).strip()
    if not story_text:
        story_text = "Story text goes here."

    story_sents: list[str] = []
    for p in pages:
        if (p.text or "").strip():
            story_sents.append(p.text.strip())
    if not story_sents:
        story_sents = ["Story sentence goes here."]

    def _pair_sent(idx: int) -> dict:
        """从故事页取真句子做选项，避免 'Not X' 这种生硬兜底。"""
        correct = story_sents[idx % len(story_sents)]
        distractor_idx = (idx + 1) % len(story_sents)
        if len(story_sents) > 1 and distractor_idx == idx % len(story_sents):
            distractor_idx = (idx + 2) % len(story_sents)
        distractor = story_sents[distractor_idx]
        return {"options": [correct, distractor], "correct": 0}

    return {
        "match_pairs": [
            {"word": w, "def": f"definition of {w}"} for w in words
        ],
        "word_bank": list(words),
        "fill_blanks": [
            {"sentence": f"I _______ when I see {w}.", "answer": w} for w in words
        ],
        "sentence_mcs": [_pair_sent(i) for i in range(min(4, len(story_sents)))],
        "reading_text": story_text,
        "reading_mcs": [
            {
                "q": f"Question {i + 1}?",
                "options": ["Option A", "Option B", "Option C"],
                "correct": 0,
            }
            for i in range(8)
        ],
        "writing": {
            "theme": (outline.theme or "the story").strip(),
            "title": f"{outline.title}",
            "steps": ["", "", "", "", ""],
            "step_labels": [
                "Beginning:", "First event:", "Second event:",
                "Funny event:", "Ending:",
            ],
            "min_words": 50, "max_words": 80,
        },
        "mind_map": [
            {"character": "Main character", "problem": "Problem statement.",
             "solution": "Solution statement."},
        ] * 5,
    }


# ============================================================
#  工具
# ============================================================

def _level_label(level: str) -> str:
    s = (level or "").strip()
    if not s:
        return "Level 1"
    if s.lower().startswith("smart"):
        return "Smart"
    if s.lower().startswith("level"):
        return s
    digits = "".join(ch for ch in s if ch.isdigit())
    return f"Level {digits or '1'}"


def safe_filename(name: str) -> str:
    """与 ppt_builder.safe_filename 对齐，方便复用。"""
    import re
    safe = re.sub(r"[^\w\u4e00-\u9fff -]+", "_", name).strip("_ ")
    return safe + ".pptx"
