"""组装 9 页绘本 PPT：封面 + 7 故事 + 元信息页，统一 Poppins Bold。"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config import (
    BLACK, FONT_BOLD, FONT_FAMILY, FONT_SIZE_BADGE, FONT_SIZE_BODY,
    FONT_SIZE_META_BODY, FONT_SIZE_META_HEAD, FONT_SIZE_PAGE_NUM,
    FONT_SIZE_TITLE, LIGHT_GRAY_BORDER, ORANGE_BADGE,
    PAGE_NUM_DIAMETER_IN, PAGE_NUM_MARGIN_IN,
    SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN, TEXT_BOX_PADDING_IN, TEXT_BOX_WIDTH_RATIO,
    WHITE, text_box_position,
)
from parser import BookOutline, PageSpec


def build_picturebook_pptx(
    outline: BookOutline,
    image_paths: list[Path],
    out_path: Path,
) -> Path:
    """images: 0=cover, 1-7=story, 8 ignored (元信息页用文字)."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]

    if len(image_paths) < 8:
        raise ValueError(f"需要至少 8 张图（封面+7故事），实际 {len(image_paths)}")

    # p1 封面
    _build_cover(prs.slides.add_slide(blank), outline, image_paths[0])

    # p2–p8 故事
    for i in range(1, 8):
        slide = prs.slides.add_slide(blank)
        _build_story(slide, outline.pages[i], image_paths[i], page_number=i + 1)

    # p9 元信息页
    _build_metadata(prs.slides.add_slide(blank), outline)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# ---------- 封面 ----------
def _build_cover(slide, outline: BookOutline, image_path: Path) -> None:
    sw, sh = Inches(SLIDE_WIDTH_IN), Inches(SLIDE_HEIGHT_IN)

    if Path(image_path).exists():
        pic = slide.shapes.add_picture(str(image_path), 0, 0, width=sw, height=sh)
        _send_to_back(slide, pic)

    # 书名（上方居中，靠上）
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.45), Inches(SLIDE_WIDTH_IN - 1.6), Inches(1.4)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run() if p.runs else _ensure_run(p), outline.title, Pt(FONT_SIZE_TITLE), BLACK)

    # 右上双徽章 Level X / Book N
    badge_w, badge_h = Inches(1.5), Inches(0.5)
    badge_x = Inches(SLIDE_WIDTH_IN - 1.5 - 0.4)
    _add_badge(slide, badge_x, Inches(0.35), badge_w, badge_h, f"Level {_clean_num(outline.level)}")
    _add_badge(slide, badge_x, Inches(0.95), badge_w, badge_h, f"Book {_clean_num(outline.book_number)}")


# ---------- 故事页 ----------
def _build_story(slide, page: PageSpec, image_path: Path, page_number: int) -> None:
    sw, sh = Inches(SLIDE_WIDTH_IN), Inches(SLIDE_HEIGHT_IN)

    if Path(image_path).exists():
        pic = slide.shapes.add_picture(str(image_path), 0, 0, width=sw, height=sh)
        _send_to_back(slide, pic)

    # 文字框（白底矩形 + 黑色 Poppins Bold）
    corner = page.text_corner or "top-left"
    left_in, top_in = text_box_position(corner)
    box_w_in = SLIDE_WIDTH_IN * TEXT_BOX_WIDTH_RATIO
    box_h_in = 1.6

    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in),
        Inches(box_w_in), Inches(box_h_in),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*WHITE)
    box.line.color.rgb = RGBColor(*WHITE)
    box.line.width = Pt(0.0)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(TEXT_BOX_PADDING_IN)
    tf.margin_top = tf.margin_bottom = Inches(TEXT_BOX_PADDING_IN)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_run(_ensure_run(p), page.text, Pt(FONT_SIZE_BODY), BLACK)

    # 页码：偶数页左下，奇数页右下
    _add_page_number(slide, page_number)


# ---------- 元信息页（p9 封底）----------
def _build_metadata(slide, outline: BookOutline) -> None:
    # 左侧大框
    frame_left, frame_top = Inches(0.7), Inches(0.7)
    frame_w, frame_h = Inches(6.0), Inches(6.0)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, frame_left, frame_top, frame_w, frame_h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*WHITE)
    box.line.color.rgb = RGBColor(*LIGHT_GRAY_BORDER)
    box.line.width = Pt(1.0)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.3)

    def line(text: str, *, head: bool, indent: int = 0) -> None:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.level = indent
        size = Pt(FONT_SIZE_META_HEAD if head else FONT_SIZE_META_BODY)
        _set_run(_ensure_run(p), text, size, BLACK)

    line(f"Level: {_clean_num(outline.level) or 'Smart'}", head=True)
    line(f"Book: {_clean_num(outline.book_number)}", head=True)
    line(f"CEFR: {outline.cefr or '-'}", head=True)
    line(f"Lexile: {outline.lexile or '-'}", head=True)
    line(f"Word count: {outline.total_words}", head=True)
    if outline.phonics:
        line(f"Phonics: {outline.phonics}", head=True)
    if outline.grammar_focus:
        line(f"Grammar: {outline.grammar_focus}", head=True)
    if outline.reader_type:
        line(f"Reader Type: {outline.reader_type}", head=True)
    line("Vocabulary:", head=True)

    # L0/L1/L2/Smart → 双行 Mastery + Exposure（每行 3-4 词）
    # L3-L6        → 单行 Vocabulary 4 词（lemma 原型）
    if outline.is_dual_vocab_level:
        if outline.has_double_vocab:
            line(f"Mastery:  {', '.join(outline.vocabulary_mastery) or '-'}",
                 head=False, indent=1)
            line(f"Exposure: {', '.join(outline.vocabulary_exposure) or '-'}",
                 head=False, indent=1)
        else:
            line(", ".join(outline.vocabulary_for_display) or "-",
                 head=False, indent=1)
    else:
        words = outline.vocabulary_for_display[:4]
        line(", ".join(words) or "-", head=False, indent=1)


# ---------- 通用 ----------
def _add_badge(slide, left, top, width, height, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5  # 完全圆角
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*ORANGE_BADGE)
    shape.line.color.rgb = RGBColor(*BLACK)
    shape.line.width = Pt(1.2)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(_ensure_run(p), text, Pt(FONT_SIZE_BADGE), WHITE)


def _add_page_number(slide, page_number: int) -> None:
    d = PAGE_NUM_DIAMETER_IN
    m = PAGE_NUM_MARGIN_IN
    top = SLIDE_HEIGHT_IN - m - d
    if page_number % 2 == 0:  # 偶数页 左下
        left = m
    else:                      # 奇数页 右下
        left = SLIDE_WIDTH_IN - m - d

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(d), Inches(d)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(*WHITE)
    circle.line.color.rgb = RGBColor(*BLACK)
    circle.line.width = Pt(0.75)
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.margin_left = tf.margin_right = Inches(0.0)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(_ensure_run(p), str(page_number), Pt(FONT_SIZE_PAGE_NUM), BLACK)


def _ensure_run(p):
    if p.runs:
        for r in p.runs[1:]:
            r.text = ""
        p.runs[0].text = ""
        return p.runs[0]
    return p.add_run()


def _set_run(run, text: str, size, rgb_tuple: tuple[int, int, int]) -> None:
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.bold = FONT_BOLD
    run.font.size = size
    run.font.color.rgb = RGBColor(*rgb_tuple)


def _send_to_back(slide, shape) -> None:
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def _clean_num(s: str) -> str:
    if not s:
        return "-"
    m = re.search(r"\d+", s)
    return m.group(0) if m else s


def safe_filename(title: str) -> str:
    name = re.sub(r"[^\w\s-]", "", title, flags=re.UNICODE)
    name = re.sub(r"\s+", "_", name.strip()) or "PictureBook"
    return f"{name}.pptx"
